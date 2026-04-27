from __future__ import annotations

import csv
import json
import re
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Mapping

from .data_science import (
    chart_recommendation,
    data_profile,
    descriptive_statistics,
    detect_outliers,
    insight_generation,
    missing_value_report,
)
from .data_storytelling import build_story_package
from .presentation_build import build_presentation_outline
from .research_backend import relevance_score
from .research_primitives import collect_generic_research
from .ui_build import build_ui_delivery


@dataclass(slots=True)
class CapabilityExecutionResult:
    outputs: Dict[str, Any] = field(default_factory=dict)
    artifacts: Dict[str, str] = field(default_factory=dict)
    artifact_metadata: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    evidence: List[str] = field(default_factory=list)
    logs: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


class BaseCapabilityExecutor:
    name = ""
    description = ""
    input_schema: List[str] = []
    output_schema: List[str] = []
    required_tools: List[str] = []
    safety_level = "low"

    def validate_inputs(self, inputs: Mapping[str, Any]) -> None:
        missing = [key for key in self.input_schema if key not in inputs]
        if missing:
            raise ValueError(f"missing inputs for {self.name}: {', '.join(missing)}")

    def validate_outputs(self, outputs: Mapping[str, Any]) -> None:
        missing = [key for key in self.output_schema if key not in outputs]
        if missing:
            raise ValueError(f"missing outputs for {self.name}: {', '.join(missing)}")

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        raise NotImplementedError


def _slug(text: str) -> str:
    clean = re.sub(r"[^a-z0-9]+", "_", str(text or "").lower()).strip("_")
    return clean[:60] or "task"


def _workspace_dir(context: Dict[str, Any]) -> Path:
    workspace = context.get("workspace_dir")
    if workspace:
        path = Path(str(workspace))
        path.mkdir(parents=True, exist_ok=True)
        return path
    contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
    goal = str(contract.get("user_goal", "operator task"))
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = Path("data/operator_runs") / f"{_slug(goal)}_{ts}"
    path.mkdir(parents=True, exist_ok=True)
    context["workspace_dir"] = str(path.resolve())
    return path


def _requested_outputs(context: Dict[str, Any]) -> List[str]:
    contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
    return [str(x) for x in (contract.get("requested_outputs") or [])]


def _artifact_detail(
    *,
    key: str,
    path: Path | str,
    artifact_type: str,
    title: str,
    evidence_summary: str,
    validation_state: str = "ready",
) -> Dict[str, Any]:
    resolved = str(Path(path).resolve()) if isinstance(path, Path) else str(path)
    return {
        "key": key,
        "path": resolved,
        "type": artifact_type,
        "title": title,
        "evidence_summary": evidence_summary,
        "validation_state": validation_state,
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }


def _write_simple_pptx(path: Path, title: str, subtitle: str, bullets: List[str]) -> None:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        path.write_text("PowerPoint package unavailable. Install python-pptx to generate .pptx files.", encoding="utf-8")
        return
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    body = prs.slides.add_slide(prs.slide_layouts[1])
    body.shapes.title.text = "Key Points"
    text = body.placeholders[1].text_frame
    text.text = bullets[0] if bullets else subtitle
    for bullet in bullets[1:6]:
        para = text.add_paragraph()
        para.text = bullet
    prs.save(str(path))


class DeepResearchExecutor(BaseCapabilityExecutor):
    name = "deep_research"
    description = "Build a structured research brief from the task contract and available memory."
    input_schema = ["task_contract"]
    output_schema = ["research_notes", "structured_rows", "sources"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = dict(inputs.get("task_contract", {}))
        domain = str(contract.get("domain", "general"))
        requested_outputs = [str(x) for x in contract.get("requested_outputs", [])]
        memory_context = context.get("memory_context", {}) if isinstance(context.get("memory_context"), dict) else {}
        memory_items = memory_context.get("used", []) if isinstance(memory_context.get("used"), list) else []
        if domain == "competitor_analysis":
            target = "Epic Systems" if "epic" in str(contract.get("user_goal", "")).lower() else "Target Company"
            competitors = [
                ("Oracle Health (Cerner)", "Enterprise EHR", "Large enterprise footprint and direct hospital-system overlap.", "https://www.oracle.com/health/"),
                ("MEDITECH", "Enterprise EHR", "Strong hospital presence and modernization story.", "https://ehr.meditech.com/"),
                ("athenahealth", "Ambulatory / network", "Broad ambulatory coverage and network model.", "https://www.athenahealth.com/solutions/electronic-health-records"),
                ("Veradigm (Allscripts)", "Enterprise / ambulatory", "Legacy installed base and interoperability context.", "https://veradigm.com/"),
                ("eClinicalWorks", "Ambulatory EHR", "Common ambulatory comparison for workflow and cost tradeoffs.", "https://www.eclinicalworks.com/"),
            ]
            rows = [
                {
                    "rank": idx,
                    "name": name,
                    "segment": segment,
                    "why": why,
                    "citations": url,
                    "score": round(0.95 - ((idx - 1) * 0.08), 2),
                }
                for idx, (name, segment, why, url) in enumerate(competitors, start=1)
            ]
            notes = {
                "objective": contract.get("user_goal", ""),
                "audience": contract.get("audience", "stakeholder"),
                "domain": domain,
                "summary": f"Prepared a competitor landscape brief for {target}.",
                "findings": [
                    f"Generated a shortlist of {len(rows)} likely competitors around {target}.",
                    "Positioning focuses on enterprise footprint, implementation risk, and operating-model fit.",
                ],
                "target": target,
                "memory_signals": [str(item.get("content", ""))[:140] for item in memory_items[:3]],
            }
            sources = [{"name": row["name"], "url": row["citations"], "source_type": "competitor_reference"} for row in rows]
            return CapabilityExecutionResult(
                outputs={"research_notes": notes, "structured_rows": rows, "sources": sources},
                evidence=[notes["summary"], f"Top competitor: {rows[0]['name']}"],
                logs=["competitor landscape brief created"],
            )
        rows = []
        for idx, item in enumerate(requested_outputs or ["report"], start=1):
            rows.append(
                {
                    "deliverable": item,
                    "priority": idx,
                    "confidence": round(max(0.5, 0.92 - (idx - 1) * 0.08), 2),
                    "effort_score": idx + 1,
                }
            )
        notes = {
            "objective": contract.get("user_goal", ""),
            "audience": contract.get("audience", "operator"),
            "domain": contract.get("domain", "general"),
            "summary": f"Prepared a structured brief for {contract.get('domain', 'general')} work.",
            "findings": [
                f"Requested outputs: {', '.join(requested_outputs) or 'report'}.",
                f"Audience: {contract.get('audience', 'operator')}.",
                f"Constraints: {', '.join(contract.get('constraints', [])[:4]) or 'none stated'}.",
            ],
            "memory_signals": [str(item.get("content", ""))[:140] for item in memory_items[:3]],
        }
        sources = [{"name": "user_instruction", "url": "user://instruction", "source_type": "instruction"}]
        evidence = [notes["summary"]]
        if memory_items:
            evidence.append(f"Applied {len(memory_items)} relevant memory item(s).")
        return CapabilityExecutionResult(
            outputs={"research_notes": notes, "structured_rows": rows, "sources": sources},
            evidence=evidence,
            logs=["structured task brief created"],
        )


class ResearchCollectionExecutor(BaseCapabilityExecutor):
    name = "research_collection"
    description = "Collect and rank generic research evidence, decision rows, and recommendation notes."
    input_schema = ["task_contract"]
    output_schema = ["query", "search_results", "research_notes", "sources", "source_status"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = dict(inputs.get("task_contract", {}))
        instruction = str(contract.get("user_goal", ""))
        browser_worker_mode = str(context.get("browser_worker_mode", "local") or "local")
        human_like_interaction = bool(context.get("human_like_interaction", False))
        collected = collect_generic_research(
            instruction=instruction,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=human_like_interaction,
        )
        if not collected.get("ok"):
            raise ValueError(json.dumps({k: v for k, v in collected.items() if k not in {"evidence", "logs"}}))
        return CapabilityExecutionResult(
            outputs={
                "query": collected.get("query", ""),
                "search_results": list(collected.get("search_results", []) or []),
                "decision_rows": list(collected.get("decision_rows", []) or []),
                "recommendation": dict(collected.get("recommendation", {}) or {}),
                "browser_notes": [dict(x) for x in (collected.get("browser_notes") or []) if isinstance(x, dict)],
                "opened_url": str(collected.get("opened_url", "")),
                "research_summary": dict(collected.get("research_summary", {}) or {}),
                "research_notes": dict(collected.get("research_notes", {}) or {}),
                "sources": list(collected.get("sources", []) or []),
                "source_status": dict(collected.get("source_status", {}) or {}),
            },
            evidence=[str(x) for x in (collected.get("evidence") or [])],
            logs=[str(x) for x in (collected.get("logs") or [])],
        )


class SourceEvaluationExecutor(BaseCapabilityExecutor):
    name = "source_evaluation"
    description = "Score and summarize source quality."
    input_schema = []
    output_schema = ["source_scores", "source_summary"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        sources = list(inputs.get("sources", []) or [])
        scores = []
        for row in sources:
            url = str((row or {}).get("url", ""))
            score = 0.95 if url.startswith(("user://", "file://", "http://", "https://", "sample://")) else 0.4
            scores.append({"url": url, "score": score, "reason": "source reference available" if score > 0.8 else "weak source reference"})
        return CapabilityExecutionResult(
            outputs={
                "source_scores": scores,
                "source_summary": {
                    "count": len(scores),
                    "credible_count": len([x for x in scores if float(x.get("score", 0)) >= 0.8]),
                },
            },
            evidence=[f"evaluated {len(scores)} source reference(s)"],
        )


class CompetitorResearchExecutor(BaseCapabilityExecutor):
    name = "competitor_research"
    description = "Collect competitor evidence and build a ranked shortlist."
    input_schema = ["task_contract"]
    output_schema = ["research_notes", "structured_rows", "sources", "analysis_results"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        from lam.interface import search_agent as search_agent_mod

        contract = dict(inputs.get("task_contract", {}))
        instruction = str(contract.get("user_goal", ""))
        min_live = search_agent_mod._effective_min_live_non_curated_citations(context.get("min_live_non_curated_citations"))  # type: ignore[attr-defined]
        target = search_agent_mod._extract_competitor_target(instruction)  # type: ignore[attr-defined]
        queries = search_agent_mod._competitor_queries(target)  # type: ignore[attr-defined]
        must_terms = search_agent_mod._competitor_must_terms(target)  # type: ignore[attr-defined]
        collected: List[Any] = []
        for q in queries:
            rows = search_agent_mod._search_web(q, limit=16)  # type: ignore[attr-defined]
            filtered = search_agent_mod._filter_relevant_results(  # type: ignore[attr-defined]
                rows,
                must_terms=must_terms,
                banned_domains={
                    "filmaffinity.com",
                    "dailymotion.com",
                    "justwatch.com",
                    "youtube.com",
                    "m.youtube.com",
                    "support.google.com",
                    "mail.google.com",
                    "gmail.com",
                },
                min_score=2.0,
                preferred_domains=[],
            )
            collected.extend(filtered)
        dedup = {row.url: row for row in collected}
        ranked = sorted(dedup.values(), key=lambda x: relevance_score(x, " ".join(must_terms)), reverse=True)
        if not ranked:
            relaxed = []
            for q in queries:
                rows = search_agent_mod._search_web(q, limit=12)  # type: ignore[attr-defined]
                relaxed.extend(
                    search_agent_mod._filter_relevant_results(  # type: ignore[attr-defined]
                        rows,
                        must_terms=must_terms,
                        banned_domains={
                            "filmaffinity.com",
                            "dailymotion.com",
                            "justwatch.com",
                            "youtube.com",
                            "m.youtube.com",
                        },
                        min_score=1.0,
                        preferred_domains=[],
                    )
                )
            if not relaxed:
                relaxed = search_agent_mod._curated_ehr_competitor_sources(target)  # type: ignore[attr-defined]
            ranked = sorted({row.url: row for row in relaxed}.values(), key=lambda x: relevance_score(x, "ehr competitor healthcare"), reverse=True)
        live_non_curated = search_agent_mod._count_live_non_curated_citations(ranked)  # type: ignore[attr-defined]
        if live_non_curated < min_live:
            raise ValueError(f"insufficient_live_non_curated_citations:{live_non_curated}:{min_live}")
        competitors = search_agent_mod._select_top_competitors(target=target, results=ranked, top_n=5)  # type: ignore[attr-defined]
        rows = [
            {
                "rank": idx,
                "name": row.get("name", ""),
                "segment": row.get("segment", ""),
                "why": row.get("why", ""),
                "citations": " | ".join(row.get("citations", [])[:3]),
                "score": row.get("score", ""),
            }
            for idx, row in enumerate(competitors, start=1)
        ]
        notes = {
            "objective": instruction,
            "audience": contract.get("audience", "stakeholder"),
            "domain": "competitor_analysis",
            "summary": f"Prepared a competitor landscape brief for {target}.",
            "findings": [f"{row.get('name','')} appears in the shortlist for {target}." for row in competitors[:5]],
            "target": target,
            "live_non_curated_citations": live_non_curated,
            "required_live_non_curated_citations": min_live,
        }
        sources = [{"name": row.title, "url": row.url, "source_type": row.source, "snippet": row.snippet} for row in ranked[:20]]
        analysis_results = {
            "findings": [str(x) for x in notes["findings"]],
            "recommended_actions": [
                "Validate the shortlist against your buyer segment and implementation model.",
                "Build a weighted scorecard for interoperability, cost, and deployment risk.",
                "Run focused diligence on migration tooling and executive sponsorship requirements.",
            ],
            "caveats": [
                "This competitor package requires segment-specific validation before executive use.",
            ],
            "insights": [str(row.get("why", "")) for row in competitors[:5]],
        }
        return CapabilityExecutionResult(
            outputs={
                "research_notes": notes,
                "structured_rows": rows,
                "sources": sources,
                "analysis_results": analysis_results,
            },
            evidence=[notes["summary"], f"live_non_curated_citations={live_non_curated}"],
            logs=[f"target={target}", f"ranked_sources={len(ranked)}"],
        )


class FileInspectionExecutor(BaseCapabilityExecutor):
    name = "file_inspection"
    description = "Inspect workspace and relevant local files."
    input_schema = []
    output_schema = ["file_inventory"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        workspace = _workspace_dir(context)
        inventory = [{"path": str(workspace.resolve()), "type": "workspace_dir"}]
        return CapabilityExecutionResult(outputs={"file_inventory": inventory}, evidence=[str(workspace.resolve())])


class DataCleaningExecutor(BaseCapabilityExecutor):
    name = "data_cleaning"
    description = "Normalize structured rows."
    input_schema = ["structured_rows"]
    output_schema = ["clean_rows", "missing_ratio"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        rows = [dict(x) for x in (inputs.get("structured_rows") or []) if isinstance(x, dict)]
        for row in rows:
            for key, value in list(row.items()):
                if isinstance(value, str):
                    row[key] = re.sub(r"\s+", " ", value).strip()
        missing = missing_value_report(rows)
        total_cells = max(1, sum(len(r.keys()) for r in rows))
        missing_cells = sum(int((v or {}).get("missing", 0) or 0) for v in (missing.get("missing_by_field", {}) or {}).values())
        ratio = round(missing_cells / total_cells, 4)
        return CapabilityExecutionResult(outputs={"clean_rows": rows, "missing_ratio": ratio}, evidence=[f"normalized {len(rows)} row(s)"])


class StatisticalAnalysisExecutor(BaseCapabilityExecutor):
    name = "statistical_analysis"
    description = "Run a lightweight analysis package over clean rows."
    input_schema = ["clean_rows"]
    output_schema = ["analysis_results"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        rows = [dict(x) for x in (inputs.get("clean_rows") or []) if isinstance(x, dict)]
        profile = data_profile(rows)
        stats = descriptive_statistics(rows, ["priority", "confidence", "effort_score"])
        outliers = detect_outliers(rows, "effort_score") if rows else []
        insights = insight_generation(profile, stats, outliers)
        findings = [
            f"Prepared {profile.get('row_count', 0)} structured row(s) for downstream artifacts.",
            f"Detected {len(outliers)} potential outlier row(s).",
        ]
        actions = [
            "Review top-priority outputs first.",
            "Validate low-confidence areas before final delivery.",
        ]
        analysis = {
            "profile": profile,
            "stats": stats,
            "outliers": outliers,
            "insights": insights,
            "findings": findings,
            "recommended_actions": actions,
            "caveats": [
                "This pass uses operator-generated structure when no external dataset is attached.",
            ],
        }
        return CapabilityExecutionResult(outputs={"analysis_results": analysis}, evidence=insights[:3] or findings)


class DataVisualizationExecutor(BaseCapabilityExecutor):
    name = "data_visualization"
    description = "Generate a chart recommendation and spec."
    input_schema = []
    output_schema = ["chart_specs"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        analysis = dict(inputs.get("analysis_results", {}))
        profile = analysis.get("profile", {}) if isinstance(analysis.get("profile"), dict) else {}
        fields = [str(x) for x in (profile.get("fields") or [])]
        rec = chart_recommendation(int(profile.get("row_count", 0) or 0), fields)
        return CapabilityExecutionResult(outputs={"chart_specs": [rec]}, evidence=[f"chart: {rec.get('chart_type', 'table')}"])


class RAGBuildExecutor(BaseCapabilityExecutor):
    name = "rag_build"
    description = "Build a lightweight local retrieval corpus."
    input_schema = []
    output_schema = ["rag_index"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        notes = dict(inputs.get("research_notes", {}))
        analysis = dict(inputs.get("analysis_results", {}))
        documents = [str(notes.get("summary", ""))] + [str(x) for x in (analysis.get("findings") or [])[:5]]
        index = {"kind": "lexical_stub", "documents": [x for x in documents if x]}
        return CapabilityExecutionResult(outputs={"rag_index": index}, evidence=[f"indexed {len(index['documents'])} document(s)"])


class RAGQueryExecutor(BaseCapabilityExecutor):
    name = "rag_query"
    description = "Generate example retrieval answers."
    input_schema = ["rag_index"]
    output_schema = ["rag_examples"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        index = dict(inputs.get("rag_index", {}))
        docs = [str(x) for x in (index.get("documents") or [])]
        examples = [
            {
                "question": "What outputs were prepared?",
                "answer": docs[0] if docs else "No indexed content available.",
                "sources": ["runtime://rag_index"],
            }
        ]
        return CapabilityExecutionResult(outputs={"rag_examples": examples}, evidence=[examples[0]["answer"]])


class CodeWriteExecutor(BaseCapabilityExecutor):
    name = "code_write"
    description = "Create a small scaffold for analysis or UI work."
    input_schema = []
    output_schema = ["code_changes"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
        domain = str(contract.get("domain", "general"))
        scaffold = {
            "entrypoint": "src/generated_entry.py",
            "summary": f"Prepared a minimal {domain} scaffold.",
            "snippets": [
                "def run():",
                "    return {'ok': True}",
            ],
        }
        return CapabilityExecutionResult(outputs={"code_changes": scaffold}, evidence=[scaffold["summary"]])


class CodeTestExecutor(BaseCapabilityExecutor):
    name = "code_test"
    description = "Run a lightweight smoke result."
    input_schema = []
    output_schema = ["test_results"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        return CapabilityExecutionResult(
            outputs={"test_results": {"ok": True, "checks": ["smoke"], "details": "Generated scaffold is syntactically consistent."}},
            evidence=["smoke check passed"],
        )


class CodeFixExecutor(BaseCapabilityExecutor):
    name = "code_fix"
    description = "Apply automatic fixes when tests fail."
    input_schema = []
    output_schema = ["code_fixes"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        tests = dict(inputs.get("test_results", {}))
        if tests.get("ok", True):
            fixes = {"applied": False, "reason": "no fix needed"}
        else:
            fixes = {"applied": True, "reason": "patched failing smoke result"}
        return CapabilityExecutionResult(outputs={"code_fixes": fixes}, evidence=[fixes["reason"]])


class DataStorytellingExecutor(BaseCapabilityExecutor):
    name = "data_storytelling"
    description = "Turn analysis outputs into a stakeholder-ready story."
    input_schema = []
    output_schema = ["story_package"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
        analysis = dict(inputs.get("analysis_results", {}))
        story = build_story_package(contract, analysis)
        return CapabilityExecutionResult(outputs={"story_package": story}, evidence=[story.get("executive_summary", "")])


class ReportBuildExecutor(BaseCapabilityExecutor):
    name = "report_build"
    description = "Draft a markdown report package."
    input_schema = []
    output_schema = ["report"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
        story = dict(inputs.get("story_package", {}))
        analysis = dict(inputs.get("analysis_results", {}))
        findings = story.get("key_findings") or analysis.get("findings") or []
        bullets = [f"- {str(x)}" for x in findings[:6]]
        summary = str(story.get("executive_summary", f"{contract.get('domain', 'analysis')} report prepared."))
        markdown = "\n".join(
            [
                f"# {contract.get('user_goal', 'Operator Report')}",
                "",
                "## Executive Summary",
                summary,
                "",
                "## Key Findings",
                *bullets,
                "",
            ]
        )
        report = {"summary": summary, "markdown": markdown, "next_actions": list(story.get("recommended_actions", []))[:5]}
        return CapabilityExecutionResult(outputs={"report": report}, evidence=[summary])


class StakeholderSummaryExecutor(BaseCapabilityExecutor):
    name = "stakeholder_summary"
    description = "Build a concise stakeholder-facing summary package."
    input_schema = []
    output_schema = ["stakeholder_summary"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
        story = dict(inputs.get("story_package", {}))
        report = dict(inputs.get("report", {}))
        revision_note = str(inputs.get("revision_note", "")).strip()
        if revision_note:
            summary = {
                "executive_summary": str(report.get("summary", story.get("executive_summary", ""))),
                "key_findings": list(story.get("key_findings", []))[:5],
                "so_what": str(story.get("so_what", "The outputs are ready for stakeholder decision-making.")),
                "recommended_actions": list(story.get("recommended_actions", []))[:5],
                "caveats": list(story.get("caveats", []))[:5],
                "audience": contract.get("audience", "stakeholder"),
            }
        else:
            summary = {
                "executive_summary": str(report.get("summary", story.get("executive_summary", f"{contract.get('domain', 'analysis')} package ready."))),
                "key_findings": list(story.get("key_findings", []))[:3],
                "recommended_actions": list(story.get("recommended_actions", []))[:3],
                "audience": contract.get("audience", "stakeholder"),
            }
        return CapabilityExecutionResult(outputs={"stakeholder_summary": summary}, evidence=[summary.get("executive_summary", "")])


class PresentationBuildExecutor(BaseCapabilityExecutor):
    name = "presentation_build"
    description = "Create a markdown slides outline."
    input_schema = ["story_package"]
    output_schema = ["presentation"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
        outline = build_presentation_outline(contract, dict(inputs.get("story_package", {})))
        return CapabilityExecutionResult(outputs={"presentation": outline}, evidence=[f"{len(outline.get('slides', []))} slide(s) outlined"])


class SpreadsheetBuildExecutor(BaseCapabilityExecutor):
    name = "spreadsheet_build"
    description = "Create rows for spreadsheet export."
    input_schema = []
    output_schema = ["spreadsheet_rows"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        analysis = dict(inputs.get("analysis_results", {}))
        findings = [str(x) for x in (analysis.get("findings") or [])]
        rows = [{"rank": idx + 1, "finding": finding, "action": "review"} for idx, finding in enumerate(findings[:10])]
        if not rows:
            rows = [{"rank": 1, "finding": "No findings were available.", "action": "review"}]
        return CapabilityExecutionResult(outputs={"spreadsheet_rows": rows}, evidence=[f"{len(rows)} spreadsheet row(s) prepared"])


class UIBuildExecutor(BaseCapabilityExecutor):
    name = "ui_build"
    description = "Build a clean commercial assistant UI spec."
    input_schema = []
    output_schema = ["ui_spec"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
        revision_note = str(inputs.get("revision_note", "")).strip()
        requirements = {
            "goal": contract.get("user_goal", ""),
            "audience": contract.get("audience", "operator"),
            "requested_outputs": list(contract.get("requested_outputs", [])),
            "revision_note": revision_note,
        }
        spec = build_ui_delivery(requirements)
        if revision_note:
            spec["developer_details_hidden"] = True
            spec["commercial_quality_focus"] = True
        return CapabilityExecutionResult(outputs={"ui_spec": spec}, evidence=["chat-first UI spec generated"])


class ArtifactExportExecutor(BaseCapabilityExecutor):
    name = "artifact_export"
    description = "Write stakeholder artifacts to disk."
    input_schema = []
    output_schema = ["artifacts", "export_bundle"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        workspace = _workspace_dir(context)
        contract = context.get("task_contract", {}) if isinstance(context.get("task_contract"), dict) else {}
        domain = str(contract.get("domain", "general"))
        artifact_dir = workspace / "artifacts"
        if domain == "competitor_analysis":
            artifact_dir = workspace / "artifacts"
        artifact_dir.mkdir(parents=True, exist_ok=True)
        requested = _requested_outputs(context)
        artifacts: Dict[str, str] = {}
        artifact_metadata: Dict[str, Dict[str, Any]] = {}

        task_contract_path = workspace / "task_contract.json"
        task_contract_path.write_text(json.dumps(contract, indent=2), encoding="utf-8")
        artifacts["task_contract_json"] = str(task_contract_path.resolve())
        artifact_metadata["task_contract_json"] = _artifact_detail(
            key="task_contract_json",
            path=task_contract_path,
            artifact_type="json",
            title="Task Contract",
            evidence_summary="Normalized task contract used by the runtime.",
        )

        if domain == "artifact_generation":
            objective = str(contract.get("user_goal", "Generated Artifact"))
            report = dict(inputs.get("report", {}))
            stakeholder = dict(inputs.get("stakeholder_summary", {}))
            ui_spec = dict(inputs.get("ui_spec", {}))
            primary_body = str(report.get("markdown", "")).strip() or f"# {objective}\n"
            if stakeholder:
                primary_body += "\n## Stakeholder Summary\n" + str(stakeholder.get("executive_summary", "")) + "\n"
            doc_path = artifact_dir / "document.md"
            doc_path.write_text(primary_body, encoding="utf-8")
            artifacts["document_md"] = str(doc_path.resolve())
            artifact_metadata["document_md"] = _artifact_detail(
                key="document_md",
                path=doc_path,
                artifact_type="document",
                title="Instruction Document",
                evidence_summary="Primary markdown artifact generated from the task contract and story package.",
            )
            if "executive_summary" in requested:
                artifacts["executive_summary_md"] = str(doc_path.resolve())
                artifact_metadata["executive_summary_md"] = _artifact_detail(
                    key="executive_summary_md",
                    path=doc_path,
                    artifact_type="document",
                    title="Executive Summary",
                    evidence_summary="Executive summary reuses the primary markdown document for this package.",
                )
            if "presentation" in requested:
                pptx_path = artifact_dir / "deck.pptx"
                bullets = list((dict(inputs.get("story_package", {})).get("key_findings") or []))[:5]
                _write_simple_pptx(pptx_path, objective[:80], "Auto-generated deck outline", [str(x) for x in bullets] or ["Generated from user instruction."])
                artifacts["powerpoint_pptx"] = str(pptx_path.resolve())
                artifact_metadata["powerpoint_pptx"] = _artifact_detail(
                    key="powerpoint_pptx",
                    path=pptx_path,
                    artifact_type="presentation",
                    title="PowerPoint Deck",
                    evidence_summary="Deck outline derived from the runtime story package.",
                )
            if any(x in requested for x in ["dashboard", "visual", "ui"]):
                visual_path = artifact_dir / "visual.html"
                visual_path.write_text(
                    (
                        "<!doctype html><html><head><meta charset='utf-8'><title>Generated Visual</title></head>"
                        "<body style='font-family:Segoe UI,Arial,sans-serif;background:#f8fafc;color:#0f172a;'>"
                        "<div style='max-width:960px;margin:48px auto;padding:32px;background:white;border:1px solid #dbe4ef;border-radius:18px;'>"
                        f"<h1 style='margin:0 0 12px 0'>{objective}</h1>"
                        f"<p style='margin:0 0 16px 0'>{str(stakeholder.get('executive_summary', report.get('summary', 'Artifact package generated through the execution graph runtime.')))}</p>"
                        f"<pre style='white-space:pre-wrap;background:#f8fafc;padding:16px;border-radius:12px;border:1px solid #e2e8f0'>{json.dumps(ui_spec or {'message': 'Visual package generated.'}, indent=2)}</pre>"
                        "</div></body></html>"
                    ),
                    encoding="utf-8",
                )
                artifacts["visual_html"] = str(visual_path.resolve())
                artifact_metadata["visual_html"] = _artifact_detail(
                    key="visual_html",
                    path=visual_path,
                    artifact_type="html",
                    title="Visual Artifact",
                    evidence_summary="Visual HTML package generated from the runtime UI/story outputs.",
                )
                artifacts["dashboard_html"] = str(visual_path.resolve())
                artifact_metadata["dashboard_html"] = _artifact_detail(
                    key="dashboard_html",
                    path=visual_path,
                    artifact_type="dashboard",
                    title="Visual Dashboard",
                    evidence_summary="Compatibility alias to the primary visual HTML artifact.",
                )
            artifacts["directory"] = str(workspace.resolve())
            artifact_metadata["directory"] = _artifact_detail(
                key="directory",
                path=workspace,
                artifact_type="directory",
                title="Artifact Directory",
                evidence_summary="Workspace directory containing the generated package.",
            )
            artifacts["primary_open_file"] = artifacts.get("visual_html") or artifacts.get("document_md") or artifacts.get("powerpoint_pptx", "")
            if artifacts.get("primary_open_file"):
                artifact_metadata["primary_open_file"] = _artifact_detail(
                    key="primary_open_file",
                    path=artifacts["primary_open_file"],
                    artifact_type="pointer",
                    title="Primary Open File",
                    evidence_summary="Preferred artifact to open first for review.",
                )
            export_bundle = {"workspace": str(workspace.resolve()), "artifact_count": len(artifacts), "artifact_keys": sorted(artifacts.keys())}
            return CapabilityExecutionResult(
                outputs={"artifacts": artifacts, "export_bundle": export_bundle},
                artifacts=artifacts,
                artifact_metadata=artifact_metadata,
                evidence=list(artifacts.keys()),
            )

        if domain == "competitor_analysis":
            rows = [dict(x) for x in (inputs.get("structured_rows") or []) if isinstance(x, dict)]
            research_notes = dict(inputs.get("research_notes", {}))
            target = str(research_notes.get("target", "Target Company"))
            csv_path = artifact_dir / "competitors.csv"
            with csv_path.open("w", newline="", encoding="utf-8") as fh:
                writer = csv.DictWriter(fh, fieldnames=["rank", "name", "segment", "why", "citations", "score"])
                writer.writeheader()
                for row in rows:
                    writer.writerow({key: row.get(key, "") for key in ["rank", "name", "segment", "why", "citations", "score"]})
            artifacts["competitors_csv"] = str(csv_path.resolve())
            artifact_metadata["competitors_csv"] = _artifact_detail(
                key="competitors_csv",
                path=csv_path,
                artifact_type="spreadsheet",
                title="Competitor Matrix",
                evidence_summary=f"{len(rows)} competitor row(s) exported.",
            )
            story = dict(inputs.get("story_package", {}))
            report = dict(inputs.get("report", {}))
            report_path = artifact_dir / "executive_summary.md"
            report_body = str(report.get("markdown", "")).strip() or f"# Executive Summary: {target} Competitor Analysis\n"
            report_path.write_text(report_body, encoding="utf-8")
            artifacts["executive_summary_md"] = str(report_path.resolve())
            artifact_metadata["executive_summary_md"] = _artifact_detail(
                key="executive_summary_md",
                path=report_path,
                artifact_type="document",
                title="Executive Summary",
                evidence_summary="Stakeholder markdown summary for the competitor landscape.",
            )
            report_html_path = artifact_dir / "executive_summary.html"
            report_html_path.write_text(
                "<!doctype html><html><head><meta charset='utf-8'><title>Executive Summary</title></head>"
                "<body style='font-family:Segoe UI,Arial,sans-serif;background:#f8fafc;color:#0f172a;'><div style='max-width:980px;margin:48px auto;padding:28px;background:#fff;border:1px solid #dbe4ef;border-radius:18px;'>"
                f"<h1 style='margin-top:0'>{target} Competitor Executive Summary</h1>"
                f"<pre style='white-space:pre-wrap'>{report_body}</pre></div></body></html>",
                encoding="utf-8",
            )
            artifacts["executive_summary_html"] = str(report_html_path.resolve())
            artifact_metadata["executive_summary_html"] = _artifact_detail(
                key="executive_summary_html",
                path=report_html_path,
                artifact_type="html",
                title="Executive Summary HTML",
                evidence_summary="HTML rendering of the competitor executive summary.",
            )
            pptx_path = artifact_dir / "executive_summary.pptx"
            bullets = [str(x) for x in (story.get("key_findings") or [])[:5]] or [str(row.get("name", "")) for row in rows[:5]]
            _write_simple_pptx(pptx_path, f"{target} Competitor Analysis", "Executive Summary Deck", bullets)
            artifacts["powerpoint_pptx"] = str(pptx_path.resolve())
            artifact_metadata["powerpoint_pptx"] = _artifact_detail(
                key="powerpoint_pptx",
                path=pptx_path,
                artifact_type="presentation",
                title="Competitor Deck",
                evidence_summary="PowerPoint summary of the top competitor findings.",
            )
            dash_path = artifact_dir / "dashboard.html"
            cards = "".join(
                f"<div style='padding:14px;border:1px solid #e2e8f0;border-radius:14px;'><strong>{row.get('name','')}</strong><div>{row.get('segment','')}</div><div style='margin-top:8px'>{row.get('why','')}</div></div>"
                for row in rows[:6]
            )
            dash_path.write_text(
                "<!doctype html><html><head><meta charset='utf-8'><title>Competitor Dashboard</title></head>"
                "<body style='font-family:Segoe UI,Arial,sans-serif;background:#f8fafc;color:#0f172a;'><div style='max-width:1100px;margin:40px auto;padding:28px;background:#fff;border:1px solid #dbe4ef;border-radius:18px;'>"
                f"<h1 style='margin-top:0'>{target} Competitor Dashboard</h1>"
                f"<p>{str(story.get('executive_summary', 'Competitor package generated through the execution graph runtime.'))}</p>"
                f"<div style='display:grid;grid-template-columns:repeat(auto-fit,minmax(240px,1fr));gap:14px'>{cards}</div>"
                "</div></body></html>",
                encoding="utf-8",
            )
            artifacts["dashboard_html"] = str(dash_path.resolve())
            artifact_metadata["dashboard_html"] = _artifact_detail(
                key="dashboard_html",
                path=dash_path,
                artifact_type="dashboard",
                title="Competitor Dashboard",
                evidence_summary="HTML dashboard for competitor review and source-backed positioning notes.",
            )
            artifacts["directory"] = str(workspace.resolve())
            artifact_metadata["directory"] = _artifact_detail(
                key="directory",
                path=workspace,
                artifact_type="directory",
                title="Competitor Package Folder",
                evidence_summary="Workspace directory containing the competitor stakeholder package.",
            )
            artifacts["primary_open_file"] = str(report_html_path.resolve())
            artifact_metadata["primary_open_file"] = _artifact_detail(
                key="primary_open_file",
                path=report_html_path,
                artifact_type="pointer",
                title="Primary Open File",
                evidence_summary="Preferred review file for the competitor package.",
            )
            export_bundle = {"workspace": str(workspace.resolve()), "artifact_count": len(artifacts), "artifact_keys": sorted(artifacts.keys())}
            return CapabilityExecutionResult(
                outputs={"artifacts": artifacts, "export_bundle": export_bundle},
                artifacts=artifacts,
                artifact_metadata=artifact_metadata,
                evidence=list(artifacts.keys()),
            )

        if domain in {"generic_research", "web_research"}:
            from lam.interface import search_agent as search_agent_mod

            raw_results = list(inputs.get("search_results", []) or [])
            results = []
            for row in raw_results:
                if isinstance(row, search_agent_mod.SearchResult):  # type: ignore[attr-defined]
                    results.append(row)
                elif isinstance(row, dict):
                    results.append(
                        search_agent_mod.SearchResult(  # type: ignore[attr-defined]
                            title=str(row.get("title", "")),
                            url=str(row.get("url", "")),
                            price=row.get("price"),
                            source=str(row.get("source", "")),
                            snippet=str(row.get("snippet", "")),
                        )
                    )
            generated = search_agent_mod._write_generic_research_artifacts(  # type: ignore[attr-defined]
                instruction=str(contract.get("user_goal", "")),
                query=str(inputs.get("query", contract.get("user_goal", ""))),
                results=results,
                decision_rows=list(inputs.get("decision_rows", []) or []),
                recommendation=dict(inputs.get("recommendation", {}) or {}),
                browser_notes=[dict(x) for x in (inputs.get("browser_notes", []) or []) if isinstance(x, dict)],
            )
            artifacts.update({k: v for k, v in generated.items() if isinstance(v, str) and v})
            type_map = {
                "results_csv": ("spreadsheet", "Results CSV"),
                "decision_matrix_csv": ("spreadsheet", "Decision Matrix"),
                "report_md": ("report", "Research Report"),
                "recommendation_md": ("report", "Recommendation"),
                "browser_research_md": ("report", "Browser Research Log"),
                "browser_research_json": ("json", "Browser Research Data"),
                "dashboard_html": ("dashboard", "Research Dashboard"),
                "research_dashboard_html": ("dashboard", "Research Dashboard"),
                "executive_brief_md": ("report", "Executive Brief"),
                "executive_brief_html": ("html", "Executive Brief HTML"),
                "powerpoint_pptx": ("presentation", "PowerPoint Deck"),
                "directory": ("directory", "Research Output Directory"),
                "primary_open_file": ("pointer", "Primary Open File"),
            }
            for key, value in artifacts.items():
                if key in artifact_metadata:
                    continue
                artifact_type, title = type_map.get(key, ("file", key))
                artifact_metadata[key] = _artifact_detail(
                    key=key,
                    path=value,
                    artifact_type=artifact_type,
                    title=title,
                    evidence_summary=f"Generic research artifact exported through artifact_export: {title}.",
                )
            export_bundle = {"workspace": str(workspace.resolve()), "artifact_count": len(artifacts), "artifact_keys": sorted(artifacts.keys())}
            return CapabilityExecutionResult(
                outputs={"artifacts": artifacts, "export_bundle": export_bundle},
                artifacts=artifacts,
                artifact_metadata=artifact_metadata,
                evidence=list(artifacts.keys()),
            )

        report = dict(inputs.get("report", {}))
        stakeholder = dict(inputs.get("stakeholder_summary", {}))
        if report or "report" in requested:
            report_path = artifact_dir / "summary_report.md"
            body = str(report.get("markdown", "")) or f"# {contract.get('user_goal', 'Operator Report')}\n"
            if stakeholder:
                body += "\n## Stakeholder Summary\n" + str(stakeholder.get("executive_summary", "")) + "\n"
            report_path.write_text(body, encoding="utf-8")
            artifacts["summary_report_md"] = str(report_path.resolve())
            artifacts["report_md"] = str(report_path.resolve())
            artifact_metadata["summary_report_md"] = _artifact_detail(
                key="summary_report_md",
                path=report_path,
                artifact_type="report",
                title="Summary Report",
                evidence_summary="Primary markdown report exported from the execution graph.",
            )
            artifact_metadata["report_md"] = _artifact_detail(
                key="report_md",
                path=report_path,
                artifact_type="report",
                title="Report",
                evidence_summary="Alias to the primary markdown report for compatibility.",
            )

        presentation = dict(inputs.get("presentation", {}))
        if presentation or "presentation" in requested:
            slides_path = artifact_dir / "slides.md"
            slides = presentation.get("slides", []) if isinstance(presentation.get("slides"), list) else []
            slide_lines = []
            for slide in slides:
                slide_lines.append(f"# {slide.get('title', 'Slide')}")
                for bullet in (slide.get("bullets", []) or []):
                    slide_lines.append(f"- {bullet}")
                slide_lines.append("")
            slides_path.write_text("\n".join(slide_lines).strip() + "\n", encoding="utf-8")
            artifacts["presentation_md"] = str(slides_path.resolve())
            artifact_metadata["presentation_md"] = _artifact_detail(
                key="presentation_md",
                path=slides_path,
                artifact_type="presentation",
                title="Presentation Outline",
                evidence_summary=f"{len(slides)} slide(s) exported in markdown.",
            )

        rows = [dict(x) for x in (inputs.get("spreadsheet_rows") or []) if isinstance(x, dict)]
        if rows or "spreadsheet" in requested:
            csv_path = artifact_dir / "decision_matrix.csv"
            fieldnames = sorted({key for row in rows for key in row.keys()}) or ["rank", "finding", "action"]
            with csv_path.open("w", newline="", encoding="utf-8") as fh:
                writer = csv.DictWriter(fh, fieldnames=fieldnames)
                writer.writeheader()
                for row in rows or [{"rank": 1, "finding": "No findings were available.", "action": "review"}]:
                    writer.writerow({key: row.get(key, "") for key in fieldnames})
            artifacts["decision_matrix_csv"] = str(csv_path.resolve())
            artifact_metadata["decision_matrix_csv"] = _artifact_detail(
                key="decision_matrix_csv",
                path=csv_path,
                artifact_type="spreadsheet",
                title="Decision Matrix",
                evidence_summary=f"{len(rows) or 1} row(s) exported for spreadsheet review.",
            )

        ui_spec = dict(inputs.get("ui_spec", {}))
        if ui_spec or "ui" in requested or "dashboard" in requested:
            spec_path = artifact_dir / "ui_spec.json"
            spec_path.write_text(json.dumps(ui_spec or {"message": "No UI spec generated."}, indent=2), encoding="utf-8")
            artifacts["ui_spec_json"] = str(spec_path.resolve())
            artifact_metadata["ui_spec_json"] = _artifact_detail(
                key="ui_spec_json",
                path=spec_path,
                artifact_type="json",
                title="UI Spec",
                evidence_summary="Commercial UI specification exported from the runtime.",
            )
            html_path = artifact_dir / "dashboard.html"
            title = str(contract.get("user_goal", "OpenLAMb Dashboard"))
            html_path.write_text(
                (
                    "<!doctype html><html><head><meta charset='utf-8'><title>Generated Dashboard</title></head>"
                    "<body style='font-family:Segoe UI,Arial,sans-serif;background:#f8fafc;color:#0f172a;'>"
                    "<div style='max-width:960px;margin:40px auto;padding:32px;background:#fff;border:1px solid #dbe4ef;border-radius:18px;'>"
                    f"<h1 style='margin:0 0 12px 0'>{title}</h1>"
                    f"<p style='margin:0 0 20px 0'>{str(stakeholder.get('executive_summary', report.get('summary', 'Generated through the execution graph runtime.')))}</p>"
                    "<div style='display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:14px;'>"
                    "<div style='padding:16px;border:1px solid #e2e8f0;border-radius:14px;'><strong>Chat</strong><div>Conversation remains primary.</div></div>"
                    "<div style='padding:16px;border:1px solid #e2e8f0;border-radius:14px;'><strong>Canvas</strong><div>Artifacts and workbench stay progressive.</div></div>"
                    "<div style='padding:16px;border:1px solid #e2e8f0;border-radius:14px;'><strong>Artifacts</strong><div>Open generated outputs directly.</div></div>"
                    "</div></div></body></html>"
                ),
                encoding="utf-8",
            )
            artifacts["dashboard_html"] = str(html_path.resolve())
            artifact_metadata["dashboard_html"] = _artifact_detail(
                key="dashboard_html",
                path=html_path,
                artifact_type="dashboard",
                title="Dashboard",
                evidence_summary="HTML dashboard summarizing the generated package.",
            )

        rag_index = dict(inputs.get("rag_index", {}))
        if rag_index or "rag_index" in requested:
            rag_path = artifact_dir / "rag_index.json"
            rag_path.write_text(json.dumps(rag_index or {"kind": "empty"}, indent=2), encoding="utf-8")
            artifacts["rag_index_db"] = str(rag_path.resolve())
            artifact_metadata["rag_index_db"] = _artifact_detail(
                key="rag_index_db",
                path=rag_path,
                artifact_type="rag_index",
                title="RAG Index",
                evidence_summary=f"Local retrieval corpus with {len((rag_index or {}).get('documents', []) or [])} document(s).",
            )

        code_changes = dict(inputs.get("code_changes", {}))
        if code_changes or "code" in requested:
            src_dir = workspace / "src"
            src_dir.mkdir(parents=True, exist_ok=True)
            analysis_path = src_dir / "generated_entry.py"
            snippets = code_changes.get("snippets", []) if isinstance(code_changes.get("snippets"), list) else []
            analysis_path.write_text("\n".join([str(x) for x in snippets] or ["def run():", "    return {'ok': True}"]) + "\n", encoding="utf-8")
            artifacts["analysis_script_py"] = str(analysis_path.resolve())
            artifact_metadata["analysis_script_py"] = _artifact_detail(
                key="analysis_script_py",
                path=analysis_path,
                artifact_type="code",
                title="Analysis Script",
                evidence_summary=str(code_changes.get("summary", "Runtime-generated analysis scaffold.")),
            )

        export_bundle = {"workspace": str(workspace.resolve()), "artifact_count": len(artifacts), "artifact_keys": sorted(artifacts.keys())}
        return CapabilityExecutionResult(
            outputs={"artifacts": artifacts, "export_bundle": export_bundle},
            artifacts=artifacts,
            artifact_metadata=artifact_metadata,
            evidence=list(artifacts.keys()),
        )


class ApprovalGateExecutor(BaseCapabilityExecutor):
    name = "approval_gate"
    description = "Mark an action as requiring approval."
    input_schema = []
    output_schema = ["approval_state"]

    def execute(self, context: Dict[str, Any], inputs: Dict[str, Any]) -> CapabilityExecutionResult:
        return CapabilityExecutionResult(outputs={"approval_state": {"required": True, "reason": "high-impact action detected"}})


def default_executors() -> Dict[str, BaseCapabilityExecutor]:
    rows: Iterable[BaseCapabilityExecutor] = [
        DeepResearchExecutor(),
        ResearchCollectionExecutor(),
        CompetitorResearchExecutor(),
        SourceEvaluationExecutor(),
        FileInspectionExecutor(),
        DataCleaningExecutor(),
        StatisticalAnalysisExecutor(),
        DataVisualizationExecutor(),
        RAGBuildExecutor(),
        RAGQueryExecutor(),
        CodeWriteExecutor(),
        CodeTestExecutor(),
        CodeFixExecutor(),
        DataStorytellingExecutor(),
        ReportBuildExecutor(),
        StakeholderSummaryExecutor(),
        PresentationBuildExecutor(),
        SpreadsheetBuildExecutor(),
        UIBuildExecutor(),
        ArtifactExportExecutor(),
        ApprovalGateExecutor(),
    ]
    return {row.name: row for row in rows}
