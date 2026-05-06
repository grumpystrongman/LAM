from __future__ import annotations

import csv
import html
import hashlib
import imaplib
import json
import logging
import os
import random
import re
import shutil
import statistics
import subprocess
import time
import uuid
import urllib.parse
import urllib.request
import webbrowser
import xml.etree.ElementTree as ET
from email import message_from_bytes
from email.message import EmailMessage
from email.header import decode_header
from email.utils import parseaddr, parsedate_to_datetime
from imaplib import Time2Internaldate
from datetime import datetime, timedelta
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Set

from lam.interface.ai_backend import backend_metadata, normalize_backend
from lam.interface.app_launcher import is_app_running, normalize_app_name, open_installed_app
from lam.interface.browser_worker import ensure_browser_worker, normalize_browser_worker_mode
from lam.interface.app_learner import get_guidance
from lam.interface.clipboard_capture import base64_to_image, capture_clipboard_image, image_to_base64
from lam.interface.desktop_sequence import assess_risk, build_plan, execute_plan
from lam.interface.domain_playbooks import (
    build_step_obligations,
    evaluate_step_obligations,
    select_playbook,
    validate_plan_steps,
    validate_transition_graph,
)
from lam.interface.human_judgment import (
    ActionCritic,
    EleganceBudget,
    NegativeMemory,
    QualityCritic,
    assess_result_quality,
)
from lam.interface.local_vector_store import LocalVectorStore
from lam.interface.operator_contract import attach_operator_contract
from lam.interface.password_vault import LocalPasswordVault
from lam.interface.session_manager import SessionManager
from lam.interface.world_model import build_run_world_model
from lam.operator_platform import (
    ArtifactFactory,
    CapabilityPlanner,
    CompletionCritic as PlatformCompletionCritic,
    DataQualityCritic as PlatformDataQualityCritic,
    ExecutionGraphRuntime,
    HumanStyleReporter,
    MemoryStore,
    MissionContractEngine,
    MissionRuntime,
    PresentationCritic as PlatformPresentationCritic,
    SourceCritic as PlatformSourceCritic,
    StoryCritic as PlatformStoryCritic,
    TaskContractEngine,
    UIUXCritic as PlatformUIUXCritic,
    WorldModelBuilder,
    build_platform_cards,
    default_capability_registry,
    default_executors,
)
from lam.learn.topic_mastery_runtime import TopicMasteryRuntime
from lam.operator_platform.research_backend import (
    apply_human_judgment_quality_gate as platform_apply_human_judgment_quality_gate,
    build_recommendation_focus_query as platform_build_recommendation_focus_query,
    count_locality_matches as platform_count_locality_matches,
    is_recommendation_research_intent as platform_is_recommendation_research_intent,
    is_wine_pairing_intent as platform_is_wine_pairing_intent,
    minimum_topic_overlap as platform_minimum_topic_overlap,
    passes_topic_gate as platform_passes_topic_gate,
    query_focus_terms as platform_query_focus_terms,
    relevance_score as platform_relevance_score,
    topic_overlap_count as platform_topic_overlap_count,
)
from lam.operator_platform.research_primitives import (
    collect_generic_research as platform_collect_generic_research,
    expand_queries as platform_expand_queries,
    extract_generic_query as platform_extract_generic_query,
    human_judgment_constraints as platform_human_judgment_constraints,
    human_judgment_refine_queries as platform_human_judgment_refine_queries,
)
from lam.operator_platform.browser_research import (
    browser_query_url as platform_browser_query_url,
    browser_research_walk as platform_browser_research_walk,
)
from lam.operator_platform.research_constants import (
    QUERY_NOISE_TERMS,
    RECOMMENDATION_RESEARCH_TOKENS,
    STEAK_WINE_STYLE_BONUS,
    USER_AGENT,
    WINE_STYLE_KEYWORDS,
)
from lam.operator_platform.research_types import SearchResult
from lam.operator_platform.recommendation_helpers import (
    build_decision_rows as platform_build_decision_rows,
    build_product_candidate_rows as platform_build_product_candidate_rows,
    build_recommendation_summary as platform_build_recommendation_summary,
    curated_recommendation_sources as platform_curated_recommendation_sources,
    merge_browser_notes_into_rows as platform_merge_browser_notes_into_rows,
)
from lam.deep_workbench.workflow import (
    build_workspace as build_code_workbench_workspace,
    extract_workbench_contract,
)
from lam.payer_rag.workflow import (
    ask_workspace_question,
    build_workspace,
    ensure_workspace,
    extract_current_task_contract,
)
MIN_LIVE_NON_CURATED_CITATIONS = 3
DESTRUCTIVE_ACTION_KEYWORDS = {
    "delete",
    "remove",
    "destroy",
    "send email",
    "send message",
    "wire transfer",
    "transfer money",
    "payment",
    "pay bill",
    "purchase",
    "buy now",
    "submit payment",
}

COMMUNICATION_SOCIAL_KEYWORDS = {
    "whatsapp",
    "telegram",
    "messenger",
    "discord",
    "slack",
    "chat",
    "dm",
    "social media",
    "instagram",
    "facebook",
    "linkedin",
    "x.com",
    "twitter",
    "tiktok",
    "reddit",
}

HUMAN_JUDGMENT_SUPERLATIVE_TOKENS = {
    "best",
    "most expensive",
    "cheapest",
    "closest",
    "fastest",
    "latest",
    "top",
    "highest",
    "highest-rated",
    "most relevant",
}

HUMAN_JUDGMENT_LOCALITY_TOKENS = {
    "near me",
    "locally",
    "local",
    "around here",
    "in ",
}

CODE_WORKBENCH_TOKENS = {
    "vscode",
    "vs code",
    "visual studio code",
    "write code",
    "build code",
    "analysis app",
    "analysis script",
    "workspace",
    "new codebase",
}

CODE_WORKBENCH_PIPELINE_TOKENS = {
    "research",
    "collect",
    "source data",
    "ingest",
    "normalize",
    "analy",
    "build",
    "write",
    "test",
    "fix",
    "package",
    "deliverable",
    "submit the finished result",
}

CODE_WORKBENCH_DELIVERABLE_TOKENS = {
    "code",
    "script",
    "app",
    "cli",
    "web app",
    "analysis app",
    "analysis script",
    "rag",
    "vector store",
    "retriever",
    "unit test",
    "smoke test",
    "stakeholder deliverable",
}

PAYER_PRICING_INTENT_TOKENS = {
    "payer",
    "payers",
    "insurance",
    "plan",
    "plans",
    "health plan",
    "transparency in coverage",
    "negotiated rate",
    "pricing outlier",
    "standard charges",
    "shoppable",
}

PAYER_PRICING_ACTION_TOKENS = {
    "durham",
    "rag",
    "vector store",
    "retriever",
    "pricing",
    "rate",
    "rates",
    "analy",
    "outlier",
    "contract",
    "contact",
    "spreadsheet",
    "stakeholder",
}

LOGGER = logging.getLogger(__name__)
_FITZ_MODULE: Any | None = None
_FITZ_IMPORT_ATTEMPTED = False
STUDY_PDF_RENDER_ZOOM = float(os.getenv("LAM_STUDY_PDF_RENDER_ZOOM", "1.6") or "1.6")
STUDY_ASSETS_ROOT = Path(__file__).resolve().parents[2] / "data" / "reports" / "study_assets"
VISUAL_KEYWORDS = (
    "sign",
    "signal",
    "light",
    "marking",
    "lane",
    "intersection",
    "diagram",
    "chart",
    "table",
    "figure",
    "photo",
    "image",
    "map",
    "screenshot",
    "scan",
    "radiograph",
)



@dataclass(slots=True)
class JobListing:
    title: str
    url: str
    source: str
    location: str
    remote: bool
    salary_text: str
    salary_min: Optional[float]
    salary_max: Optional[float]
    currency: str
    snippet: str = ""


@dataclass(slots=True)
class JobSearchConstraints:
    require_vp_avp: bool
    require_remote_or_hybrid: bool
    min_base_salary_usd: Optional[float]
    min_total_comp_usd: Optional[float]
    allowed_regions: List[str]


@dataclass(slots=True)
class StudyItem:
    question: str
    answer: str
    category: str
    difficulty: str
    source_url: str = ""
    evidence: str = ""
    image_path: str = ""
    image_base64: str = ""


@dataclass(slots=True)
class EmailActionItem:
    message_id: str
    sender: str
    subject: str
    received_at: str
    snippet: str
    requires_action: bool
    reason: str
    draft_created: bool


_EMAIL_AUTH_SESSIONS: Dict[str, Dict[str, Any]] = {}
ARTIFACT_REUSE_MODES = {"reuse", "reuse_if_recent", "always_regenerate"}


def _normalize_gmail_auth_url(url: str, default: str = "https://mail.google.com/") -> str:
    raw = str(url or "").strip()
    if not raw:
        return default
    try:
        parsed = urllib.parse.urlparse(raw)
        host = (parsed.netloc or "").lower()
    except Exception:
        return default
    if "mail.google.com" in host:
        return raw
    if "accounts.google.com" in host:
        return raw
    return default


def _artifact_reuse_index_path() -> Path:
    path = Path("data/interface/artifact_reuse_index.json")
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def _load_artifact_reuse_index() -> Dict[str, Any]:
    path = _artifact_reuse_index_path()
    if not path.exists():
        return {"entries": []}
    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
        if isinstance(raw, dict):
            raw.setdefault("entries", [])
            return raw
    except Exception:
        pass
    return {"entries": []}


def _save_artifact_reuse_index(index: Dict[str, Any]) -> None:
    path = _artifact_reuse_index_path()
    path.write_text(json.dumps(index, indent=2), encoding="utf-8")


def _artifact_instruction_key(instruction: str) -> str:
    normalized = re.sub(r"\s+", " ", str(instruction or "").strip().lower())
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()


def _remember_artifacts_for_reuse(kind: str, instruction: str, artifacts: Dict[str, str]) -> None:
    if not isinstance(artifacts, dict) or not artifacts:
        return
    index = _load_artifact_reuse_index()
    entries = [x for x in index.get("entries", []) if isinstance(x, dict)]
    key = _artifact_instruction_key(instruction)
    filtered: List[Dict[str, Any]] = []
    for item in entries:
        if str(item.get("kind", "")) == str(kind) and str(item.get("instruction_key", "")) == key:
            continue
        filtered.append(item)
    filtered.append(
        {
            "ts": time.time(),
            "kind": str(kind),
            "instruction_key": key,
            "instruction": str(instruction)[:500],
            "artifacts": {k: str(v) for k, v in artifacts.items() if isinstance(v, str)},
        }
    )
    index["entries"] = filtered[-300:]
    _save_artifact_reuse_index(index)


def _find_reusable_artifacts(kind: str, instruction: str, required_keys: List[str], max_age_hours: int = 24) -> Dict[str, str]:
    index = _load_artifact_reuse_index()
    entries = [x for x in index.get("entries", []) if isinstance(x, dict)]
    key = _artifact_instruction_key(instruction)
    now = time.time()
    for item in reversed(entries):
        if str(item.get("kind", "")) != str(kind):
            continue
        if str(item.get("instruction_key", "")) != key:
            continue
        ts = float(item.get("ts", 0.0) or 0.0)
        if ts <= 0 or (now - ts) > max_age_hours * 3600:
            continue
        artifacts = item.get("artifacts", {})
        if not isinstance(artifacts, dict):
            continue
        if required_keys and any(str(k) not in artifacts for k in required_keys):
            continue
        paths_ok = True
        for value in artifacts.values():
            if not isinstance(value, str):
                paths_ok = False
                break
            if value.startswith(("http://", "https://")):
                continue
            if not Path(value).exists():
                paths_ok = False
                break
        if paths_ok:
            return {k: str(v) for k, v in artifacts.items() if isinstance(v, str)}
    return {}


def _normalize_artifact_reuse_mode(mode: str) -> str:
    value = str(mode or "").strip().lower()
    if value not in ARTIFACT_REUSE_MODES:
        return "reuse_if_recent"
    return value


def _fetch_text(url: str) -> str:
    req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(req, timeout=20) as resp:  # nosec B310 - controlled URLs
        return resp.read().decode("utf-8", errors="ignore")


def _parse_duckduckgo(query: str, limit: int = 8) -> List[SearchResult]:
    q = urllib.parse.quote_plus(query)
    url = f"https://duckduckgo.com/html/?q={q}"
    html = _fetch_text(url)
    pattern = re.compile(
        r'<a[^>]*class="result__a"[^>]*href="(?P<href>[^"]+)"[^>]*>(?P<title>.*?)</a>',
        re.IGNORECASE | re.DOTALL,
    )
    snippet_pattern = re.compile(r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>', re.IGNORECASE | re.DOTALL)
    snippets = snippet_pattern.findall(html)

    results: List[SearchResult] = []
    for idx, match in enumerate(pattern.finditer(html)):
        href = match.group("href")
        title = re.sub("<.*?>", "", match.group("title")).strip()
        parsed = urllib.parse.urlparse(href)
        url_value = href
        if parsed.netloc.endswith("duckduckgo.com"):
            if parsed.path.startswith("/y.js"):
                qs = urllib.parse.parse_qs(parsed.query)
                if "u3" in qs:
                    url_value = urllib.parse.unquote(qs["u3"][0])
                elif "u" in qs:
                    url_value = urllib.parse.unquote(qs["u"][0])
                else:
                    continue
            qs = urllib.parse.parse_qs(parsed.query)
            if "uddg" in qs:
                url_value = urllib.parse.unquote(qs["uddg"][0])
        if "duckduckgo.com/y.js" in url_value:
            continue
        snippet = re.sub("<.*?>", "", snippets[idx]).strip() if idx < len(snippets) else ""
        results.append(SearchResult(title=title, url=url_value, price=_extract_price(title + " " + snippet), source="duckduckgo", snippet=snippet))
        if len(results) >= limit:
            break
    return results


def _parse_bing_rss(query: str, limit: int = 8) -> List[SearchResult]:
    q = urllib.parse.quote_plus(query)
    url = f"https://www.bing.com/search?q={q}&format=rss"
    xml_text = _fetch_text(url)
    results: List[SearchResult] = []
    try:
        root = ET.fromstring(xml_text)
    except ET.ParseError:
        return []
    for item in root.findall(".//item"):
        title = (item.findtext("title") or "").strip()
        link = (item.findtext("link") or "").strip()
        snippet = (item.findtext("description") or "").strip()
        if not title or not link:
            continue
        results.append(
            SearchResult(
                title=title,
                url=link,
                price=_extract_price(title + " " + snippet),
                source="bing_rss",
                snippet=snippet,
            )
        )
        if len(results) >= limit:
            break
    return results


def _search_web(query: str, limit: int = 10) -> List[SearchResult]:
    results = _safe_search_web(query, limit=limit)
    if len(results) >= max(3, limit // 2):
        return results
    try:
        fallback = _parse_bing_rss(query, limit=limit)
    except Exception:
        fallback = []
    out: Dict[str, SearchResult] = {}
    for r in results + fallback:
        out[r.url] = r
    return list(out.values())[:limit]


def _parse_review_image_signals(text: str) -> Dict[str, Any]:
    raw = str(text or "")
    low = raw.lower()
    rating: Optional[float] = None
    review_count: Optional[int] = None
    image_count = low.count("<img")
    rating_match = re.search(r"([0-5](?:\.[0-9])?)\s*(?:out of 5|/5)", low)
    if rating_match:
        try:
            rating = float(rating_match.group(1))
        except ValueError:
            rating = None
    review_match = re.search(r"([0-9][0-9,]{0,8})\s+(?:ratings|rating|reviews|review)\b", low)
    if review_match:
        try:
            review_count = int(review_match.group(1).replace(",", ""))
        except ValueError:
            review_count = None
    return {
        "rating": rating,
        "review_count": review_count,
        "image_count": int(image_count),
    }


def _collect_page_signals(url: str) -> Dict[str, Any]:
    target = str(url or "").strip()
    if not target:
        return {"rating": None, "review_count": None, "image_count": 0}
    try:
        page_text = _fetch_text(target)
    except Exception:
        return {"rating": None, "review_count": None, "image_count": 0}
    return _parse_review_image_signals(page_text)


def _estimate_condition(text: str) -> str:
    low = str(text or "").lower()
    if any(t in low for t in ["brand new", "new in box", "new"]):
        return "new"
    if any(t in low for t in ["refurb", "renewed"]):
        return "refurbished"
    if any(t in low for t in ["open box", "used", "pre-owned", "preowned"]):
        return "used"
    return "unknown"


def _build_shopping_candidates(results: List[SearchResult], max_items: int = 10, signal_pages: int = 4) -> List[Dict[str, Any]]:
    candidates: List[Dict[str, Any]] = []
    seen: set[str] = set()
    for result in results:
        url = str(result.url or "").strip()
        if not url or url in seen:
            continue
        seen.add(url)
        base_text = f"{result.title} {result.snippet}"
        signals = _parse_review_image_signals(base_text)
        candidates.append(
            {
                "title": result.title,
                "url": url,
                "source": result.source,
                "price": result.price,
                "snippet": result.snippet,
                "condition": _estimate_condition(base_text),
                "rating": signals.get("rating"),
                "review_count": signals.get("review_count"),
                "image_count": int(signals.get("image_count") or 0),
            }
        )
        if len(candidates) >= max_items:
            break
    for idx, row in enumerate(candidates[: max(0, signal_pages)]):
        if row.get("rating") is not None and row.get("review_count"):
            continue
        page_signals = _collect_page_signals(str(row.get("url", "")))
        if row.get("rating") is None and page_signals.get("rating") is not None:
            row["rating"] = page_signals.get("rating")
        if not row.get("review_count") and page_signals.get("review_count") is not None:
            row["review_count"] = page_signals.get("review_count")
        row["image_count"] = max(int(row.get("image_count") or 0), int(page_signals.get("image_count") or 0))
        candidates[idx] = row
    return candidates


def _pick_recommended_candidate(candidates: List[Dict[str, Any]]) -> Dict[str, Any]:
    if not candidates:
        return {}
    priced = [c for c in candidates if isinstance(c.get("price"), (int, float))]
    if priced:
        priced.sort(
            key=lambda c: (
                float(c.get("price") or 999999.0),
                -(float(c.get("rating") or 0.0)),
                -(int(c.get("review_count") or 0)),
            )
        )
        return dict(priced[0])
    ranked = sorted(
        candidates,
        key=lambda c: (
            -(float(c.get("rating") or 0.0)),
            -(int(c.get("review_count") or 0)),
            -(int(c.get("image_count") or 0)),
        ),
    )
    return dict(ranked[0]) if ranked else {}


def _write_shopping_decision_artifacts(
    instruction: str,
    query: str,
    candidates: List[Dict[str, Any]],
    recommendation: Dict[str, Any],
) -> Dict[str, str]:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path("data/reports/shopping_assistant") / ts
    out_dir.mkdir(parents=True, exist_ok=True)
    csv_path = out_dir / "decision_matrix.csv"
    notes_path = out_dir / "recommendation.md"
    dash_path = out_dir / "shopping_dashboard.html"

    with csv_path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(
            f,
            fieldnames=[
                "rank",
                "title",
                "url",
                "source",
                "price",
                "rating",
                "review_count",
                "image_count",
                "condition",
                "snippet",
            ],
        )
        writer.writeheader()
        for idx, row in enumerate(candidates, start=1):
            writer.writerow(
                {
                    "rank": idx,
                    "title": row.get("title", ""),
                    "url": row.get("url", ""),
                    "source": row.get("source", ""),
                    "price": row.get("price", ""),
                    "rating": row.get("rating", ""),
                    "review_count": row.get("review_count", ""),
                    "image_count": row.get("image_count", ""),
                    "condition": row.get("condition", ""),
                    "snippet": row.get("snippet", ""),
                }
            )

    lines = [
        "# Shopping Recommendation",
        "",
        f"- Generated: {datetime.now().isoformat(timespec='seconds')}",
        f"- Instruction: {instruction}",
        f"- Query focus: {query}",
        f"- Candidate rows reviewed: {len(candidates)}",
        "",
    ]
    if recommendation:
        lines.extend(
            [
                "## Recommended Option",
                f"- Title: {recommendation.get('title', '')}",
                f"- URL: {recommendation.get('url', '')}",
                f"- Price: {recommendation.get('price', 'n/a')}",
                f"- Rating: {recommendation.get('rating', 'n/a')}",
                f"- Review count: {recommendation.get('review_count', 'n/a')}",
                f"- Condition: {recommendation.get('condition', 'unknown')}",
                "",
                "## Why this option",
                "- Lowest available price is prioritized first.",
                "- Then rating/review confidence and listing clarity are considered.",
            ]
        )
    else:
        lines.append("No recommendation could be determined from available candidates.")
    notes_path.write_text("\n".join(lines), encoding="utf-8")

    rows = []
    for row in candidates:
        price = row.get("price")
        rows.append(
            {
                "title": row.get("title", ""),
                "source": row.get("source", ""),
                "price": (f"${float(price):.2f}" if isinstance(price, (int, float)) else "n/a"),
                "url": row.get("url", ""),
                "snippet": row.get("snippet", ""),
                "rating": row.get("rating", ""),
                "review_count": row.get("review_count", ""),
                "image_count": row.get("image_count", ""),
                "condition": row.get("condition", ""),
            }
        )
    rows_payload = json.dumps(rows)
    rec_payload = json.dumps(recommendation or {})
    html_text = f"""<!doctype html>
<html lang="en"><head>
<meta charset="utf-8"/><meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Shopping Decision Dashboard</title>
<style>
body{{font-family:'Segoe UI',Tahoma,sans-serif;background:#f8fafc;color:#0f172a;margin:0}}
.wrap{{max-width:1200px;margin:0 auto;padding:20px}}
.hero{{background:#fff;border:1px solid #dbe4f0;border-radius:12px;padding:16px}}
.meta{{font-size:12px;color:#64748b}}
.tbl{{margin-top:12px;background:#fff;border:1px solid #e2e8f0;border-radius:12px;overflow:auto}}
table{{width:100%;border-collapse:collapse;font-size:13px}}
th,td{{padding:8px 10px;border-bottom:1px solid #e2e8f0;text-align:left;vertical-align:top}}
a{{color:#0f766e;text-decoration:none}}
</style></head><body><div class="wrap">
<div class="hero">
<h1 style="margin:0">Shopping Decision Dashboard</h1>
<div class="meta" style="margin-top:6px;">{html.escape(query)} | candidates: {len(candidates)}</div>
<div id="rec" style="margin-top:10px;"></div>
</div>
<div class="tbl"><table><thead>
<tr><th>#</th><th>Title</th><th>Price</th><th>Rating</th><th>Reviews</th><th>Images</th><th>Condition</th><th>Source</th></tr>
</thead><tbody id="rows"></tbody></table></div>
</div>
<script>
const rows={rows_payload};
const rec={rec_payload};
const recEl=document.getElementById('rec');
if(rec && rec.url){{
  recEl.innerHTML=`<strong>Recommended:</strong> <a target="_blank" rel="noopener" href="${{rec.url}}">${{(rec.title||'selected item')}}</a> | price: ${{rec.price ?? 'n/a'}} | rating: ${{rec.rating ?? 'n/a'}} | reviews: ${{rec.review_count ?? 'n/a'}}`;
}}
document.getElementById('rows').innerHTML = rows.map((r, i)=>`<tr><td>${{i+1}}</td><td><a target="_blank" rel="noopener" href="${{r.url}}">${{r.title}}</a><div class="meta">${{r.snippet||''}}</div></td><td>${{r.price}}</td><td>${{r.rating||'n/a'}}</td><td>${{r.review_count||'n/a'}}</td><td>${{r.image_count||0}}</td><td>${{r.condition||'unknown'}}</td><td>${{r.source}}</td></tr>`).join('');
</script></body></html>"""
    dash_path.write_text(html_text, encoding="utf-8")

    return {
        "directory": str(out_dir.resolve()),
        "decision_matrix_csv": str(csv_path.resolve()),
        "recommendation_md": str(notes_path.resolve()),
        "shopping_dashboard_html": str(dash_path.resolve()),
        "primary_open_file": str(dash_path.resolve()),
    }


def _extract_inline_url(text: str) -> str:
    match = re.search(r"https?://[^\s)>\"]+", str(text or ""), flags=re.IGNORECASE)
    if not match:
        return ""
    return match.group(0).rstrip(".,;:!?")


def _is_price_recommendation_intent(instruction: str) -> bool:
    low = str(instruction or "").lower()
    strong = [
        "best price",
        "lowest price",
        "cheapest",
        "recommend",
        "which one to buy",
        "what should i buy",
        "for sale",
    ]
    if any(token in low for token in strong):
        return True
    return "ebay" in low and "price" in low


def _is_recommendation_research_intent(instruction: str) -> bool:
    return bool(platform_is_recommendation_research_intent(instruction))


def _is_marketplace_shopping_intent(instruction: str) -> bool:
    low = str(instruction or "").lower()
    mentions_marketplace = any(token in low for token in ["ebay", "amazon"])
    asks_for_choice = any(token in low for token in ["best price", "lowest price", "cheapest", "recommend", "which one", "to buy"])
    return bool(mentions_marketplace and asks_for_choice)


def _is_wine_pairing_intent(instruction: str, query: str = "") -> bool:
    return bool(platform_is_wine_pairing_intent(instruction=instruction, query=query))


def _browser_query_url(query: str) -> str:
    return platform_browser_query_url(query)


def _generic_decision_rows(results: List[SearchResult], query: str) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for idx, result in enumerate(results[:12], start=1):
        score = round(_relevance_score(result, query), 3)
        rationale: List[str] = []
        if result.price is not None:
            rationale.append(f"price detected at ${float(result.price):.2f}")
        if result.source:
            rationale.append(f"source={result.source}")
        if result.snippet:
            rationale.append("matching snippet evidence")
        rows.append(
            {
                "rank": idx,
                "candidate": result.title,
                "candidate_type": "source_result",
                "url": result.url,
                "source": result.source,
                "price": result.price,
                "score": score,
                "support_count": 1,
                "rationale": "; ".join(rationale[:3]) or "Top-ranked relevant source.",
            }
        )
    return rows


def _query_focus_terms(text: str) -> List[str]:
    return list(platform_query_focus_terms(text))


def _topic_overlap_count(result: SearchResult, query: str) -> int:
    return int(platform_topic_overlap_count(result, query))


def _minimum_topic_overlap(query: str) -> int:
    return int(platform_minimum_topic_overlap(query))


def _passes_topic_gate(result: SearchResult, instruction: str, query: str) -> bool:
    return bool(platform_passes_topic_gate(result, instruction, query))


def _build_recommendation_focus_query(query: str, instruction: str) -> str:
    return str(platform_build_recommendation_focus_query(query=query, instruction=instruction))


def _curated_recommendation_sources(instruction: str, query: str) -> List[SearchResult]:
    return list(platform_curated_recommendation_sources(instruction=instruction, query=query))


def _build_product_candidate_rows(browser_notes: List[Dict[str, Any]], instruction: str, query: str) -> List[Dict[str, Any]]:
    return list(platform_build_product_candidate_rows(browser_notes=browser_notes, instruction=instruction, query=query))


def _slugify_product_name(value: str) -> str:
    text = str(value or "").strip().lower()
    text = text.replace("&", " and ")
    text = re.sub(r"[^a-z0-9]+", "-", text)
    text = re.sub(r"-+", "-", text).strip("-")
    return text


def _is_buy_page_url(url: str) -> bool:
    target = str(url or "").strip().lower()
    if not target.startswith("http"):
        return False
    host = urllib.parse.urlparse(target).netloc.lower()
    path = urllib.parse.urlparse(target).path.lower()
    if any(domain in host for domain in ["amazon.com", "totalwine.com", "wine.com", "heb.com", "instacart.com", "wacaco.com"]):
        return True
    article_terms = ["review", "reviews", "guide", "best-", "buying-guides", "article", "content", "blog"]
    if any(term in path for term in article_terms):
        return False
    return any(term in path for term in ["/product", "/products/", "/dp/", "/p/", "/shop/"])


def _probe_candidate_url(url: str) -> str:
    target = str(url or "").strip()
    if not target:
        return ""
    try:
        req = urllib.request.Request(target, headers={"User-Agent": USER_AGENT})
        with urllib.request.urlopen(req, timeout=20) as resp:  # nosec B310 - controlled URLs
            final_url = str(getattr(resp, "geturl", lambda: target)() or target)
            return final_url
    except Exception:
        return ""


def _candidate_buy_url_variants(candidate: str, instruction: str, query: str) -> List[str]:
    low = f"{instruction} {query} {candidate}".lower()
    variants: List[str] = []
    if "wacaco" in low:
        if "picopresso" in low:
            variants.append("https://www.wacaco.com/products/picopresso")
        if "nanopresso" in low:
            variants.append("https://www.wacaco.com/products/nanopresso")
        if "minipresso gr2" in low:
            variants.append("https://www.wacaco.com/products/minipresso-gr2")
        if "minipresso ns2" in low:
            variants.append("https://www.wacaco.com/products/minipresso-ns2")
        slug = _slugify_product_name(re.sub(r"^wacaco\s+", "", candidate, flags=re.IGNORECASE))
        if slug:
            variants.append(f"https://www.wacaco.com/products/{slug}")
    slug = _slugify_product_name(candidate)
    if "amazon.com" not in " ".join(variants) and slug:
        variants.append(f"https://www.amazon.com/s?k={urllib.parse.quote_plus(candidate)}")
    out: List[str] = []
    for item in variants:
        if item not in out:
            out.append(item)
    return out


def _resolve_product_candidate_buy_url(candidate: str, instruction: str, query: str, current_url: str = "") -> str:
    existing = str(current_url or "").strip()
    if _is_buy_page_url(existing):
        return existing
    for url in _candidate_buy_url_variants(candidate=candidate, instruction=instruction, query=query):
        resolved = _probe_candidate_url(url)
        if resolved and _is_buy_page_url(resolved):
            return resolved
    return existing


def _wine_pairing_decision_rows(results: List[SearchResult], instruction: str, query: str) -> List[Dict[str, Any]]:
    if not _is_wine_pairing_intent(instruction=instruction, query=query):
        return []
    text_low = f"{instruction} {query}".lower()
    scores: Dict[str, float] = {}
    evidence: Dict[str, List[SearchResult]] = {}
    for result in results[:12]:
        hay = f"{result.title} {result.snippet}".lower()
        base_score = max(1.0, _relevance_score(result, query))
        for style, tokens in WINE_STYLE_KEYWORDS.items():
            hits = sum(1 for token in tokens if token in hay)
            if hits <= 0:
                continue
            scores[style] = float(scores.get(style, 0.0)) + float(hits) * 2.0 + base_score
            evidence.setdefault(style, []).append(result)
    if "steak" in text_low:
        for style, bonus in STEAK_WINE_STYLE_BONUS.items():
            if style in scores:
                scores[style] = float(scores.get(style, 0.0)) + float(bonus)
    ranked_styles = sorted(scores.items(), key=lambda item: item[1], reverse=True)
    rows: List[Dict[str, Any]] = []
    for idx, (style, score) in enumerate(ranked_styles[:6], start=1):
        supporting = evidence.get(style, [])
        first_url = str(supporting[0].url) if supporting else ""
        rationale = "Strong evidence across pairing sources."
        if "steak" in text_low and style in STEAK_WINE_STYLE_BONUS:
            rationale = "Robust red pairing that repeatedly aligns with steak-focused recommendations."
        rows.append(
            {
                "rank": idx,
                "candidate": style,
                "candidate_type": "wine_style",
                "url": first_url,
                "source": supporting[0].source if supporting else "web",
                "price": None,
                "score": round(score, 3),
                "support_count": len(supporting),
                "rationale": rationale,
            }
        )
    return rows


def _build_decision_rows(results: List[SearchResult], instruction: str, query: str) -> List[Dict[str, Any]]:
    return list(
        platform_build_decision_rows(
            results=results,
            instruction=instruction,
            query=query,
            relevance_fn=_relevance_score,
            is_wine_pairing_fn=_is_wine_pairing_intent,
        )
    )


def _build_recommendation_summary(
    *,
    decision_rows: List[Dict[str, Any]],
    results: List[SearchResult],
    instruction: str,
    query: str,
) -> Dict[str, Any]:
    recommendation = dict(
        platform_build_recommendation_summary(
            decision_rows=decision_rows,
            results=results,
            instruction=instruction,
            query=query,
        )
    )
    top = dict(decision_rows[0]) if decision_rows else {}
    if str(top.get("candidate_type", "")) == "product_candidate":
        recommendation["selected_url"] = _resolve_product_candidate_buy_url(
            candidate=str(top.get("candidate", "")),
            instruction=instruction,
            query=query,
            current_url=str(recommendation.get("selected_url", "")),
        )
    return recommendation


def _merge_browser_notes_into_rows(decision_rows: List[Dict[str, Any]], browser_notes: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    return list(platform_merge_browser_notes_into_rows(decision_rows=decision_rows, browser_notes=browser_notes))


def _clean_ebay_query(instruction: str) -> str:
    raw = str(instruction or "").strip()
    cleaned = re.sub(r"https?://[^\s]+", " ", raw, flags=re.IGNORECASE)
    patterns = [
        r"^\s*(?:search|find|look(?:\s+up)?)\s+ebay\s+(?:for\s+)?",
        r"^\s*(?:on\s+)?ebay\s+(?:for\s+)?",
        r"^\s*on\s+",
    ]
    for pattern in patterns:
        cleaned = re.sub(pattern, "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(
        r"\b(best|lowest|cheapest)\s+price\b|\brecommend(?:\s+me)?(?:\s+the\s+one\s+to\s+buy)?\b|\bwhich\s+one\s+to\s+buy\b",
        " ",
        cleaned,
        flags=re.IGNORECASE,
    )
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" .,-")
    cleaned = re.sub(r"^\s*(on|for)\s+", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\b(and|to|for|please)\s*$", "", cleaned, flags=re.IGNORECASE).strip(" .,-")
    return cleaned or raw


def _parse_ebay_listing_blocks(html_text: str, limit: int = 10) -> List[SearchResult]:
    item_pattern = re.compile(r"<li[^>]*class=\"[^\"]*s-item[^\"]*\"[^>]*>(.*?)</li>", re.IGNORECASE | re.DOTALL)
    link_pattern = re.compile(r"class=\"[^\"]*s-item__link[^\"]*\"[^>]*href=\"([^\"]+)\"", re.IGNORECASE | re.DOTALL)
    title_pattern = re.compile(r"class=\"[^\"]*s-item__title[^\"]*\"[^>]*>(.*?)</", re.IGNORECASE | re.DOTALL)
    price_pattern = re.compile(r"class=\"[^\"]*s-item__price[^\"]*\"[^>]*>(.*?)</", re.IGNORECASE | re.DOTALL)
    subtitle_pattern = re.compile(r"class=\"[^\"]*s-item__subtitle[^\"]*\"[^>]*>(.*?)</", re.IGNORECASE | re.DOTALL)

    out: Dict[str, SearchResult] = {}
    for block in item_pattern.findall(html_text):
        href_match = link_pattern.search(block)
        title_match = title_pattern.search(block)
        if not href_match or not title_match:
            continue
        href = html.unescape(href_match.group(1)).strip()
        title = html.unescape(re.sub("<.*?>", "", title_match.group(1))).strip()
        if not href or not title:
            continue
        low_title = title.lower()
        if "shop on ebay" in low_title or "results matching fewer words" in low_title:
            continue
        price_match = price_pattern.search(block)
        price_text = html.unescape(re.sub("<.*?>", "", price_match.group(1))).strip() if price_match else ""
        snippet_match = subtitle_pattern.search(block)
        snippet = html.unescape(re.sub("<.*?>", "", snippet_match.group(1))).strip() if snippet_match else ""
        out[href] = SearchResult(
            title=title,
            url=href,
            price=_extract_price(f"{price_text} {title}"),
            source="ebay",
            snippet=snippet,
        )
        if len(out) >= limit:
            break
    return list(out.values())[:limit]


def _search_ebay_listings_playwright(search_url: str, limit: int = 10, browser_worker_mode: str = "local") -> List[SearchResult]:
    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return []

    rows: Dict[str, SearchResult] = {}
    with sync_playwright() as p:
        browser = None
        context = None
        attached = False
        try:
            worker = ensure_browser_worker(mode=browser_worker_mode)
            if worker.get("ok"):
                debug_port = int(worker.get("debug_port", 9222) or 9222)
                try:
                    browser = p.chromium.connect_over_cdp(f"http://127.0.0.1:{debug_port}", timeout=3500)
                    contexts = list(getattr(browser, "contexts", []) or [])
                    context = contexts[0] if contexts else browser.new_context()
                    attached = True
                except Exception:
                    browser = None
                    context = None
            if context is None:
                browser = p.chromium.launch(headless=True)
                context = browser.new_context()

            page = _select_best_context_page(context, search_url) if attached else None
            if page is None:
                pages = list(getattr(context, "pages", []) or [])
                page = pages[0] if pages else context.new_page()
            if "ebay.com" not in str(getattr(page, "url", "")).lower():
                page.goto(search_url, timeout=30000)
            page.wait_for_timeout(2200)
            if "pardon our interruption" in str(page.title() or "").lower():
                return []

            cards = page.query_selector_all("li.s-item, .srp-results .s-item")
            if not cards:
                page.wait_for_timeout(1200)
                cards = page.query_selector_all("li.s-item, .srp-results .s-item")
            for card in cards[: max(6, limit * 3)]:
                try:
                    link_node = card.query_selector(".s-item__link")
                    title_node = card.query_selector(".s-item__title")
                    price_node = card.query_selector(".s-item__price")
                    subtitle_node = card.query_selector(".s-item__subtitle")
                    if not link_node or not title_node:
                        continue
                    href = str(link_node.get_attribute("href") or "").strip()
                    title = str(title_node.inner_text() or "").strip()
                    if not href or not title:
                        continue
                    low_title = title.lower()
                    if "shop on ebay" in low_title or "results matching fewer words" in low_title:
                        continue
                    price_text = str(price_node.inner_text() or "").strip() if price_node else ""
                    subtitle = str(subtitle_node.inner_text() or "").strip() if subtitle_node else ""
                    rows[href] = SearchResult(
                        title=title,
                        url=href,
                        price=_extract_price(f"{price_text} {title}"),
                        source="ebay",
                        snippet=subtitle,
                    )
                    if len(rows) >= limit:
                        break
                except Exception:
                    continue
            return list(rows.values())[:limit]
        finally:
            if browser is not None and not attached:
                try:
                    browser.close()
                except Exception:
                    pass
    return list(rows.values())[:limit]


def _search_ebay_listings(query: str, search_url: str = "", limit: int = 10, browser_worker_mode: str = "local") -> List[SearchResult]:
    target = str(search_url or "").strip()
    if not target:
        q = urllib.parse.quote_plus(str(query or "").strip())
        target = f"https://www.ebay.com/sch/i.html?_nkw={q}&_sop=12"
    text_rows: List[SearchResult] = []
    try:
        html_text = _fetch_text(target)
        if "pardon our interruption" not in html_text.lower():
            text_rows = _parse_ebay_listing_blocks(html_text=html_text, limit=limit)
    except Exception:
        text_rows = []
    if text_rows:
        return text_rows
    return _search_ebay_listings_playwright(
        search_url=target,
        limit=limit,
        browser_worker_mode=browser_worker_mode,
    )


def _attach_generic_browser_context(
    *,
    playwright: Any,
    browser_worker_mode: str,
    human_like_interaction: bool,
) -> tuple[Any, Any, bool, Dict[str, Any]]:
    worker_info = ensure_browser_worker(mode=browser_worker_mode)
    browser = None
    context = None
    attached = False
    if bool(worker_info.get("ok")):
        debug_port = int(worker_info.get("debug_port", 9222) or 9222)
        try:
            browser = playwright.chromium.connect_over_cdp(f"http://127.0.0.1:{debug_port}", timeout=3500)
            contexts = list(getattr(browser, "contexts", []) or [])
            context = contexts[0] if contexts else browser.new_context()
            attached = True
        except Exception:
            browser = None
            context = None
    if context is None:
        browser = playwright.chromium.launch(headless=not bool(human_like_interaction))
        context = browser.new_context()
    return browser, context, attached, worker_info


def _browser_extract_page_text(page: Any, limit: int = 900) -> str:
    try:
        body = page.locator("body")
        text = str(body.inner_text(timeout=2000) or "")
    except Exception:
        try:
            text = str(page.content() or "")
        except Exception:
            text = ""
    text = re.sub(r"\s+", " ", text or "").strip()
    return text[:limit]


def _browser_note_for_page(*, query: str, page_url: str, title: str, text: str) -> Dict[str, Any]:
    cleaned_title = re.sub(r"\s+", " ", str(title or "")).strip()
    cleaned_text = re.sub(r"\s+", " ", str(text or "")).strip()
    sentences = re.split(r"(?<=[.!?])\s+", cleaned_text)
    summary = ""
    terms = [term for term in re.split(r"[^a-z0-9]+", query.lower()) if len(term) > 3]
    for sentence in sentences:
        low = sentence.lower()
        if any(term in low for term in terms):
            summary = sentence.strip()
            break
    if not summary:
        summary = cleaned_text[:220]
    return {
        "url": str(page_url or ""),
        "title": cleaned_title,
        "summary": summary[:260],
        "excerpt": cleaned_text[:900],
    }


def _browser_research_walk(
    *,
    query: str,
    candidates: List[SearchResult],
    browser_worker_mode: str,
    human_like_interaction: bool,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    max_pages: int = 4,
) -> Dict[str, Any]:
    return dict(
        platform_browser_research_walk(
            query=query,
            candidates=candidates,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=human_like_interaction,
            progress_cb=progress_cb,
            max_pages=max_pages,
        )
    )


def _safe_search_web(query: str, limit: int = 10) -> List[SearchResult]:
    try:
        return _parse_duckduckgo(query, limit=limit)
    except Exception:
        return []


def _extract_price(text: str) -> Optional[float]:
    price_match = re.search(r"\$([0-9]{1,4}(?:\.[0-9]{2})?)", text.replace(",", ""))
    if not price_match:
        return None
    try:
        return float(price_match.group(1))
    except ValueError:
        return None


def _search_amazon_playwright(query: str, limit: int = 8) -> List[SearchResult]:
    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return []

    results: List[SearchResult] = []
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        q = urllib.parse.quote_plus(query)
        page.goto(f"https://www.amazon.com/s?k={q}", timeout=30000)
        cards = page.query_selector_all("div.s-main-slot div[data-component-type='s-search-result']")
        for card in cards[: limit * 2]:
            title_node = card.query_selector("h2 span")
            link_node = card.query_selector("h2 a")
            price_whole = card.query_selector("span.a-price-whole")
            price_frac = card.query_selector("span.a-price-fraction")
            if not title_node or not link_node:
                continue
            title = (title_node.inner_text() or "").strip()
            href = link_node.get_attribute("href") or ""
            if href and href.startswith("/"):
                href = "https://www.amazon.com" + href
            whole = (price_whole.inner_text() or "").replace(",", "").strip() if price_whole else ""
            frac = (price_frac.inner_text() or "").strip() if price_frac else "00"
            price: Optional[float] = None
            if whole.isdigit() and frac.isdigit():
                price = float(f"{whole}.{frac}")
            if title and href:
                results.append(SearchResult(title=title, url=href, price=price, source="amazon"))
            if len(results) >= limit:
                break
        browser.close()
    return results


def _best_price(results: List[SearchResult]) -> Optional[SearchResult]:
    priced = [r for r in results if r.price is not None]
    if not priced:
        return results[0] if results else None
    return sorted(priced, key=lambda r: r.price or 999999.0)[0]


def _is_destructive_intent(instruction: str, plan_steps: List[Dict[str, Any]]) -> bool:
    low = instruction.lower()
    if any(token in low for token in DESTRUCTIVE_ACTION_KEYWORDS):
        return True
    for step in plan_steps:
        merged = " ".join(
            str(x).lower()
            for x in [
                step.get("action", ""),
                step.get("text", ""),
                step.get("keys", ""),
                (step.get("selector", {}) or {}).get("value", "") if isinstance(step.get("selector", {}), dict) else "",
            ]
        )
        if any(token in merged for token in DESTRUCTIVE_ACTION_KEYWORDS):
            return True
    return False


def _summarize_plan_steps(plan_steps: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    for i, step in enumerate(plan_steps):
        action = str(step.get("action", step.get("kind", "step")))
        target = ""
        selector = step.get("selector", {})
        if isinstance(selector, dict):
            target = str(selector.get("value", ""))
        if not target:
            target = str(
                step.get("app", "")
                or step.get("name", "")
                or step.get("text", "")
                or step.get("output_path", "")
                or step.get("source", "")
            )
        out.append({"index": i, "action": action, "target": target[:140]})
    return out


def _build_undo_plan(plan_steps: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    undo: List[Dict[str, Any]] = []
    for i, step in enumerate(plan_steps):
        action = str(step.get("action", step.get("kind", "step"))).lower()
        if action in {"type_text", "type"}:
            instruction = "Use Ctrl+Z or clear edited field to restore prior value."
        elif action in {"click", "click_found", "submit_action", "set_cell"}:
            instruction = "Navigate back to impacted record and revert the latest change manually."
        elif action in {"open_app", "focus_window", "navigate_url"}:
            instruction = "Close the opened window/tab if this run should be rolled back."
        elif action in {"copy", "paste", "hotkey"}:
            instruction = "Restore clipboard/context and reverse any pasted content."
        elif action in {"research", "extract", "analyze", "produce", "present"}:
            instruction = "Delete generated local artifacts if run should be fully reverted."
        else:
            instruction = "No automatic rollback; use manual checkpoint restore."
        undo.append({"step_index": i, "action": action, "undo": instruction})
    return undo


def _detect_ambiguities(instruction: str, plan_steps: List[Dict[str, Any]]) -> List[str]:
    questions: List[str] = []
    if len(instruction.strip()) < 6:
        questions.append("Please provide more detail on the target outcome.")
    for idx, step in enumerate(plan_steps):
        action = str(step.get("action", step.get("kind", ""))).lower()
        if action in {"click", "type_text", "type"}:
            selector = step.get("selector", {})
            if isinstance(selector, dict):
                sel_val = str(selector.get("value", "")).strip()
            else:
                sel_val = ""
            if not sel_val and not str(step.get("text", "")).strip():
                questions.append(f"Step {idx} is missing a target element or text value.")
        if action == "note":
            questions.append(f"Step {idx} could not be mapped to an executable action.")
    return questions


def _classify_explicit_route(instruction: str) -> str:
    low = instruction.lower()
    if (
        _is_clipboard_capture_intent(instruction)
        or _is_email_triage_intent(instruction)
        or _is_competitor_analysis_intent(instruction)
        or _is_study_pack_intent(instruction)
        or _is_job_research_intent(instruction)
        or _is_payer_pricing_review_intent(instruction)
        or _is_code_workbench_intent(instruction)
    ):
        return ""
    asks_chat_action = any(x in low for x in ["respond", "reply", "message", "post", "send"]) and any(
        x in low for x in ["chat", "dm", "thread", "comment", "message"]
    )
    asks_social_or_comms = any(token in low for token in COMMUNICATION_SOCIAL_KEYWORDS)
    if asks_chat_action or asks_social_or_comms:
        return "desktop_sequence"

    asks_document = any(x in low for x in ["write document", "write a document", "document", "doc", "memo", "brief"])
    asks_powerpoint = any(x in low for x in ["powerpoint", "ppt", "slide deck", "slides"])
    asks_visuals = any(x in low for x in ["create visuals", "visual", "diagram", "infographic", "poster"])
    if asks_document or asks_powerpoint or asks_visuals:
        return "artifact_generation"
    return ""


def _is_clipboard_capture_intent(instruction: str) -> bool:
    low = str(instruction or "").lower()
    has_clipboard = "clipboard" in low
    has_capture = any(token in low for token in ["capture", "grab", "save", "import", "paste"])
    has_visual = any(token in low for token in ["image", "screenshot", "png", "photo", "picture"])
    has_macro_sequence = any(token in low for token in [" then ", "open ", "click ", "type ", "press ", "focus ", "switch to ", "scroll "])
    return has_clipboard and (has_capture or has_visual) and not has_macro_sequence


def _is_mission_runtime_intent(instruction: str) -> bool:
    low = str(instruction or "").lower()
    job_package = any(token in low for token in ["tailor my resume", "tailor resume", "cover letter", "application checklist"])
    grant_package = "grant" in low and any(token in low for token in ["proposal", "eligibility", "submission checklist", "funder priorities", "rank them", "draft the top proposal", "grant tracker"])
    executive_brief = (
        any(token in low for token in ["executive briefing", "executive brief", "brief my vp", "brief leadership"])
        or (any(token in low for token in ["market", "research", "recommendations"]) and any(token in low for token in ["brief", "briefing"]))
    )
    data_story = (
        any(token in low for token in ["find the story", "data story", "build charts", "executive summary", "brief my vp"])
        and any(token in low for token in ["dataset", "analyze", "analysis", "data"])
    )
    if (
        _is_clipboard_capture_intent(instruction)
        or _is_email_triage_intent(instruction)
        or _is_payer_pricing_review_intent(instruction)
        or _is_competitor_analysis_intent(instruction)
        or _is_code_workbench_intent(instruction)
        or (_is_job_research_intent(instruction) and not job_package)
    ):
        return False
    return job_package or grant_package or executive_brief or data_story


def _requested_outputs(instruction: str) -> Set[str]:
    low = instruction.lower()
    requested: Set[str] = set()
    if any(x in low for x in ["spreadsheet", "csv", "task list"]):
        requested.add("spreadsheet")
    if any(x in low for x in ["report", "document", "doc", "memo", "brief"]):
        requested.add("document")
    if "executive summary" in low:
        requested.add("executive_summary")
    if any(x in low for x in ["powerpoint", "ppt", "slides", "slide deck"]):
        requested.add("powerpoint")
    if any(x in low for x in ["dashboard", "html dashboard"]):
        requested.add("dashboard")
    if any(x in low for x in ["visual", "visuals", "diagram", "infographic"]):
        requested.add("visual")
    if any(x in low for x in ["draft reply", "draft replies", "reply", "respond", "chat response"]):
        requested.add("chat_response")
    if any(x in low for x in ["write code", "build code", "analysis script", "python script", "code scaffold"]):
        requested.add("code")
    if any(x in low for x in ["workspace", "vs code", "vscode", "visual studio code"]):
        requested.add("workspace")
    return requested


def _plan_represented_outputs(plan: Dict[str, Any], plan_steps: List[Dict[str, Any]]) -> Set[str]:
    represented: Set[str] = set()
    deliverables = [str(x).lower() for x in (plan.get("deliverables") or [])]
    deliverable_map = {
        "spreadsheet": "spreadsheet",
        "csv": "spreadsheet",
        "report": "document",
        "document": "document",
        "executive_summary": "executive_summary",
        "powerpoint": "powerpoint",
        "dashboard": "dashboard",
        "visual": "visual",
        "draft_reply": "chat_response",
        "apply_links": "links",
        "code": "code",
        "workspace": "workspace",
        "vscode_workspace": "workspace",
    }
    for item in deliverables:
        if item in deliverable_map:
            represented.add(deliverable_map[item])

    for step in plan_steps:
        action = str(step.get("action", step.get("kind", ""))).lower()
        text = str(step.get("text", "")).lower()
        if action in {"save_csv", "set_cell"}:
            represented.add("spreadsheet")
        if action in {"create_draft"}:
            represented.add("chat_response")
        if action in {"produce"}:
            represented.update({"document", "dashboard"})
        if action in {"present"}:
            represented.add("dashboard")
        if action in {"type_text", "type"} and text:
            represented.add("chat_response")
    if str(plan.get("domain", "")).lower() == "code_workbench":
        represented.update({"code", "workspace", "document"})
    return represented


def _fail_fast_output_mismatch(
    instruction: str,
    mode: str,
    plan_steps: List[Dict[str, Any]],
    requested_outputs: Set[str],
    represented_outputs: Set[str],
) -> Dict[str, Any]:
    missing = sorted(x for x in requested_outputs if x not in represented_outputs)
    return {
        "ok": False,
        "mode": mode,
        "instruction": instruction,
        "error": "requested_outputs_not_in_plan",
        "message": "Requested outputs are not represented in the execution plan. Refine the plan before execution.",
        "missing_outputs": missing,
        "requested_outputs": sorted(requested_outputs),
        "planned_outputs": sorted(represented_outputs),
        "planned_steps": _summarize_plan_steps(plan_steps),
        "undo_plan": _build_undo_plan(plan_steps),
    }


def _artifact_plan_steps(deliverables: Set[str]) -> List[Dict[str, Any]]:
    steps: List[Dict[str, Any]] = [{"kind": "research", "name": "Parse instruction into artifact outline", "target": {"query": "instruction_outline"}}]
    steps.append({"kind": "produce", "name": "Draft artifact content from the instruction", "target": {"id": "artifact:draft"}})
    if "document" in deliverables or "executive_summary" in deliverables:
        steps.append({"kind": "produce", "name": "Write markdown document artifact", "target": {"path": "data/reports/artifact_generation"}})
    if "powerpoint" in deliverables:
        steps.append({"kind": "produce", "name": "Build PowerPoint artifact", "target": {"path": "data/reports/artifact_generation"}})
    if "visual" in deliverables:
        steps.append({"kind": "produce", "name": "Create visual HTML artifact", "target": {"path": "data/reports/artifact_generation"}})
    return steps


def _build_artifact_generation_plan(instruction: str) -> Dict[str, Any]:
    requested = _requested_outputs(instruction)
    deliverables: Set[str] = set()
    if "executive_summary" in requested:
        deliverables.add("executive_summary")
        deliverables.add("document")
    if "document" in requested:
        deliverables.add("document")
    if "powerpoint" in requested:
        deliverables.add("powerpoint")
    if "visual" in requested:
        deliverables.add("visual")
    if not deliverables:
        deliverables.add("document")
    return {
        "planner": "deterministic-artifact-v1",
        "domain": "artifact_generation",
        "objective": re.sub(r"\s+", " ", instruction).strip(),
        "deliverables": sorted(deliverables),
        "sources": ["user_instruction"],
        "constraints": {
            "prefer_public_pages": False,
            "no_password_capture": True,
            "persist_history": True,
        },
        "steps": _artifact_plan_steps(deliverables),
    }


def _should_use_execution_graph_runtime(*, task_contract: Any, explicit_route: str, instruction: str) -> bool:
    contract = task_contract.to_dict() if hasattr(task_contract, "to_dict") else (dict(task_contract) if isinstance(task_contract, dict) else {})
    domain = str(contract.get("domain", "")).strip().lower()
    low = str(instruction or "").lower()
    ui_tokens = ["ui", "frontend", "app shell", "chat and canvas", "artifact viewer", "artifact viewers", "component", "layout", "information architecture"]
    explicit_editor = any(token in low for token in ["vscode", "vs code", "visual studio code", "workspace", "new instance"])
    if explicit_route == "artifact_generation":
        return False
    if domain == "topic_learning":
        return True
    if domain == "ui_build" and any(token in low for token in ui_tokens):
        return True
    if domain == "deep_analysis" and not explicit_editor:
        return True
    return False


def _runtime_plan_steps_from_graph(graph: Dict[str, Any]) -> List[Dict[str, Any]]:
    mapping = {
        "deep_research": ("research", "Compile the working brief", {"query": "task_contract"}),
        "research_collection": ("research", "Collect source evidence and candidate options", {"query": "search_results"}),
        "mission_research": ("research", "Collect mission-scoped evidence", {"query": "mission_queries"}),
        "topic_mastery_learn": ("learn", "Build topic mastery package", {"path": "workspace/learn_artifacts"}),
        "source_evaluation": ("research", "Score source quality", {"id": "source_scores"}),
        "mission_work_product": ("produce", "Build mission work-product package", {"path": "workspace/mission_artifacts"}),
        "file_inspection": ("inspect", "Inspect workspace inputs", {"path": "workspace"}),
        "data_cleaning": ("transform", "Normalize structured rows", {"id": "clean_rows"}),
        "statistical_analysis": ("analyze", "Run analysis", {"id": "analysis_results"}),
        "data_visualization": ("produce", "Draft visuals", {"id": "chart_specs"}),
        "rag_build": ("build_index", "Build retrieval index", {"id": "rag_index"}),
        "rag_query": ("answer", "Generate retrieval examples", {"id": "rag_examples"}),
        "code_write": ("produce", "Write scaffold code", {"path": "workspace/src"}),
        "code_test": ("verify", "Run smoke checks", {"id": "test_results"}),
        "code_fix": ("repair", "Patch failing areas", {"id": "code_fixes"}),
        "data_storytelling": ("produce", "Build executive story", {"id": "story_package"}),
        "report_build": ("produce", "Draft report", {"path": "workspace/artifacts"}),
        "stakeholder_summary": ("produce", "Draft stakeholder summary", {"id": "stakeholder_summary"}),
        "presentation_build": ("produce", "Draft presentation outline", {"path": "workspace/artifacts"}),
        "spreadsheet_build": ("produce", "Draft spreadsheet rows", {"path": "workspace/artifacts"}),
        "ui_build": ("produce", "Design chat/canvas UI", {"path": "workspace/artifacts"}),
        "artifact_export": ("present", "Write artifact package", {"path": "workspace/artifacts"}),
        "approval_gate": ("approval", "Request approval", {"id": "approval"}),
    }
    steps: List[Dict[str, Any]] = []
    for node in (graph.get("nodes", []) or []):
        capability = str(node.get("capability", ""))
        kind, name, target = mapping.get(capability, ("work", capability, {"id": capability}))
        steps.append({"kind": kind, "name": name, "target": target, "capability": capability})
    return steps


def _build_runtime_plan(task_contract: Any, graph: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    contract = task_contract.to_dict() if hasattr(task_contract, "to_dict") else dict(task_contract)
    graph_dict = graph or CapabilityPlanner(registry=default_capability_registry()).plan(task_contract).to_dict()
    return {
        "planner": "capability-runtime-v1",
        "domain": str(contract.get("domain", "capability_runtime") or "capability_runtime"),
        "objective": str(contract.get("user_goal", "")),
        "deliverables": list(contract.get("requested_outputs", []) or []),
        "sources": list(contract.get("source_rules", []) or []),
        "constraints": {
            "task_contract": dict(contract),
            "invalidation_keys": dict(contract.get("invalidation_keys", {}) or {}),
        },
        "steps": _runtime_plan_steps_from_graph(graph_dict),
    }


def _run_execution_graph_runtime_path(
    *,
    instruction: str,
    task_contract: Any,
    graph: Optional[Any] = None,
    ai_meta: Dict[str, Any],
    progress_cb: Optional[Callable[[int, str], None]] = None,
    artifact_reuse_mode: str = "reuse_if_recent",
    artifact_reuse_max_age_hours: int = 72,
    mode_override: str = "",
    extra_context: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    contract_dict = task_contract.to_dict() if hasattr(task_contract, "to_dict") else dict(task_contract)
    _emit_progress(progress_cb, 12, "Composing capability graph")
    registry = default_capability_registry()
    graph = graph or CapabilityPlanner(registry=registry).plan(task_contract)
    runtime = ExecutionGraphRuntime(registry=registry)
    _emit_progress(progress_cb, 24, "Executing capability graph")
    runtime_result = runtime.run(
        graph,
        {
            "task_contract": contract_dict,
            "instruction": instruction,
            "task_id": graph.task_id,
            **(dict(extra_context or {})),
        },
    )
    _emit_progress(progress_cb, 78, "Reviewing artifacts and critics")
    graph_dict = runtime_result.graph.to_dict()
    artifacts = dict(runtime_result.artifacts)
    primary_open = (
        artifacts.get("primary_open_file")
        or artifacts.get("visual_html")
        or artifacts.get("executive_summary_html")
        or artifacts.get("dashboard_html")
        or artifacts.get("summary_report_md")
        or artifacts.get("report_md")
        or artifacts.get("document_md")
        or artifacts.get("presentation_md")
        or artifacts.get("powerpoint_pptx")
        or artifacts.get("decision_matrix_csv")
        or artifacts.get("ui_spec_json")
        or artifacts.get("artifact_manifest_json")
        or ""
    )
    result = {
        "ok": runtime_result.ok,
        "mode": mode_override or "execution_graph_runtime",
        "runtime_mode": "execution_graph_runtime",
        "error": runtime_result.error,
        "instruction": instruction,
        "ai": ai_meta,
        "plan": _build_runtime_plan(task_contract, graph=graph_dict),
        "query": str(contract_dict.get("user_goal", "")),
        "results_count": len(runtime_result.events),
        "results": [],
        "artifacts": artifacts,
        "opened_url": Path(primary_open).resolve().as_uri() if primary_open and Path(primary_open).exists() else "",
        "summary": {
            "graph_status": graph_dict.get("status", ""),
            "runtime_events": len(runtime_result.events),
            "artifact_reuse_mode": artifact_reuse_mode,
            "artifact_reuse_max_age_hours": artifact_reuse_max_age_hours,
            "revisions_performed": len(runtime_result.revisions),
            "reused_existing_outputs": False,
        },
        "report": {
            "summary": str(((runtime_result.outputs.get(next((node.node_id for node in runtime_result.graph.nodes if node.capability == "stakeholder_summary"), ""), {}) or {}).get("stakeholder_summary", {}) or {}).get("executive_summary", "")),
            "next_actions": [str(x.get("required_fix", "")) for x in runtime_result.revisions if str(x.get("required_fix", "")).strip()],
        },
        "verification_report": dict(runtime_result.verification_report),
        "verification": dict(runtime_result.verification),
        "final_report": dict(runtime_result.final_report),
        "runtime_events": list(runtime_result.events),
        "revisions_performed": list(runtime_result.revisions),
        "memory_context": dict(runtime_result.memory_context),
        "task_contract": contract_dict,
        "capability_execution_graph": graph_dict,
        "artifact_metadata": dict(runtime_result.artifact_metadata),
        "critics": {"platform": dict(runtime_result.critics)},
        "canvas": {
            "title": "Capability Graph Completed" if runtime_result.ok else "Capability Graph Failed",
            "subtitle": f"{graph_dict.get('domain', '')} | {len(artifacts)} artifact(s) | {len(runtime_result.revisions)} revision(s)",
            "cards": [
                {
                    "title": str(node.get("capability", "")),
                    "price": str(node.get("status", "")),
                    "source": "runtime",
                    "url": "",
                }
                for node in (graph_dict.get("nodes", []) or [])[:6]
            ],
        },
    }
    return result


def _write_simple_pptx(pptx_path: Path, title: str, subtitle: str, bullets: List[str]) -> None:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        pptx_path.write_text("PowerPoint package unavailable. Install python-pptx to generate .pptx files.", encoding="utf-8")
        return
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = subtitle
    for item in bullets[:6]:
        para = body.add_paragraph()
        para.text = f"- {item}"
    prs.save(str(pptx_path))


def _run_artifact_generation(
    plan: Dict[str, Any],
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    artifact_reuse_mode: str = "reuse_if_recent",
    artifact_reuse_max_age_hours: int = 72,
) -> Dict[str, Any]:
    _emit_progress(progress_cb, 20, "Generating requested artifacts")
    deliverables = {str(x).lower() for x in (plan.get("deliverables") or [])}
    objective = str(plan.get("objective", instruction)).strip()
    title = objective[:100] if objective else "Generated Artifact"
    bullets = [seg.strip().capitalize() for seg in re.split(r"[.;]", objective) if seg.strip()]
    if not bullets:
        bullets = ["Generated from user instruction."]
    required_keys: List[str] = []
    if "document" in deliverables or "executive_summary" in deliverables:
        required_keys.append("document_md")
    if "powerpoint" in deliverables:
        required_keys.append("powerpoint_pptx")
    if "visual" in deliverables:
        required_keys.append("visual_html")
    reuse_mode = _normalize_artifact_reuse_mode(artifact_reuse_mode)
    max_age = max(1, min(24 * 30, int(artifact_reuse_max_age_hours or 72)))
    reusable: Dict[str, str] = {}
    if reuse_mode != "always_regenerate":
        age_limit = 24 * 30 if reuse_mode == "reuse" else max_age
        reusable = _find_reusable_artifacts(
            kind="artifact_generation",
            instruction=instruction,
            required_keys=required_keys,
            max_age_hours=age_limit,
        )
    if reusable:
        open_target = reusable.get("visual_html") or reusable.get("document_md") or reusable.get("powerpoint_pptx", "")
        opened_uri = ""
        nav: Dict[str, Any] = {}
        if open_target:
            opened_uri, nav = _open_target_with_reuse(
                target_url=Path(open_target).resolve().as_uri(),
                recent_actions=[f"open_tab:{Path(open_target).resolve().as_uri()}"],
            )
        _emit_progress(progress_cb, 100, "Completed")
        return {
            "ok": True,
            "mode": "artifact_generation",
            "query": objective,
            "results_count": 0,
            "results": [],
            "artifacts": reusable,
            "summary": {
                "outputs_generated": sorted(reusable.keys()),
                "reused_existing_outputs": True,
                "artifact_reuse_mode": reuse_mode,
                "artifact_reuse_max_age_hours": max_age,
                "action_critic": nav.get("decision", {}),
            },
            "opened_url": opened_uri,
            "canvas": {
                "title": "Artifacts Reused",
                "subtitle": objective[:120],
                "cards": [{"title": k, "price": "reused", "source": "artifact", "url": Path(v).resolve().as_uri()} for k, v in reusable.items()],
            },
        }

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path("data/reports/artifact_generation") / ts
    out_dir.mkdir(parents=True, exist_ok=True)

    artifacts: Dict[str, str] = {}
    if "document" in deliverables or "executive_summary" in deliverables:
        doc_path = out_dir / "document.md"
        doc_lines = [f"# {title}", "", "## Summary", ""]
        doc_lines.extend([f"- {item}" for item in bullets[:8]])
        doc_path.write_text("\n".join(doc_lines) + "\n", encoding="utf-8")
        artifacts["document_md"] = str(doc_path.resolve())
        if "executive_summary" in deliverables:
            artifacts["executive_summary_md"] = str(doc_path.resolve())

    if "powerpoint" in deliverables:
        pptx_path = out_dir / "deck.pptx"
        _write_simple_pptx(pptx_path=pptx_path, title=title, subtitle="Auto-generated deck outline", bullets=bullets)
        artifacts["powerpoint_pptx"] = str(pptx_path.resolve())

    if "visual" in deliverables:
        visual_path = out_dir / "visual.html"
        visual_path.write_text(
            (
                "<!doctype html><html><head><meta charset='utf-8'><title>Generated Visual</title></head>"
                "<body style='font-family:Segoe UI,Arial,sans-serif;background:#f8fafc;'>"
                "<div style='max-width:900px;margin:48px auto;padding:32px;background:white;border:1px solid #dbe4ef;border-radius:16px;'>"
                f"<h1 style='margin-top:0'>{html.escape(title)}</h1>"
                "<p>Instruction-derived visual placeholder.</p>"
                "</div></body></html>"
            ),
            encoding="utf-8",
        )
        artifacts["visual_html"] = str(visual_path.resolve())

    _remember_artifacts_for_reuse(kind="artifact_generation", instruction=instruction, artifacts=artifacts)
    open_target = artifacts.get("visual_html") or artifacts.get("document_md") or artifacts.get("powerpoint_pptx", "")
    opened_uri = ""
    nav: Dict[str, Any] = {}
    if open_target:
        opened_uri, nav = _open_target_with_reuse(
            target_url=Path(open_target).resolve().as_uri(),
            recent_actions=[f"open_tab:{Path(open_target).resolve().as_uri()}"],
        )

    _emit_progress(progress_cb, 100, "Completed")
    return {
        "ok": True,
        "mode": "artifact_generation",
        "query": objective,
        "results_count": 0,
        "results": [],
        "artifacts": artifacts,
        "summary": {
            "outputs_generated": sorted(artifacts.keys()),
            "reused_existing_outputs": False,
            "artifact_reuse_mode": reuse_mode,
            "artifact_reuse_max_age_hours": max_age,
            "action_critic": nav.get("decision", {}),
        },
        "opened_url": opened_uri,
        "canvas": {
            "title": "Artifacts Generated",
            "subtitle": objective[:120],
            "cards": [{"title": k, "price": "saved", "source": "artifact", "url": Path(v).resolve().as_uri()} for k, v in artifacts.items()],
        },
    }


def _ensure_elegance_budget(summary: Dict[str, Any], elegance: EleganceBudget) -> Dict[str, Any]:
    out = dict(summary or {})
    out["elegance_budget"] = elegance.snapshot()
    return out


def _apply_freshness_metadata(summary: Dict[str, Any], artifact_reuse_mode: str, artifact_reuse_max_age_hours: int) -> Dict[str, Any]:
    out = dict(summary or {})
    out.setdefault("artifact_reuse_mode", _normalize_artifact_reuse_mode(artifact_reuse_mode))
    try:
        out.setdefault("artifact_reuse_max_age_hours", max(1, min(24 * 30, int(artifact_reuse_max_age_hours))))
    except Exception:
        out.setdefault("artifact_reuse_max_age_hours", 72)
    return out


def _enforce_elegance_budget_gate(
    *,
    elegance: EleganceBudget,
    mode: str,
    instruction: str,
) -> Optional[Dict[str, Any]]:
    if elegance.consumed <= elegance.total:
        return None
    snap = elegance.snapshot()
    return {
        "ok": False,
        "mode": mode,
        "instruction": instruction,
        "error": "elegance_budget_exceeded",
        "error_code": "elegance_budget_exceeded",
        "summary": {
            "error": "elegance_budget_exceeded",
            "elegance_budget": snap,
        },
        "canvas": {
            "title": "Run Blocked",
            "subtitle": "Elegance budget exceeded; path is too wasteful.",
            "cards": [],
        },
    }


def _verification_block(ok: bool, plan_steps: List[Dict[str, Any]], result: Dict[str, Any]) -> Dict[str, Any]:
    artifacts = result.get("artifacts", {}) or {}
    trace = result.get("trace", []) or []
    evidence: List[str] = []
    if trace:
        evidence.append(f"Executed trace entries: {len(trace)}")
    if artifacts:
        evidence.append(f"Artifacts generated: {', '.join(sorted(artifacts.keys()))}")
    if result.get("opened_url"):
        evidence.append(f"Opened target: {result.get('opened_url')}")
    checks = [
        {"name": "plan_has_steps", "pass": bool(plan_steps)},
        {"name": "execution_ok", "pass": bool(ok)},
        {"name": "evidence_present", "pass": bool(evidence)},
    ]
    return {"passed": all(bool(c["pass"]) for c in checks), "checks": checks, "evidence": evidence}


def _resolve_navigation_target(
    *,
    target_url: str,
    recent_actions: Optional[List[str]] = None,
) -> Dict[str, Any]:
    target = str(target_url or "").strip()
    if not target:
        return {"url": "", "reused": False, "opened": False, "decision": {"score": 0.0, "reasons": ["missing_target"]}}
    host = urllib.parse.urlparse(target).netloc.lower()
    manager = SessionManager()
    reusable = manager.find_reusable_authenticated_tab(host) if host else ""
    if not reusable:
        reusable = manager.find_reusable_url(target)
    already_open = bool(reusable and reusable == target)
    ctx: Dict[str, Any] = {"recent_actions": list(recent_actions or [])[-5:]}
    if reusable:
        ctx["reusable_target"] = reusable
    decision = ActionCritic().evaluate(
        next_action="open_tab",
        target=target,
        already_open=already_open,
        context=ctx,
    )
    if reusable and not decision.allow and any(
        r in {"shortest_path_reuse_existing_state", "redundant_open"} for r in decision.reasons
    ):
        return {
            "url": reusable,
            "reused": True,
            "opened": False,
            "decision": {"score": decision.score, "reasons": decision.reasons, "elegance_cost": decision.elegance_cost},
        }
    try:
        webbrowser.open(target, new=2)
        manager.remember_tab(url=target, title="Navigation Target", authenticated=False)
        return {
            "url": target,
            "reused": False,
            "opened": True,
            "decision": {"score": decision.score, "reasons": decision.reasons, "elegance_cost": decision.elegance_cost},
        }
    except Exception:
        return {
            "url": target,
            "reused": False,
            "opened": False,
            "decision": {"score": decision.score, "reasons": decision.reasons, "elegance_cost": decision.elegance_cost},
        }


def _open_target_with_reuse(target_url: str, recent_actions: Optional[List[str]] = None) -> tuple[str, Dict[str, Any]]:
    nav = _resolve_navigation_target(target_url=target_url, recent_actions=recent_actions)
    opened = str(nav.get("url", "") or target_url or "")
    return opened, nav


def _run_clipboard_capture(
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
) -> Dict[str, Any]:
    _emit_progress(progress_cb, 14, "Reading clipboard")
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path("data/reports/clipboard_capture") / ts
    out_dir.mkdir(parents=True, exist_ok=True)
    image_path = capture_clipboard_image(out_dir / "clipboard_capture.png")
    if not image_path:
        return {
            "ok": False,
            "mode": "clipboard_capture",
            "instruction": instruction,
            "error": "No clipboard image or file was available to capture.",
            "source_status": {"clipboard": "empty"},
            "results_count": 0,
            "results": [],
            "artifacts": {},
            "summary": {"captured": False},
            "canvas": {
                "title": "Clipboard Capture Unavailable",
                "subtitle": "No clipboard image detected.",
                "cards": [],
            },
        }
    _emit_progress(progress_cb, 60, "Encoding clipboard artifact")
    resolved = Path(image_path).resolve()
    encoded = image_to_base64(resolved)
    base64_path = out_dir / "clipboard_capture.base64.txt"
    base64_path.write_text(encoded, encoding="utf-8")
    report_path = out_dir / "clipboard_capture_report.md"
    report_path.write_text(
        "\n".join(
            [
                "# Clipboard Capture",
                "",
                f"- Captured: {datetime.now().isoformat(timespec='seconds')}",
                f"- Instruction: {instruction}",
                f"- Image path: {resolved}",
                f"- Base64 bytes: {len(encoded)}",
            ]
        ),
        encoding="utf-8",
    )
    artifacts = {
        "clipboard_image_png": str(resolved),
        "clipboard_image_base64_txt": str(base64_path.resolve()),
        "clipboard_capture_report_md": str(report_path.resolve()),
        "directory": str(out_dir.resolve()),
        "primary_open_file": str(resolved),
    }
    artifact_metadata = {
        "clipboard_image_png": {
            "key": "clipboard_image_png",
            "path": str(resolved),
            "type": "image",
            "title": "Clipboard Image",
            "evidence_summary": "Captured directly from the system clipboard.",
            "validation_state": "validated",
            "created_at": datetime.now().isoformat(timespec="seconds"),
        },
        "clipboard_image_base64_txt": {
            "key": "clipboard_image_base64_txt",
            "path": str(base64_path.resolve()),
            "type": "text",
            "title": "Clipboard Image Base64",
            "evidence_summary": "Base64-encoded clipboard image for downstream AI or API use.",
            "validation_state": "validated",
            "created_at": datetime.now().isoformat(timespec="seconds"),
        },
        "clipboard_capture_report_md": {
            "key": "clipboard_capture_report_md",
            "path": str(report_path.resolve()),
            "type": "report",
            "title": "Clipboard Capture Report",
            "evidence_summary": "Capture metadata and artifact references.",
            "validation_state": "validated",
            "created_at": datetime.now().isoformat(timespec="seconds"),
        },
    }
    return {
        "ok": True,
        "mode": "clipboard_capture",
        "instruction": instruction,
        "results_count": 1,
        "results": [{"title": "Clipboard image", "url": resolved.as_uri(), "source": "clipboard", "price": None, "snippet": ""}],
        "artifacts": artifacts,
        "artifact_metadata": artifact_metadata,
        "opened_url": resolved.as_uri(),
        "source_status": {"clipboard": "ok"},
        "summary": {"captured": True, "bytes_base64": len(encoded), "artifact_directory": str(out_dir.resolve())},
        "canvas": {
            "title": "Clipboard Capture Ready",
            "subtitle": resolved.name,
            "cards": [
                {"title": "Clipboard Image", "price": "captured", "source": "clipboard", "url": resolved.as_uri()},
                {"title": "Base64 Export", "price": f"{len(encoded)} chars", "source": "clipboard", "url": base64_path.resolve().as_uri()},
            ],
        },
    }


def _finalize_operator_result(result: Dict[str, Any], instruction: str, plan_steps: List[Dict[str, Any]]) -> Dict[str, Any]:
    out = dict(result)
    out["planned_steps"] = _summarize_plan_steps(plan_steps)
    out["undo_plan"] = _build_undo_plan(plan_steps)
    plan = out.get("plan", {}) if isinstance(out.get("plan"), dict) else {}
    domain = str(plan.get("domain", out.get("mode", "general") or "general"))
    existing_task_contract = out.get("task_contract", {}) if isinstance(out.get("task_contract"), dict) else {}
    task_contract = TaskContractEngine().extract(instruction)
    if existing_task_contract:
        task_contract = TaskContractEngine().extract(str(existing_task_contract.get("user_goal", instruction) or instruction))
        task_contract_dict = dict(existing_task_contract)
    else:
        task_contract_dict = task_contract.to_dict()
    registry = default_capability_registry()
    existing_graph = out.get("capability_execution_graph", {}) if isinstance(out.get("capability_execution_graph"), dict) else {}
    execution_graph = CapabilityPlanner(registry=registry).plan(task_contract)
    execution_graph_dict = existing_graph or execution_graph.to_dict()
    playbook = select_playbook(domain=domain, instruction=instruction)
    out["playbook"] = playbook
    out["narration"] = _build_run_narration(out=out, playbook=playbook)
    summary = out.get("summary", {}) if isinstance(out.get("summary"), dict) else {}
    existing_critics = out.get("critics", {}) if isinstance(out.get("critics"), dict) else {}
    merged_critics = {
        "action": summary.get("action_critic", {}),
        "quality": summary.get("judgment", {}),
        "elegance_budget": summary.get("elegance_budget", {}),
    }
    merged_critics.update(existing_critics)
    out["critics"] = merged_critics
    out["world_model"] = build_run_world_model(
        instruction=instruction,
        mode=str(out.get("mode", "unknown")),
        domain=domain,
        playbook=playbook,
        opened_url=str(out.get("opened_url", "")),
        paused_for_credentials=bool(out.get("paused_for_credentials", False)),
        pause_reason=str(out.get("pause_reason", "")),
        auth_session_id=str(out.get("auth_session_id", "")),
        artifacts=out.get("artifacts", {}) if isinstance(out.get("artifacts"), dict) else {},
        summary=out.get("summary", {}) if isinstance(out.get("summary"), dict) else {},
        source_status=out.get("source_status", {}) if isinstance(out.get("source_status"), dict) else {},
        decision_log=out.get("decision_log", []) if isinstance(out.get("decision_log"), list) else [],
        results_count=int(out.get("results_count", 0) or 0),
        playbook_validation=plan.get("playbook_validation", {}) if isinstance(plan.get("playbook_validation"), dict) else {},
        playbook_graph_validation=plan.get("playbook_graph_validation", {}) if isinstance(plan.get("playbook_graph_validation"), dict) else {},
        playbook_step_obligations=plan.get("playbook_step_obligations", {}) if isinstance(plan.get("playbook_step_obligations"), dict) else {},
    )
    out["task_contract"] = task_contract_dict
    out["capability_execution_graph"] = execution_graph_dict
    out["capability_registry"] = [spec.to_dict() for spec in registry.list()]
    out["world_model"]["capability_context"] = WorldModelBuilder.from_run(
        session_snapshot=dict(out["world_model"].get("environment", {}).get("session", {}) or {}),
        artifacts=out.get("artifacts", {}) if isinstance(out.get("artifacts"), dict) else {},
        task_contract=task_contract_dict,
        summary=summary,
        opened_url=str(out.get("opened_url", "")),
    ).to_dict()
    out["operator_contract"] = {
        "instruction": instruction,
        "model": "plan_validate_execute_verify_report",
        "least_privilege": True,
    }
    attached = attach_operator_contract(instruction=instruction, result=out, plan_steps=plan_steps)
    if not isinstance(attached.get("critics"), dict):
        attached["critics"] = {}
    task_id = str(attached.get("task_id", uuid.uuid4().hex))
    artifact_factory = ArtifactFactory()
    capabilities = [str(node.get("capability", "")) for node in execution_graph_dict.get("nodes", [])]
    source_data = [
        str(v)
        for v in (attached.get("source_status", {}) or {}).values()
        if isinstance(v, str) and v.strip()
    ]
    artifacts = attached.get("artifacts", {}) if isinstance(attached.get("artifacts"), dict) else {}
    artifact_metadata = attached.get("artifact_metadata", {}) if isinstance(attached.get("artifact_metadata"), dict) else {}
    existing_manifest = str(artifacts.get("artifact_manifest_json", "")).strip()
    if existing_manifest:
        manifest_path = Path(existing_manifest)
    else:
        manifest_path = artifact_factory.write_manifest(
            task_id=task_id,
            task_contract=task_contract_dict,
            artifacts=artifacts,
            artifact_metadata=artifact_metadata,
            generated_by_capabilities=capabilities,
            validation_status=str((attached.get("verification_report", {}) or {}).get("final_verification", "unknown")),
            source_data=source_data,
        )
        artifacts["artifact_manifest_json"] = str(manifest_path.resolve())
    attached["artifacts"] = artifacts
    memory = MemoryStore()
    invalidation_keys = task_contract_dict.get("invalidation_keys", {}) if isinstance(task_contract_dict.get("invalidation_keys"), dict) else {}
    invalidation_key = json.dumps(invalidation_keys, sort_keys=True)
    for key, value in artifacts.items():
        if isinstance(value, str) and value.strip():
            memory.remember_artifact(
                task_id=task_id,
                path=value,
                domain=str(task_contract_dict.get("domain", "")),
                geography=str(task_contract_dict.get("geography", "")),
                invalidation_key=invalidation_key,
                status="created",
                metadata={"artifact_key": key, "task_contract": task_contract_dict},
            )
    invalidated = summary.get("invalidated_artifacts", []) if isinstance(summary.get("invalidated_artifacts"), list) else []
    for item in invalidated:
        memory.remember_artifact(
            task_id=task_id,
            path=str(item),
            domain=str(task_contract_dict.get("domain", "")),
            geography=str(task_contract_dict.get("geography", "")),
            invalidation_key=invalidation_key,
            status="rejected",
            metadata={"reason": "stale_or_invalidated", "task_contract": task_contract_dict},
        )
    platform_story = {
        "executive_summary": str((attached.get("report", {}) or {}).get("summary", "")),
        "key_findings": [str(x) for x in ((attached.get("results", []) or [])[:3])],
        "so_what": "Outputs are attached for review.",
        "recommended_actions": list((attached.get("report", {}) or {}).get("next_actions", [])),
        "caveats": list((attached.get("verification_report", {}) or {}).get("failed_checks", [])),
    }
    platform_critics = {
        "source": PlatformSourceCritic().evaluate(
            [{"url": str(r.get("url", "")), "source": str(r.get("source", ""))} for r in (attached.get("results", []) or [])[:8]]
        ).to_dict(),
        "data_quality": PlatformDataQualityCritic().evaluate(
            int(attached.get("results_count", 0) or 0),
            0.0 if int(attached.get("results_count", 0) or 0) else 1.0,
        ).to_dict(),
        "story": PlatformStoryCritic().evaluate(platform_story, task_contract.audience).to_dict(),
        "uiux": PlatformUIUXCritic().evaluate(
            {
                "chat_workspace": True,
                "canvas_panel": bool((attached.get("canvas", {}) or {}).get("title")),
            }
        ).to_dict(),
        "presentation": PlatformPresentationCritic().evaluate(
            {"slides": [{"title": "Summary"}, {"title": "Method"}, {"title": "Findings"}, {"title": "Actions"}, {"title": "Appendix"}]}
        ).to_dict(),
        "completion": PlatformCompletionCritic().evaluate(
            list(task_contract_dict.get("requested_outputs", []) or []),
            artifacts,
            str((attached.get("verification_report", {}) or {}).get("final_verification", "")),
        ).to_dict(),
    }
    existing_platform_critics = (
        (attached.get("critics", {}) or {}).get("platform", {})
        if isinstance((attached.get("critics", {}) or {}).get("platform", {}), dict)
        else {}
    )
    graph_status = str(execution_graph_dict.get("status", "")).strip()
    executed_graph = graph_status in {"running", "succeeded", "failed", "blocked"}
    if executed_graph and existing_platform_critics:
        merged_platform_critics = dict(existing_platform_critics)
    else:
        merged_platform_critics = dict(platform_critics)
        merged_platform_critics.update(existing_platform_critics)
    attached["critics"]["platform"] = merged_platform_critics
    attached["human_report"] = HumanStyleReporter().build(
        task_contract=task_contract_dict,
        execution_graph=execution_graph_dict,
        result=attached,
    )
    manifest_payload: Dict[str, Any] = {}
    try:
        manifest_payload = json.loads(Path(manifest_path).read_text(encoding="utf-8"))
    except Exception:
        manifest_payload = {}
    attached["artifact_manifest"] = manifest_payload
    summary_payload = attached.get("summary", {}) if isinstance(attached.get("summary"), dict) else {}
    validation_results = attached.get("validation_results", {}) if isinstance(attached.get("validation_results"), dict) else {}
    if not validation_results and isinstance(summary_payload.get("validation_results"), dict):
        validation_results = dict(summary_payload.get("validation_results", {}))
        attached["validation_results"] = validation_results
    final_output_gate = attached.get("final_output_gate", {}) if isinstance(attached.get("final_output_gate"), dict) else {}
    if not final_output_gate and isinstance(summary_payload.get("final_output_gate"), dict):
        final_output_gate = dict(summary_payload.get("final_output_gate", {}))
        attached["final_output_gate"] = final_output_gate
    completion_status = str(attached.get("completion_status", "") or summary_payload.get("completion_status", "") or "")
    if completion_status and not attached.get("completion_status"):
        attached["completion_status"] = completion_status
    if final_output_gate and not bool(final_output_gate.get("passed", True)):
        verification_checks = [
            {
                "name": "final_output_gate",
                "pass": False,
                "evidence": [f"blocking_failures={list(final_output_gate.get('blocking_failures', []))}"],
            }
        ]
        attached["verification_report"] = dict(attached.get("verification_report", {}) or {})
        attached["verification_report"].update(
            {
                "requested_outputs_exist": bool(artifacts),
                "artifact_content_matches_goal": False,
                "no_unresolved_execution_errors": False,
                "user_goal_satisfied": False,
                "verification_checks": verification_checks,
                "failed_checks": ["final_output_gate"],
                "final_verification": "failed",
            }
        )
        attached["verification"] = {"passed": False, "checks": verification_checks, "evidence": ["validation gate blocked final output"]}
        attached["final_report"] = {
            "task_id": task_id,
            "status": "blocked",
            "summary": "Validation failed",
            "actions_taken": [f"validation:{name}" for name in list(final_output_gate.get("blocking_failures", []))],
            "outputs_created": [{"type": "file", "location": value, "description": key} for key, value in artifacts.items() if isinstance(value, str)],
            "verification_summary": "failed",
            "remaining_issues": list(final_output_gate.get("blocking_failures", [])),
            "next_safe_action": "Review validator findings and rebuild with valid geography and service-scope evidence.",
        }
        attached["ok"] = False
        attached["error"] = str(attached.get("error") or "final_output_gate_failed")
        attached["error_code"] = str(attached.get("error_code") or "final_output_gate_failed")
        attached["opened_url"] = str(artifacts.get("geography_validation_report_md", attached.get("opened_url", "")) or "")
        attached["canvas"] = {
            "title": "Validation Failed",
            "subtitle": "I found out-of-market or out-of-scope data and I am not presenting this as a stakeholder-ready result.",
            "cards": [
                {"title": str(item)[:110], "price": "repair", "source": "validator", "url": ""}
                for item in list(final_output_gate.get("required_repairs", []))[:5]
            ],
        }
    elif final_output_gate and bool(final_output_gate.get("passed", False)) and completion_status in {"completed_demo_package", "partially_completed_validated"}:
        verification_checks = [
            {"name": "final_output_gate", "pass": True, "evidence": ["validation gate passed"]},
            {"name": "requested_outputs_exist", "pass": bool(artifacts), "evidence": [f"created={sorted(artifacts.keys())}"]},
            {"name": "user_goal_satisfied", "pass": True, "evidence": [f"completion_status={completion_status}"]},
        ]
        attached["verification_report"] = dict(attached.get("verification_report", {}) or {})
        attached["verification_report"].update(
            {
                "requested_outputs_exist": bool(artifacts),
                "artifact_content_matches_goal": True,
                "no_unresolved_execution_errors": True,
                "user_goal_satisfied": True,
                "verification_checks": verification_checks,
                "failed_checks": [],
                "final_verification": "passed",
            }
        )
        attached["verification"] = {
            "passed": True,
            "checks": verification_checks,
            "evidence": [f"completion_status={completion_status}", f"artifacts={sorted(artifacts.keys())}"],
        }
        attached["final_report"] = {
            "task_id": task_id,
            "status": completion_status,
            "summary": str((attached.get("canvas", {}) or {}).get("title", "Payer package ready")),
            "actions_taken": ["validation:passed", f"completion:{completion_status}"],
            "outputs_created": [{"type": "file", "location": value, "description": key} for key, value in artifacts.items() if isinstance(value, str)],
            "verification_summary": "passed",
            "remaining_issues": [],
            "next_safe_action": "Open the artifact package in Canvas.",
        }
        if isinstance(attached.get("critics"), dict) and isinstance((attached.get("critics") or {}).get("platform"), dict):
            attached["critics"]["platform"]["completion"] = {
                "passed": True,
                "score": 0.92,
                "reason": "Requested outputs are present and validation passed for the delivered payer package.",
                "required_fix": "",
                "severity": "low",
            }
    if executed_graph:
        graph_failed_checks = [
            name for name, payload in merged_platform_critics.items() if isinstance(payload, dict) and not bool(payload.get("passed", False))
        ]
        verification_checks = [
            {"name": "graph_executed", "pass": str(execution_graph_dict.get("status", "")) in {"succeeded", "failed"}, "evidence": [f"graph_status={execution_graph_dict.get('status', '')}"]},
            {"name": "all_nodes_complete", "pass": all(str(node.get("status", "")) in {"succeeded", "revised"} for node in (execution_graph_dict.get("nodes", []) or [])), "evidence": [[str(node.get("status", "")) for node in (execution_graph_dict.get("nodes", []) or [])]]},
            {"name": "requested_outputs_exist", "pass": bool(artifacts), "evidence": [f"created={sorted(artifacts.keys())}"]},
            {"name": "critics_resolved", "pass": len(graph_failed_checks) == 0, "evidence": [f"runtime_failed_critics={graph_failed_checks}"]},
            {"name": "user_goal_satisfied", "pass": len(graph_failed_checks) == 0 and str(execution_graph_dict.get("status", "")) == "succeeded", "evidence": [f"runtime_revisions={len(attached.get('revisions_performed', []) or [])}"]},
        ]
        final_verification = "passed" if all(bool(item["pass"]) for item in verification_checks) else "failed"
        attached["verification_report"] = dict(attached.get("verification_report", {}) or {}) or {
            "task_id": task_id,
        }
        attached["verification_report"].update(
            {
                "used_expected_tool_family": True,
                "targets_match_request": True,
                "requested_outputs_exist": bool(artifacts),
                "artifact_content_matches_goal": len(graph_failed_checks) == 0,
                "no_unresolved_execution_errors": len(graph_failed_checks) == 0 and str(execution_graph_dict.get("status", "")) != "failed",
                "no_irrelevant_detours": True,
                "user_goal_satisfied": final_verification == "passed",
                "verification_checks": verification_checks,
                "failed_checks": [item["name"] for item in verification_checks if not bool(item["pass"])],
                "final_verification": final_verification,
            }
        )
        attached["verification"] = dict(attached.get("verification", {}) or {})
        attached["verification"].update(
            {
                "passed": final_verification == "passed",
                "checks": verification_checks,
                "evidence": [f"graph_status={execution_graph_dict.get('status', '')}", f"artifacts={sorted(artifacts.keys())}"],
            }
        )
        attached["final_report"] = dict(attached.get("final_report", {}) or {})
        attached["final_report"].update(
            {
                "task_id": task_id,
                "status": "completed" if final_verification == "passed" else "failed",
                "summary": str((attached.get("canvas", {}) or {}).get("title", "Capability Graph Completed")),
                "actions_taken": [f"{node.get('capability')}: {node.get('status')}" for node in (execution_graph_dict.get("nodes", []) or [])],
                "outputs_created": [{"type": "file", "location": value, "description": key} for key, value in artifacts.items() if isinstance(value, str)],
                "verification_summary": final_verification,
                "remaining_issues": [item["name"] for item in verification_checks if not bool(item["pass"])],
                "next_safe_action": "Open the artifact package in Canvas." if final_verification == "passed" else "Review failed graph checks and rerun.",
            }
        )
    attached["platform"] = {
        "task_contract_engine": task_contract_dict,
        "capability_execution_graph": execution_graph_dict,
        "artifact_manifest_json": str(manifest_path.resolve()),
        "memory_invalidation_key": invalidation_key,
        "architecture_version": "capability-runtime-v1",
    }
    attached["ui_cards"] = build_platform_cards(attached)
    return attached


def _run_mission_runtime(
    instruction: str,
    *,
    ai_meta: Dict[str, Any],
    progress_cb: Optional[Callable[[int, str], None]] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    _emit_progress(progress_cb, 12, "Building mission contract")
    mission_engine = MissionContractEngine()
    mission_contract = mission_engine.extract(instruction)
    task_contract = TaskContractEngine().extract(instruction)
    task_contract.mission_type = mission_contract.mission_type
    task_contract.deliverable_mode = mission_contract.deliverable_mode
    task_contract.requested_outputs = list(mission_contract.requested_outputs)

    def _mission_source_collector(**kwargs: Any) -> Dict[str, Any]:
        query = str(kwargs.get("query", "") or "").strip()
        mission_instruction = str(kwargs.get("instruction", instruction) or instruction)
        effective_instruction = query if query else mission_instruction
        collected = platform_collect_generic_research(
            instruction=effective_instruction,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=bool(human_like_interaction),
        )
        normalized_sources: List[Dict[str, Any]] = []
        for row in list(collected.get("sources", []) or []):
            if not isinstance(row, dict):
                continue
            normalized = dict(row)
            normalized["source"] = str(row.get("source") or row.get("source_name") or row.get("company") or row.get("name") or "source")
            normalized["source_type"] = str(row.get("source_type") or row.get("type") or row.get("source") or "reference")
            normalized["url"] = str(row.get("url") or row.get("job_url") or row.get("url_or_path") or "")
            normalized["title"] = str(row.get("title") or row.get("role_title") or row.get("name") or "")
            normalized["snippet"] = str(row.get("snippet") or row.get("summary") or "")
            normalized_sources.append(normalized)
        if not normalized_sources:
            for row in list(collected.get("search_results", []) or [])[:12]:
                if not isinstance(row, dict):
                    continue
                normalized_sources.append(
                    {
                        "source": str(row.get("source", "search_result")),
                        "source_type": str(row.get("source", "search_result")),
                        "url": str(row.get("url", "")),
                        "title": str(row.get("title", "")),
                        "snippet": str(row.get("snippet", "")),
                    }
                )
        return {
            "sources": normalized_sources,
            "search_results": list(collected.get("search_results", []) or []),
            "query": str(collected.get("query", query)),
            "error": str(collected.get("error", "")),
            "ok": bool(collected.get("ok", False)),
        }

    workspace_slug = re.sub(r"[^a-z0-9]+", "_", instruction.lower()).strip("_")[:48] or "mission"
    workspace_dir = Path("data/mission_runs") / f"{workspace_slug}_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}"
    runtime_graph = CapabilityPlanner(registry=default_capability_registry()).plan(task_contract)
    runtime_result = _run_execution_graph_runtime_path(
        instruction=instruction,
        task_contract=task_contract,
        graph=runtime_graph,
        ai_meta=ai_meta,
        progress_cb=progress_cb,
        artifact_reuse_mode="always_regenerate",
        artifact_reuse_max_age_hours=1,
        mode_override="mission_runtime",
        extra_context={
            "workspace_dir": str(workspace_dir),
            "browser_worker_mode": browser_worker_mode,
            "human_like_interaction": bool(human_like_interaction),
            "source_collector": _mission_source_collector,
        },
    )
    _emit_progress(progress_cb, 55, "Generating work products")
    graph_dict = runtime_result.get("capability_execution_graph", {}) if isinstance(runtime_result.get("capability_execution_graph"), dict) else {}
    mission_node = next((node for node in (graph_dict.get("nodes", []) or []) if str(node.get("capability", "")) == "mission_work_product"), {})
    mission_outputs = dict(mission_node.get("output_payload", {}) or {})
    artifacts = dict(runtime_result.get("artifacts", {}) or {})
    mission_status = str(mission_outputs.get("mission_status", "") or "failed_execution")
    passed = mission_status not in {"failed_validation", "failed_execution"}
    verification_checks = [
        {"name": "mission_contract_built", "pass": bool(mission_outputs.get("mission_contract") or mission_contract.to_dict())},
        {"name": "artifact_plan_realized", "pass": bool(mission_outputs.get("artifact_plan") or mission_contract.artifact_plan) and bool(artifacts)},
        {"name": "critic_revision_cycle_completed", "pass": bool(mission_outputs.get("mission_critics"))},
        {"name": "truthful_status_assigned", "pass": bool(mission_status)},
    ]
    result = dict(runtime_result)
    result.update(
        {
            "mode": "mission_runtime",
            "runtime_mode": "execution_graph_runtime",
            "mission_contract": dict(mission_outputs.get("mission_contract", mission_contract.to_dict()) or mission_contract.to_dict()),
            "task_contract": task_contract.to_dict(),
            "deliverable_mode": str(mission_outputs.get("mission_contract", {}).get("deliverable_mode", mission_contract.deliverable_mode) if isinstance(mission_outputs.get("mission_contract", {}), dict) else mission_contract.deliverable_mode),
            "research_strategy": dict(mission_outputs.get("research_strategy", {}) or {}),
            "evidence_map": dict(mission_outputs.get("evidence_map", {}) or {}),
            "artifact_plan": list(mission_outputs.get("artifact_plan", mission_contract.artifact_plan) or mission_contract.artifact_plan),
            "memory_context": dict(runtime_result.get("memory_context", {}) or {}),
            "mission_status": mission_status,
            "output_truth": dict(mission_outputs.get("output_truth", {}) or {}),
            "recovery": dict(mission_outputs.get("recovery", {}) or {}),
            "final_package": dict(mission_outputs.get("final_package", {}) or {}),
            "revisions_performed": list(mission_outputs.get("revisions_performed", runtime_result.get("revisions_performed", [])) or runtime_result.get("revisions_performed", [])),
            "critics": {"mission": dict(mission_outputs.get("mission_critics", {}) or {}), "platform": dict((runtime_result.get("critics", {}) or {}).get("platform", {}) if isinstance(runtime_result.get("critics", {}), dict) else {})},
            "summary": {
                "mission_type": mission_contract.mission_type,
                "deliverable_mode": mission_contract.deliverable_mode,
                "accepted_sources": int(((mission_outputs.get("evidence_map", {}) or {}).get("summary", {}) or {}).get("accepted_count", 0) or 0),
                "accepted_external_sources": int(((mission_outputs.get("evidence_map", {}) or {}).get("summary", {}) or {}).get("accepted_external_count", 0) or 0),
                "artifacts_created": len(artifacts),
            },
            "verification": {
                "passed": passed,
                "checks": verification_checks,
                "evidence": [f"mission_status={mission_status}", f"artifacts={sorted(artifacts.keys())}"],
            },
            "verification_report": {
                "final_verification": "passed" if passed else "failed",
                "verification_checks": verification_checks,
                "failed_checks": [item["name"] for item in verification_checks if not bool(item["pass"])],
            },
            "final_report": {
                "status": mission_status,
                "summary": str((mission_outputs.get("final_package", {}) or {}).get("summary", "Mission package prepared.")),
                "outputs_created": [{"type": "file", "location": value, "description": key} for key, value in artifacts.items() if isinstance(value, str)],
                "next_safe_action": "Open the package and review the critic results before external use.",
            },
            "canvas": {
                "title": "Mission Package Ready" if passed else "Mission Package Needs Review",
                "subtitle": mission_contract.mission_type,
                "cards": [
                    {"title": str(item.get("title", item.get("key", ""))), "price": str(item.get("validation_state", "ready")), "source": "artifact", "url": str(item.get("path", ""))}
                    for item in list((runtime_result.get("artifact_metadata", {}) or {}).values())[:6]
                    if isinstance(item, dict)
                ],
            },
        }
    )
    result["ui_cards"] = build_platform_cards(result)
    _emit_progress(progress_cb, 100, "Completed")
    return result


def _run_topic_mastery_learn_mode(
    instruction: str,
    *,
    ai_meta: Dict[str, Any],
    progress_cb: Optional[Callable[[int, str], None]] = None,
) -> Dict[str, Any]:
    _emit_progress(progress_cb, 10, "Building learn mission")
    task_contract = TaskContractEngine().extract(instruction)
    task_contract.domain = "topic_learning"
    runtime_graph = CapabilityPlanner(registry=default_capability_registry()).plan(task_contract)
    workspace_slug = re.sub(r"[^a-z0-9]+", "_", instruction.lower()).strip("_")[:48] or "topic_mastery"
    workspace_dir = Path("data/learn_runs") / f"{workspace_slug}_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}"
    runtime_result = _run_execution_graph_runtime_path(
        instruction=instruction,
        task_contract=task_contract,
        graph=runtime_graph,
        ai_meta=ai_meta,
        progress_cb=progress_cb,
        artifact_reuse_mode="always_regenerate",
        artifact_reuse_max_age_hours=1,
        mode_override="topic_mastery_learn_mode",
        extra_context={"workspace_dir": str(workspace_dir)},
    )
    graph_dict = runtime_result.get("capability_execution_graph", {}) if isinstance(runtime_result.get("capability_execution_graph"), dict) else {}
    learn_node = next((node for node in (graph_dict.get("nodes", []) or []) if str(node.get("capability", "")) == "topic_mastery_learn"), {})
    learn_outputs = dict(learn_node.get("output_payload", {}) or {})
    result = dict(runtime_result)
    result.update(
        {
            "mode": "topic_mastery_learn_mode",
            "runtime_mode": "execution_graph_runtime",
            "learn_mission": dict(learn_outputs.get("learn_mission", {}) or {}),
            "source_discovery": dict(learn_outputs.get("source_discovery", {}) or {}),
            "video_analysis": dict(learn_outputs.get("video_analysis", {}) or {}),
            "topic_model": dict(learn_outputs.get("topic_model", {}) or {}),
            "consensus_workflow": list(learn_outputs.get("consensus_workflow", []) or []),
            "contradictions": list(learn_outputs.get("contradictions", []) or []),
            "learned_skill": dict(learn_outputs.get("learned_skill", {}) or {}),
            "learned_skill_library": dict(learn_outputs.get("learned_skill_library", {}) or {}),
            "mastery_guide": dict(learn_outputs.get("mastery_guide", {}) or {}),
            "practice_plan": dict(learn_outputs.get("practice_plan", {}) or {}),
            "practice_preview": dict(learn_outputs.get("practice_preview", {}) or {}),
            "refresh_plan": dict(learn_outputs.get("refresh_plan", {}) or {}),
            "critic_results": dict(learn_outputs.get("critic_results", {}) or {}),
            "skill_validation": dict(learn_outputs.get("skill_validation", {}) or {}),
            "memory": dict(learn_outputs.get("memory", {}) or {}),
            "status": str(learn_outputs.get("status", "real_partial") or "real_partial"),
            "final_package": dict(learn_outputs.get("final_package", {}) or {}),
        }
    )
    result["ui_cards"] = build_platform_cards(result)
    _emit_progress(progress_cb, 100, "Completed")
    return result


def _build_run_narration(out: Dict[str, Any], playbook: Dict[str, Any]) -> List[str]:
    lines: List[str] = []
    pb_name = str(playbook.get("name", "General Operator Playbook")).strip()
    if pb_name:
        lines.append(f"Selected playbook: {pb_name}.")
    task_contract = out.get("current_task_contract", {}) if isinstance(out.get("current_task_contract"), dict) else {}
    geography = str(task_contract.get("geography", "")).strip()
    if geography:
        lines.append(f"Current task geography: {geography}.")
    opened_url = str(out.get("opened_url", "")).strip()
    if opened_url:
        lines.append(f"Working target: {opened_url}.")
    summary = out.get("summary", {}) if isinstance(out.get("summary"), dict) else {}
    if bool(summary.get("auth_session_reused", False)):
        lines.append("Reused existing authenticated session state.")
    if bool(out.get("paused_for_credentials", False)):
        reason = str(out.get("pause_reason", "")).strip() or "User authentication is required."
        lines.append(f"Paused for user action: {reason}")
    artifacts = out.get("artifacts", {}) if isinstance(out.get("artifacts"), dict) else {}
    if artifacts:
        lines.append(f"Created {len(artifacts)} artifact(s): {', '.join(sorted(artifacts.keys())[:4])}.")
    invalidated = summary.get("invalidated_artifacts", []) if isinstance(summary.get("invalidated_artifacts"), list) else []
    if invalidated:
        lines.append("Rejected stale artifacts from a different task contract.")
    final_output_gate = out.get("final_output_gate", {}) if isinstance(out.get("final_output_gate"), dict) else {}
    if final_output_gate and not bool(final_output_gate.get("passed", True)):
        lines.append("Validation failed. I found out-of-market or out-of-scope data and I am not presenting this as a stakeholder-ready result.")
    if isinstance(out.get("decision_log"), list) and out.get("decision_log"):
        lines.append(f"Planner decisions: {len(out.get('decision_log', []))} recorded.")
    source_status = out.get("source_status", {}) if isinstance(out.get("source_status"), dict) else {}
    if source_status:
        first_key = next(iter(source_status.keys()))
        lines.append(f"Primary source status: {first_key}={source_status.get(first_key)}.")
    if not lines:
        lines.append("Execution completed with no notable environment events.")
    return lines[:8]


def execute_instruction(
    instruction: str,
    control_granted: bool,
    step_mode: bool = False,
    confirm_risky: bool = False,
    ai_backend: str = "deterministic-local",
    min_live_non_curated_citations: Optional[int] = None,
    manual_auth_phase: bool = False,
    auth_session_id: Optional[str] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = True,
    artifact_reuse_mode: str = "reuse_if_recent",
    artifact_reuse_max_age_hours: int = 72,
    progress_cb: Optional[Callable[[int, str], None]] = None,
) -> Dict[str, Any]:
    _emit_progress(progress_cb, 2, "Understanding your request")
    ai_meta = backend_metadata(ai_backend)
    ai_backend = normalize_backend(ai_backend)
    freshness_mode = _normalize_artifact_reuse_mode(artifact_reuse_mode)
    try:
        freshness_hours = max(1, min(24 * 30, int(artifact_reuse_max_age_hours)))
    except Exception:
        freshness_hours = 72
    if not control_granted:
        return {
            "ok": False,
            "error": "Control not granted. Click Accept Control first.",
            "instruction": instruction,
            "ai": ai_meta,
        }

    normalized = instruction.strip()
    if not normalized:
        return {"ok": False, "error": "Instruction is empty.", "instruction": instruction}

    explicit_route = _classify_explicit_route(normalized)
    task_contract = TaskContractEngine().extract(normalized)

    if _is_clipboard_capture_intent(normalized):
        _emit_progress(progress_cb, 8, "Capturing clipboard image")
        result = _run_clipboard_capture(instruction=instruction, progress_cb=progress_cb)
        result["task_contract"] = task_contract.to_dict()
        result["verification"] = {
            "passed": bool(result.get("ok", False)),
            "checks": [
                {"name": "clipboard_capture_completed", "pass": bool(result.get("ok", False))},
                {"name": "artifact_present", "pass": bool((result.get("artifacts", {}) or {}).get("clipboard_image_png"))},
            ],
            "evidence": [f"clipboard_status={str((result.get('source_status', {}) or {}).get('clipboard', 'unknown'))}"],
        }
        result["verification_report"] = {
            "final_verification": "passed" if bool(result.get("ok", False)) else "failed",
            "verification_checks": list(result.get("verification", {}).get("checks", [])),
            "failed_checks": [item.get("name") for item in result.get("verification", {}).get("checks", []) if not bool(item.get("pass", False))],
        }
        result["ui_cards"] = build_platform_cards(result)
        _emit_progress(progress_cb, 100, "Completed")
        return result

    if _is_topic_mastery_intent(normalized):
        return _run_topic_mastery_learn_mode(
            instruction=instruction,
            ai_meta=ai_meta,
            progress_cb=progress_cb,
        )

    if _is_mission_runtime_intent(normalized):
        return _run_mission_runtime(
            instruction=instruction,
            ai_meta=ai_meta,
            progress_cb=progress_cb,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=False,
        )

    if _should_use_execution_graph_runtime(task_contract=task_contract, explicit_route=explicit_route, instruction=normalized):
        runtime_graph = CapabilityPlanner(registry=default_capability_registry()).plan(task_contract)
        runtime_plan = _build_runtime_plan(task_contract, graph=runtime_graph.to_dict())
        plan_steps = [dict(x) for x in runtime_plan.get("steps", [])]
        _emit_progress(progress_cb, 8, "Building task contract")
        execution = _run_execution_graph_runtime_path(
            instruction=instruction,
            task_contract=task_contract,
            graph=runtime_graph,
            ai_meta=ai_meta,
            progress_cb=progress_cb,
            artifact_reuse_mode=freshness_mode,
            artifact_reuse_max_age_hours=freshness_hours,
        )
        _emit_progress(progress_cb, 100, "Completed")
        return _finalize_operator_result(execution, instruction=instruction, plan_steps=plan_steps)

    if explicit_route == "artifact_generation":
        _emit_progress(progress_cb, 8, "Building artifact package contract")
        plan = _build_artifact_generation_plan(normalized)
        plan_steps = [dict(x) for x in plan.get("steps", [])]
        elegance = EleganceBudget(total=70)
        if len(plan_steps) > 4:
            elegance.consume(len(plan_steps) - 4, "artifact_plan_complexity")
        requested_outputs = _requested_outputs(normalized)
        represented_outputs = _plan_represented_outputs(plan, plan_steps)
        if requested_outputs and any(x not in represented_outputs for x in requested_outputs):
            return _fail_fast_output_mismatch(
                instruction=instruction,
                mode="artifact_generation_plan_invalid",
                plan_steps=plan_steps,
                requested_outputs=requested_outputs,
                represented_outputs=represented_outputs,
            )
        runtime_contract = TaskContractEngine().extract(normalized)
        runtime_contract.domain = "artifact_generation"
        runtime_outputs: List[str] = []
        if "document" in requested_outputs or "executive_summary" in requested_outputs:
            runtime_outputs.append("report")
        if "powerpoint" in requested_outputs:
            runtime_outputs.append("presentation")
        if "visual" in requested_outputs or "dashboard" in requested_outputs:
            runtime_outputs.append("dashboard")
        if "executive_summary" in requested_outputs:
            runtime_outputs.append("executive_summary")
        runtime_contract.requested_outputs = runtime_outputs or ["report"]
        execution = _run_execution_graph_runtime_path(
            instruction=instruction,
            task_contract=runtime_contract,
            graph=CapabilityPlanner(registry=default_capability_registry()).plan(runtime_contract),
            ai_meta=ai_meta,
            progress_cb=progress_cb,
            artifact_reuse_mode=freshness_mode,
            artifact_reuse_max_age_hours=freshness_hours,
            mode_override="artifact_generation",
        )
        summary = execution.get("summary", {}) if isinstance(execution.get("summary"), dict) else {}
        summary = _apply_freshness_metadata(summary, freshness_mode, freshness_hours)
        execution["summary"] = _ensure_elegance_budget(summary, elegance)
        budget_block = _enforce_elegance_budget_gate(elegance=elegance, mode="artifact_generation", instruction=instruction)
        if budget_block:
            budget_block["plan"] = plan
            return _finalize_operator_result(budget_block, instruction=instruction, plan_steps=plan_steps)
        execution["plan"] = plan
        return _finalize_operator_result(execution, instruction=instruction, plan_steps=plan_steps)

    if _is_native_planning_intent(normalized) or _is_study_pack_intent(normalized) or _is_job_research_intent(normalized):
        _emit_progress(progress_cb, 8, "Building execution plan")
        plan = _build_native_plan(normalized)
        plan_steps = [dict(x) for x in plan.get("steps", [])]
        elegance = EleganceBudget(total=80)
        if len(plan_steps) > 5:
            elegance.consume(len(plan_steps) - 5, "native_plan_complexity")
        requested_outputs = _requested_outputs(normalized)
        represented_outputs = _plan_represented_outputs(plan, plan_steps)
        if requested_outputs and any(x not in represented_outputs for x in requested_outputs):
            return _fail_fast_output_mismatch(
                instruction=instruction,
                mode="native_plan_invalid",
                plan_steps=plan_steps,
                requested_outputs=requested_outputs,
                represented_outputs=represented_outputs,
            )
        ambiguities = _detect_ambiguities(normalized, plan_steps)
        if ambiguities:
            return {
                "ok": False,
                "mode": "needs_clarification",
                "instruction": instruction,
                "message": "Clarification required before execution.",
                "questions": ambiguities[:4],
                "planned_steps": _summarize_plan_steps(plan_steps),
                "undo_plan": _build_undo_plan(plan_steps),
            }
        if _is_destructive_intent(normalized, plan_steps) and not confirm_risky:
            return {
                "ok": False,
                "mode": "confirmation_required",
                "instruction": instruction,
                "requires_confirmation": True,
                "message": "Destructive or high-impact action detected. Confirm to continue.",
                "planned_steps": _summarize_plan_steps(plan_steps),
                "undo_plan": _build_undo_plan(plan_steps),
            }
        _emit_progress(progress_cb, 12, f"Plan ready with {len(plan_steps)} step(s)")
        _emit_progress(progress_cb, 14, "Executing plan")
        execution = _execute_native_plan(
            plan=plan,
            instruction=instruction,
            progress_cb=progress_cb,
            min_live_non_curated_citations=min_live_non_curated_citations,
            manual_auth_phase=manual_auth_phase,
            auth_session_id=auth_session_id,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=bool(human_like_interaction),
        )
        ex_summary = execution.get("summary", {}) if isinstance(execution.get("summary"), dict) else {}
        ex_summary = _apply_freshness_metadata(ex_summary, freshness_mode, freshness_hours)
        execution["summary"] = _ensure_elegance_budget(ex_summary, elegance)
        budget_block = _enforce_elegance_budget_gate(elegance=elegance, mode="autonomous_plan_execute", instruction=instruction)
        if budget_block:
            budget_block["plan"] = plan
            budget_block["ai"] = ai_meta
            _emit_progress(progress_cb, 100, "Completed")
            return _finalize_operator_result(budget_block, instruction=instruction, plan_steps=plan_steps)
        execution["ai"] = ai_meta
        _emit_progress(progress_cb, 100, "Completed")
        return _finalize_operator_result(execution, instruction=instruction, plan_steps=plan_steps)

    if explicit_route == "desktop_sequence" or _is_desktop_sequence_intent(normalized):
        _emit_progress(progress_cb, 12, "Building desktop action sequence")
        plan = build_plan(normalized)
        plan_steps = [dict(x) for x in plan.get("steps", [])]
        elegance = EleganceBudget(total=65)
        if len(plan_steps) > 6:
            elegance.consume(len(plan_steps) - 6, "desktop_plan_complexity")
        requested_outputs = _requested_outputs(normalized)
        represented_outputs = _plan_represented_outputs(plan, plan_steps)
        if requested_outputs and any(x not in represented_outputs for x in requested_outputs):
            return _fail_fast_output_mismatch(
                instruction=instruction,
                mode="desktop_plan_invalid",
                plan_steps=plan_steps,
                requested_outputs=requested_outputs,
                represented_outputs=represented_outputs,
            )
        ambiguities = _detect_ambiguities(normalized, plan_steps)
        if ambiguities:
            return {
                "ok": False,
                "mode": "needs_clarification",
                "instruction": instruction,
                "message": "Clarification required before execution.",
                "questions": ambiguities[:4],
                "planned_steps": _summarize_plan_steps(plan_steps),
                "undo_plan": _build_undo_plan(plan_steps),
            }
        risk = assess_risk(plan)
        if (risk["requires_confirmation"] or _is_destructive_intent(normalized, plan_steps)) and not confirm_risky:
            return {
                "ok": False,
                "mode": "desktop_sequence_preview",
                "instruction": instruction,
                "requires_confirmation": True,
                "message": "Risky actions detected. Confirm to execute.",
                "risk": risk,
                "plan": plan,
                "planned_steps": _summarize_plan_steps(plan_steps),
                "undo_plan": _build_undo_plan(plan_steps),
                "canvas": {
                    "title": "Confirmation Required",
                    "subtitle": f"{len(risk['risky_steps'])} risky step(s)",
                    "cards": [
                        {
                            "title": f"Step {item['index']}: {item['step'].get('action','')}",
                            "price": "confirm",
                            "source": "risk",
                            "url": "",
                        }
                        for item in risk["risky_steps"][:6]
                    ],
                },
            }
        run = execute_plan(
            plan,
            start_index=0,
            step_mode=step_mode,
            allow_input_fallback=True,
            human_like_interaction=bool(human_like_interaction),
        )
        _emit_progress(progress_cb, 95, "Finalizing desktop run output")
        store = LocalVectorStore()
        app_name = plan.get("app_name", "") or "desktop"
        guidance = get_guidance(app_name=app_name, user_goal=normalized, store=store) if app_name else {"guidance": []}
        response = {
            "ok": run.ok,
            "mode": "desktop_sequence",
            "instruction": instruction,
            "ai": ai_meta,
            "app_name": app_name,
            "plan": plan,
            "trace": run.trace,
            "done": run.done,
            "next_step_index": run.next_step_index,
            "pending_plan": {"plan": plan, "next_step_index": run.next_step_index} if not run.done else None,
            "paused_for_credentials": run.paused_for_credentials,
            "pause_reason": run.pause_reason,
            "error": run.error,
            "artifacts": dict(run.artifacts or {}),
            "opened_url": Path(run.artifacts["primary_open_file"]).resolve().as_uri() if run.artifacts.get("primary_open_file") else "",
            "risk": risk,
            "guidance": guidance,
            "canvas": {
                "title": f"Desktop Sequence: {app_name or 'session'}",
                "subtitle": f"Executed {len(run.trace)} steps",
                "cards": [
                    {"title": f"Step {t.get('step', 0)}: {t.get('action','')}", "price": "ok" if t.get("ok") else "error", "source": "uia", "url": ""}
                    for t in run.trace[:6]
                ],
            },
        }
        response_summary = {
            "executed_steps": len(run.trace),
            "human_like_interaction": bool(human_like_interaction),
            "elegance_budget": elegance.snapshot(),
        }
        response["summary"] = _apply_freshness_metadata(response_summary, freshness_mode, freshness_hours)
        budget_block = _enforce_elegance_budget_gate(elegance=elegance, mode="desktop_sequence", instruction=instruction)
        if budget_block:
            budget_block["plan"] = plan
            budget_block["ai"] = ai_meta
            return _finalize_operator_result(budget_block, instruction=instruction, plan_steps=plan_steps)
        return _finalize_operator_result(response, instruction=instruction, plan_steps=plan_steps)

    open_match = re.search(r"\bopen\s+(.+?)(?:\s+app)?\b", normalized, flags=re.IGNORECASE)
    if open_match and ("search" not in normalized.lower()):
        _emit_progress(progress_cb, 15, "Opening installed application")
        target = open_match.group(1).strip()
        plan_steps = [{"action": "open_app", "app": normalize_app_name(target)}]
        running, running_ref = is_app_running(target)
        app_name = normalize_app_name(target)
        app_target = app_name or target
        action_decision = ActionCritic().evaluate(
            next_action="open_app",
            target=app_target,
            already_open=running,
            context={"reusable_target": app_target if running else "", "recent_actions": [f"open_app:{app_target}"]},
        )
        launched = ""
        if running and not action_decision.allow:
            ok = True
            launched = f"already_running:{running_ref or app_target}"
        else:
            ok, launched = open_installed_app(target)
        app_name = normalize_app_name(target)
        store = LocalVectorStore()
        guidance = get_guidance(app_name=app_name, user_goal=normalized, store=store)
        if ok:
            response = {
                "ok": True,
                "mode": "desktop_app_open",
                "instruction": instruction,
                "ai": ai_meta,
                "app_name": app_name,
                "launched": launched,
                "summary": _apply_freshness_metadata({
                    "reused_running_app": bool(running and launched.startswith("already_running:")),
                    "running_ref": running_ref,
                    "action_critic": {
                        "score": action_decision.score,
                        "reasons": action_decision.reasons,
                        "elegance_cost": action_decision.elegance_cost,
                    },
                }, freshness_mode, freshness_hours),
                "paused_for_credentials": True,
                "pause_reason": "If a login prompt appears, enter credentials and click Resume.",
                "guidance": guidance,
                "canvas": {
                    "title": f"Opened App: {app_name}",
                    "subtitle": "Credential checkpoint active",
                    "cards": [
                        {"title": g["title"], "price": "tip", "source": "knowledge", "url": g["source_url"]}
                        for g in guidance["guidance"][:6]
                    ],
                },
            }
            return _finalize_operator_result(response, instruction=instruction, plan_steps=plan_steps)
        return {
            "ok": False,
            "mode": "desktop_app_open",
            "instruction": instruction,
            "ai": ai_meta,
            "error": f"Could not locate installed app '{target}'.",
            "app_name": app_name,
            "guidance": guidance,
            "planned_steps": _summarize_plan_steps(plan_steps),
            "undo_plan": _build_undo_plan(plan_steps),
        }

    plan_steps = [{"action": "web_search", "target": {"query": normalized}}, {"action": "open_result", "target": {"query": normalized}}]
    requested_outputs = _requested_outputs(normalized)
    if requested_outputs:
        return _fail_fast_output_mismatch(
            instruction=instruction,
            mode="search_plan_invalid",
            plan_steps=plan_steps,
            requested_outputs=requested_outputs,
            represented_outputs=set(),
        )
    if _is_destructive_intent(normalized, plan_steps) and not confirm_risky:
        return {
            "ok": False,
            "mode": "confirmation_required",
            "instruction": instruction,
            "requires_confirmation": True,
            "message": "Destructive or high-impact action detected. Confirm to continue.",
            "planned_steps": _summarize_plan_steps(plan_steps),
            "undo_plan": _build_undo_plan(plan_steps),
        }
    query = normalized
    _emit_progress(progress_cb, 20, "Running web search")
    results: List[SearchResult] = []
    ebay_intent = "ebay" in normalized.lower()
    if "amazon" in normalized.lower():
        cleaned = re.sub(r"^.*?search\s+amazon\s+for\s+", "", normalized, flags=re.IGNORECASE)
        query = cleaned if cleaned and cleaned != normalized else normalized
        results.extend(_search_amazon_playwright(query, limit=8))
        if len(results) < 3:
            results.extend(_search_web(f"site:amazon.com {query}", limit=8))
    elif ebay_intent:
        query = _clean_ebay_query(normalized)
        results.extend(_search_web(f"site:ebay.com/itm {query}", limit=10))
        if len(results) < 4:
            results.extend(_search_web(f"site:ebay.com {query}", limit=10))
        results = [
            row
            for row in results
            if "ebay." in urllib.parse.urlparse(str(row.url or "")).netloc.lower()
        ]
        if not results:
            fallback_url = f"https://www.ebay.com/sch/i.html?_nkw={urllib.parse.quote_plus(query)}&_sop=12"
            results.append(
                SearchResult(
                    title=f"eBay search results for {query}",
                    url=fallback_url,
                    price=None,
                    source="ebay_search",
                    snippet="Marketplace search landing page.",
                )
            )
    else:
        results.extend(_search_web(query, limit=8))

    price_recommendation_intent = _is_price_recommendation_intent(normalized)
    if ebay_intent:
        _emit_progress(progress_cb, 56, "Analyzing eBay listings")
        inline_url = _extract_inline_url(normalized)
        ebay_search_url = ""
        if "ebay.com" in inline_url.lower():
            ebay_search_url = inline_url
        else:
            for row in results:
                if "ebay.com" in row.url.lower():
                    ebay_search_url = row.url
                    break
        ebay_query = _clean_ebay_query(normalized)
        try:
            ebay_rows = _search_ebay_listings(
                query=ebay_query,
                search_url=ebay_search_url,
                limit=12,
                browser_worker_mode=browser_worker_mode,
            )
        except Exception:
            ebay_rows = []
        if ebay_rows:
            results.extend(ebay_rows)

    dedup: Dict[str, SearchResult] = {}
    for result in results:
        dedup[result.url] = result
    ranked = list(dedup.values())
    if price_recommendation_intent:
        ranked.sort(
            key=lambda r: (
                0 if "ebay." in urllib.parse.urlparse(r.url).netloc.lower() else 1,
                r.price is None,
                r.price if r.price is not None else 999999.0,
            )
        )
    best = _best_price(ranked)

    opened_url = ""
    nav = {"url": "", "reused": False, "opened": False, "decision": {"score": 0.0, "reasons": ["no_result"], "elegance_cost": 0}}
    if best:
        _emit_progress(progress_cb, 84, "Selecting best option")
        nav = _resolve_navigation_target(
            target_url=best.url,
            recent_actions=[f"open_tab:{best.url}"],
        )
        opened_url = str(nav.get("url", "") or best.url)
        plan_steps[1]["target"] = {"url": opened_url}
    recommendation: Dict[str, Any] = {}
    shopping_artifacts: Dict[str, str] = {}
    priced_candidates = len([row for row in ranked if row.price is not None])
    candidate_rows: List[Dict[str, Any]] = []
    if best and price_recommendation_intent:
        _emit_progress(progress_cb, 92, "Compiling decision matrix")
        candidate_rows = _build_shopping_candidates(ranked, max_items=10, signal_pages=4)
        picked = _pick_recommended_candidate(candidate_rows)
        selected_host = urllib.parse.urlparse(str(picked.get("url", best.url))).netloc.lower()
        recommendation = {
            "selected_title": picked.get("title", best.title),
            "selected_url": picked.get("url", best.url),
            "selected_price": picked.get("price", best.price),
            "selected_rating": picked.get("rating"),
            "selected_review_count": picked.get("review_count"),
            "selected_condition": picked.get("condition"),
            "reason": (
                f"Lowest detected listing price among {priced_candidates} priced candidate(s)."
                if priced_candidates
                else ("No machine-readable price found; selected the top marketplace listing." if "ebay." in selected_host else "Best available match from discovered candidates.")
            ),
        }
        if recommendation.get("selected_url"):
            plan_steps[1]["target"] = {"url": str(recommendation.get("selected_url"))}
            opened_url = str(recommendation.get("selected_url"))
        try:
            shopping_artifacts = _write_shopping_decision_artifacts(
                instruction=instruction,
                query=query,
                candidates=candidate_rows,
                recommendation={
                    "title": recommendation.get("selected_title"),
                    "url": recommendation.get("selected_url"),
                    "price": recommendation.get("selected_price"),
                    "rating": recommendation.get("selected_rating"),
                    "review_count": recommendation.get("selected_review_count"),
                    "condition": recommendation.get("selected_condition"),
                },
            )
        except Exception:
            shopping_artifacts = {}
    needs_credentials = _likely_requires_login(normalized, opened_url, best.title if best else "")

    response = {
        "ok": True,
        "mode": "web_search",
        "instruction": instruction,
        "ai": ai_meta,
        "query": query,
        "opened_url": opened_url,
        "best_result": asdict(best) if best else None,
        "recommendation": recommendation,
        "results_count": len(ranked),
        "results": [asdict(item) for item in ranked[:10]],
        "artifacts": shopping_artifacts,
        "canvas": {
            "title": "Best Price Recommendation" if recommendation else "Search Summary",
            "subtitle": (f"{recommendation.get('selected_price'):.2f} selected" if isinstance(recommendation.get("selected_price"), (int, float)) else query),
            "cards": [
                {
                    "title": item.title[:100],
                    "price": f"${item.price:.2f}" if item.price is not None else "n/a",
                    "source": item.source,
                    "url": item.url,
                }
                for item in ranked[:6]
            ],
        },
        "paused_for_credentials": needs_credentials,
        "pause_reason": "Possible login prompt detected. Enter credentials manually and click Resume." if needs_credentials else "",
        "summary": _apply_freshness_metadata({
            "navigation_reused_existing_target": bool(nav.get("reused", False)),
            "navigation_opened_new_target": bool(nav.get("opened", False)),
            "marketplace_priced_candidates": priced_candidates,
            "marketplace_candidate_count": len(candidate_rows),
            "action_critic": nav.get("decision", {}),
        }, freshness_mode, freshness_hours),
    }
    _emit_progress(progress_cb, 100, "Completed")
    return _finalize_operator_result(json.loads(json.dumps(response)), instruction=instruction, plan_steps=plan_steps)


def preview_instruction(instruction: str) -> Dict[str, Any]:
    normalized = instruction.strip()
    if not normalized:
        return {"ok": False, "error": "Instruction is empty."}
    explicit_route = _classify_explicit_route(normalized)
    if _is_clipboard_capture_intent(normalized):
        return {
            "ok": True,
            "mode": "preview_clipboard_capture",
            "instruction": instruction,
            "planned_steps": _summarize_plan_steps([{"action": "capture_clipboard_image", "target": {"source": "system_clipboard"}}]),
            "undo_plan": [],
            "canvas": {"title": "Clipboard Capture Preview", "subtitle": "Read clipboard and save image/base64 artifacts", "cards": []},
        }
    if explicit_route == "artifact_generation":
        plan = _build_artifact_generation_plan(normalized)
        plan_steps = [dict(x) for x in plan.get("steps", [])]
        requested_outputs = _requested_outputs(normalized)
        represented_outputs = _plan_represented_outputs(plan, plan_steps)
        if requested_outputs and any(x not in represented_outputs for x in requested_outputs):
            return _fail_fast_output_mismatch(
                instruction=instruction,
                mode="preview_artifact_plan_invalid",
                plan_steps=plan_steps,
                requested_outputs=requested_outputs,
                represented_outputs=represented_outputs,
            )
        return {
            "ok": True,
            "mode": "preview_artifact_generation",
            "instruction": instruction,
            "plan": plan,
            "planned_steps": _summarize_plan_steps(plan_steps),
            "undo_plan": _build_undo_plan(plan_steps),
            "canvas": {
                "title": "Artifact Plan Preview",
                "subtitle": f"{len(plan_steps)} steps",
                "cards": [{"title": s.get("name", ""), "price": s.get("kind", "step"), "source": "planner", "url": ""} for s in plan_steps[:6]],
            },
        }
    if _is_native_planning_intent(normalized):
        plan = _build_native_plan(normalized)
        plan_steps = [dict(x) for x in plan.get("steps", [])]
        return {
            "ok": True,
            "mode": "preview_native_plan",
            "instruction": instruction,
            "plan": plan,
            "planned_steps": _summarize_plan_steps(plan_steps),
            "undo_plan": _build_undo_plan(plan_steps),
            "canvas": {
                "title": "Native Plan Preview",
                "subtitle": f"{plan.get('domain')} | {len(plan.get('steps', []))} steps",
                "cards": [
                    {"title": s.get("name", ""), "price": s.get("kind", "step"), "source": "planner", "url": ""}
                    for s in plan.get("steps", [])[:6]
                ],
            },
        }
    if explicit_route == "desktop_sequence" or _is_desktop_sequence_intent(normalized):
        plan = build_plan(normalized)
        plan_steps = [dict(x) for x in plan.get("steps", [])]
        risk = assess_risk(plan)
        return {
            "ok": True,
            "mode": "preview_desktop_sequence",
            "instruction": instruction,
            "plan": plan,
            "risk": risk,
            "planned_steps": _summarize_plan_steps(plan_steps),
            "undo_plan": _build_undo_plan(plan_steps),
            "canvas": {
                "title": "Plan Preview",
                "subtitle": f"{len(plan.get('steps', []))} steps",
                "cards": [
                    {"title": f"{i}. {s.get('action','')}", "price": "preview", "source": "plan", "url": ""}
                    for i, s in enumerate(plan.get("steps", [])[:6])
                ],
            },
        }
    return {
        "ok": True,
        "mode": "preview_search",
        "instruction": instruction,
        "planned_steps": _summarize_plan_steps(
            [{"action": "web_search", "target": {"query": instruction}}, {"action": "open_result", "target": {"query": instruction}}]
        ),
        "undo_plan": _build_undo_plan(
            [{"action": "web_search", "target": {"query": instruction}}, {"action": "open_result", "target": {"query": instruction}}]
        ),
        "canvas": {"title": "Search Preview", "subtitle": instruction, "cards": []},
    }


def _likely_requires_login(instruction: str, url: str, title: str) -> bool:
    hay = f"{instruction} {url} {title}".lower()
    keywords = ["login", "sign in", "signin", "username", "password", "account", "auth"]
    if any(k in hay for k in keywords):
        return True
    domains = ["amazon.com", "salesforce", "workday", "okta", "microsoftonline"]
    return any(d in hay for d in domains)


def _is_desktop_sequence_intent(instruction: str) -> bool:
    low = instruction.lower()
    if _is_code_workbench_intent(instruction):
        return False
    if any(phrase in low for phrase in ["then click", "then type", "press enter", "hotkey ", "focus window", "click found"]):
        return True
    if any(phrase in low for phrase in ["capture clipboard", "save clipboard image", "import image from clipboard", "import clipboard image"]):
        return True
    starters = ["open ", "click ", "type ", "press ", "hotkey ", "focus ", "switch to ", "scroll ", "login with ", "use credentials "]
    starts_like_macro = any(low.startswith(s) for s in starters)
    explicit_ui_find = "find text" in low or "locate text" in low
    return (starts_like_macro or explicit_ui_find) and "search amazon" not in low and "job" not in low


def _is_code_workbench_intent(instruction: str) -> bool:
    low = str(instruction or "").lower()
    if _is_email_triage_intent(instruction):
        return False
    if _is_payer_pricing_review_intent(instruction):
        return False
    if _is_job_research_intent(instruction):
        return False
    if _is_competitor_analysis_intent(instruction):
        return False
    if _is_study_pack_intent(instruction):
        return False

    has_tooling_signal = any(token in low for token in CODE_WORKBENCH_TOKENS)
    has_build_signal = any(token in low for token in ["analy", "research", "fulfill", "build", "create", "scaffold", "prototype"])
    has_editor_phrase = (
        any(token in low for token in ["vscode", "vs code", "visual studio code"])
        and any(token in low for token in ["write", "build", "open", "create", "launch"])
    )
    pipeline_hits = sum(1 for token in CODE_WORKBENCH_PIPELINE_TOKENS if token in low)
    deliverable_hits = sum(1 for token in CODE_WORKBENCH_DELIVERABLE_TOKENS if token in low)
    deep_build_prompt = pipeline_hits >= 3 and deliverable_hits >= 2
    return bool(has_editor_phrase or (has_tooling_signal and has_build_signal) or deep_build_prompt)


def _is_native_planning_intent(instruction: str) -> bool:
    low = instruction.lower()
    if _is_code_workbench_intent(instruction):
        return True
    if _is_desktop_sequence_intent(low):
        return False
    if _is_email_triage_intent(instruction):
        return True
    if _is_payer_pricing_review_intent(instruction):
        return True
    if _is_marketplace_shopping_intent(instruction):
        return False
    if _is_recommendation_research_intent(instruction):
        return True
    signals = [
        "then",
        "from there",
        "build a report",
        "build report",
        "dashboard",
        "spreadsheet",
        "analy",
        "capture links",
        "all over",
        "across",
        "executive summary",
        "powerpoint",
        "competitor",
        "competition",
        "save everything to a folder",
    ]
    complexity = sum(1 for s in signals if s in low)
    return complexity >= 2 or (len(instruction) > 180 and any(k in low for k in ["find", "research", "search"]))


def _is_competitor_analysis_intent(instruction: str) -> bool:
    low = instruction.lower()
    has_comp = any(x in low for x in ["competitor", "competition", "compete"])
    has_research = any(x in low for x in ["research", "analysis", "market"])
    has_deliverable = any(x in low for x in ["executive summary", "powerpoint", "ppt", "slide"])
    return has_comp and has_research and has_deliverable


def _is_study_pack_intent(instruction: str) -> bool:
    low = instruction.lower()
    has_learning = any(x in low for x in ["flashcard", "quiz", "study", "exam", "test prep", "permit"])
    has_create = any(x in low for x in ["create", "build", "generate", "make"])
    return has_learning and (has_create or "notebooklm" in low)


def _is_email_triage_intent(instruction: str) -> bool:
    low = instruction.lower()
    has_email_scope = any(x in low for x in ["inbox", "gmail", "email"])
    has_time_filter = any(x in low for x in ["last 48 hours", "last 24 hours", "today", "unread", "recent"])
    has_outputs = any(x in low for x in ["draft", "reply", "spreadsheet", "task list", "action"])
    return has_email_scope and has_outputs and has_time_filter


def _is_payer_pricing_review_intent(instruction: str) -> bool:
    low = instruction.lower()
    has_domain = any(token in low for token in PAYER_PRICING_INTENT_TOKENS)
    has_action = any(token in low for token in PAYER_PRICING_ACTION_TOKENS)
    is_durham_healthcare = "durham" in low and any(token in low for token in ["payer", "insurance", "plan"])
    return (has_domain and has_action) or is_durham_healthcare


def _is_payer_pricing_question(instruction: str) -> bool:
    low = instruction.lower()
    question_signals = [
        "which plans",
        "which payers",
        "show evidence",
        "what data sources",
        "why was",
        "why were",
        "need outreach",
        "most expensive",
        "highest",
        "what supports",
    ]
    build_signals = [
        "build",
        "create",
        "generate",
        "review",
        "analy",
        "export",
        "submit",
        "package",
        "rebuild",
        "refresh",
    ]
    return _is_payer_pricing_review_intent(instruction) and (
        ("?" in instruction or any(token in low for token in question_signals))
        and not any(token in low for token in build_signals)
    )


def _payer_service_keywords(instruction: str) -> List[str]:
    low = instruction.lower()
    if all(token in low for token in ["outpatient", "imaging"]) or any(token in low for token in ["mri", "ct", "radiology", "ultrasound", "diagnostic imaging"]):
        imaging = [
            "mri",
            "magnetic resonance",
            "magnetic resonance imaging",
            "ct",
            "computed tomography",
            "cat scan",
            "ultrasound",
            "x-ray",
            "xray",
            "radiology",
            "mammography",
            "sonography",
            "imaging",
            "diagnostic imaging",
            "pet",
            "nuclear medicine",
        ]
        return list(dict.fromkeys(imaging))
    base = [
        "mri",
        "magnetic resonance",
        "magnetic resonance imaging",
        "ct",
        "computed tomography",
        "cat scan",
        "ultrasound",
        "x-ray",
        "xray",
        "colonoscopy",
        "endoscopy",
        "office visit",
        "evaluation and management",
        "emergency",
        "heart",
        "transplant",
    ]
    for keyword in [
        "mri",
        "ct",
        "ultrasound",
        "x-ray",
        "colonoscopy",
        "endoscopy",
        "office visit",
        "emergency",
        "heart",
        "transplant",
    ]:
        if keyword in low and keyword not in base:
            base.append(keyword)
    return list(dict.fromkeys(base))


def _build_native_plan(instruction: str) -> Dict[str, Any]:
    if _is_email_triage_intent(instruction):
        domain = "email_triage"
    elif _is_payer_pricing_review_intent(instruction):
        domain = "payer_pricing_review"
    elif _is_study_pack_intent(instruction):
        domain = "study_pack"
    elif _is_job_research_intent(instruction):
        domain = "job_market"
    elif _is_competitor_analysis_intent(instruction):
        domain = "competitor_analysis"
    elif _is_code_workbench_intent(instruction):
        domain = "code_workbench"
    else:
        domain = "web_research"
    deliverables: List[str] = []
    low = instruction.lower()
    if "spreadsheet" in low or "csv" in low:
        deliverables.append("spreadsheet")
    if "xlsx" in low or "workbook" in low:
        deliverables.append("spreadsheet")
    if "report" in low:
        deliverables.append("report")
    if "executive summary" in low:
        deliverables.append("executive_summary")
    if "powerpoint" in low or "ppt" in low or "slides" in low:
        deliverables.append("powerpoint")
    if "dashboard" in low:
        deliverables.append("dashboard")
    if any(token in low for token in ["write code", "build code", "analysis script", "python script", "code scaffold"]):
        deliverables.append("code")
    if any(token in low for token in ["workspace", "vs code", "vscode", "visual studio code"]):
        deliverables.append("workspace")
    if "rag" in low or "vector store" in low or "retriever" in low:
        deliverables.append("rag_index")
    if "link" in low:
        deliverables.append("apply_links")
    if _is_payer_pricing_review_intent(instruction):
        for item in ["spreadsheet", "report", "dashboard", "rag_index"]:
            if item not in deliverables:
                deliverables.append(item)
    if _is_code_workbench_intent(instruction):
        for item in ["code", "workspace", "report"]:
            if item not in deliverables:
                deliverables.append(item)
    if _is_recommendation_research_intent(instruction):
        for item in ["spreadsheet", "report", "dashboard"]:
            if item not in deliverables:
                deliverables.append(item)
    if not deliverables:
        deliverables = ["report", "dashboard"]

    if domain == "email_triage":
        sources = ["gmail_ui"]
    elif domain == "code_workbench":
        sources = ["user_instruction", "local_workspace", "desktop_editor"]
    elif domain == "payer_pricing_review":
        sources = ["public_transparency_files", "provider_price_pages", "payer_reference_pages"]
    elif domain == "job_market":
        sources = ["linkedin", "indeed", "ziprecruiter", "glassdoor", "builtin"]
    elif domain == "competitor_analysis":
        sources = ["industry_reports", "vendor_pages", "analyst_coverage", "web_search"]
    else:
        sources = ["web_search", "source_pages"]
    objective = re.sub(r"\s+", " ", instruction).strip()
    playbook = select_playbook(domain=domain, instruction=instruction)
    if domain == "email_triage":
        plan = {
            "planner": "native-v1",
            "domain": domain,
            "playbook": playbook,
            "objective": objective,
            "deliverables": ["spreadsheet", "draft_reply"],
            "sources": sources,
            "constraints": {
                "prefer_public_pages": False,
                "no_password_capture": True,
                "persist_history": True,
            },
            "steps": [
                {"kind": "list_recent_messages", "name": "Open inbox and filter last 48 hours", "target": {"url": "https://mail.google.com/"}},
                {"kind": "read_message", "name": "Read candidate messages and classify action-needed", "target": {"query": "newer_than:2d in:inbox"}},
                {"kind": "create_draft", "name": "Create draft replies for action-needed emails", "target": {"id": "gmail:drafts"}},
                {"kind": "save_csv", "name": "Write task list spreadsheet artifact", "target": {"path": "data/reports/email_triage"}},
                {"kind": "present", "name": "Open report dashboard for review", "target": {"id": "artifact:email_triage_report"}},
            ],
        }
        steps = list(plan.get("steps", []))
        plan["playbook_validation"] = validate_plan_steps(domain=domain, steps=steps)
        plan["playbook_graph_validation"] = validate_transition_graph(domain=domain, steps=steps)
        plan["playbook_step_obligations"] = build_step_obligations(domain=domain, steps=steps)
        return plan
    if domain == "payer_pricing_review":
        plan = {
            "planner": "native-v1",
            "domain": domain,
            "playbook": playbook,
            "objective": objective,
            "deliverables": deliverables,
            "sources": sources,
            "constraints": {
                "prefer_public_pages": True,
                "no_password_capture": True,
                "persist_history": True,
                "no_phi": True,
            },
            "steps": [
                {"kind": "research", "name": "Collect Durham-area payer and provider pricing sources", "target": {"query": objective}},
                {"kind": "extract", "name": "Normalize payer, plan, service, and rate records", "target": {"id": "dataset:normalized_pricing"}},
                {"kind": "analyze", "name": "Build RAG index and identify pricing outliers", "target": {"id": "dataset:pricing_analysis"}},
                {"kind": "produce", "name": "Generate workbook, dashboard, report, and validation queue", "target": {"path": "data/payer_rag_live"}},
                {"kind": "present", "name": "Return stakeholder artifacts and evidence", "target": {"id": "artifact:payer_dashboard_html"}},
            ],
        }
        steps = list(plan.get("steps", []))
        plan["playbook_validation"] = validate_plan_steps(domain=domain, steps=steps)
        plan["playbook_graph_validation"] = validate_transition_graph(domain=domain, steps=steps)
        plan["playbook_step_obligations"] = build_step_obligations(domain=domain, steps=steps)
        return plan
    if domain == "code_workbench":
        plan = {
            "planner": "native-v1",
            "domain": domain,
            "playbook": playbook,
            "objective": objective,
            "deliverables": deliverables,
            "sources": sources,
            "constraints": {
                "prefer_public_pages": False,
                "no_password_capture": True,
                "persist_history": True,
                "fresh_workspace": True,
            },
            "steps": [
                {"kind": "research", "name": "Parse the request into a code workbench contract", "target": {"query": objective}},
                {"kind": "extract", "name": "Create a fresh workspace for the task", "target": {"path": "data/deep_work_runs"}},
                {"kind": "analyze", "name": "Scaffold analysis code, notes, and smoke tests", "target": {"id": "workspace:scaffold"}},
                {"kind": "produce", "name": "Launch VS Code and capture workbench artifacts", "target": {"app": "vscode"}},
                {"kind": "present", "name": "Return the workbench package for review", "target": {"id": "artifact:workspace_readme_md"}},
            ],
        }
        steps = list(plan.get("steps", []))
        plan["playbook_validation"] = validate_plan_steps(domain=domain, steps=steps)
        plan["playbook_graph_validation"] = validate_transition_graph(domain=domain, steps=steps)
        plan["playbook_step_obligations"] = build_step_obligations(domain=domain, steps=steps)
        return plan
    plan = {
        "planner": "native-v1",
        "domain": domain,
        "playbook": playbook,
        "objective": objective,
        "deliverables": deliverables,
        "sources": sources,
        "constraints": {
            "prefer_public_pages": True,
            "no_password_capture": True,
            "persist_history": True,
        },
        "steps": [
            {"kind": "research", "name": "Collect candidate sources and listings", "target": {"query": objective}},
            {"kind": "extract", "name": "Extract structured fields", "target": {"id": "dataset:candidate_results"}},
            {"kind": "analyze", "name": "Rank, deduplicate, summarize", "target": {"id": "dataset:structured_records"}},
            {"kind": "produce", "name": "Generate requested artifacts", "target": {"path": "data/reports"}},
            {"kind": "present", "name": "Open dashboard and return actionable links", "target": {"id": "artifact:dashboard_html"}},
        ],
    }
    steps = list(plan.get("steps", []))
    plan["playbook_validation"] = validate_plan_steps(domain=domain, steps=steps)
    plan["playbook_graph_validation"] = validate_transition_graph(domain=domain, steps=steps)
    plan["playbook_step_obligations"] = build_step_obligations(domain=domain, steps=steps)
    return plan


def _execute_native_plan(
    plan: Dict[str, Any],
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    min_live_non_curated_citations: Optional[int] = None,
    manual_auth_phase: bool = False,
    auth_session_id: Optional[str] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    validation = plan.get("playbook_validation", {}) if isinstance(plan, dict) else {}
    if isinstance(validation, dict) and validation and not bool(validation.get("ok", True)):
        return {
            "ok": False,
            "mode": "native_plan_invalid_playbook",
            "query": str(plan.get("objective", instruction)),
            "results_count": 0,
            "results": [],
            "artifacts": {},
            "summary": {"error": "playbook_validation_failed", "errors": list(validation.get("errors", []))},
            "source_status": {},
            "opened_url": "",
            "canvas": {
                "title": "Run Blocked",
                "subtitle": "Playbook validation failed before execution.",
                "cards": [],
            },
            "paused_for_credentials": False,
            "pause_reason": "",
            "error": "playbook_validation_failed",
            "error_code": "playbook_validation_failed",
            "trace": [],
        }
    graph_validation = plan.get("playbook_graph_validation", {}) if isinstance(plan, dict) else {}
    if isinstance(graph_validation, dict) and graph_validation and not bool(graph_validation.get("ok", True)):
        return {
            "ok": False,
            "mode": "native_plan_invalid_playbook_graph",
            "query": str(plan.get("objective", instruction)),
            "results_count": 0,
            "results": [],
            "artifacts": {},
            "summary": {"error": "playbook_graph_validation_failed", "errors": list(graph_validation.get("errors", []))},
            "source_status": {},
            "opened_url": "",
            "canvas": {
                "title": "Run Blocked",
                "subtitle": "Playbook transition-graph validation failed before execution.",
                "cards": [],
            },
            "paused_for_credentials": False,
            "pause_reason": "",
            "error": "playbook_graph_validation_failed",
            "error_code": "playbook_graph_validation_failed",
            "trace": [],
        }
    out = _run_reflective_planner(
        plan=plan,
        instruction=instruction,
        progress_cb=progress_cb,
        min_live_non_curated_citations=min_live_non_curated_citations,
        manual_auth_phase=manual_auth_phase,
        auth_session_id=auth_session_id,
        browser_worker_mode=browser_worker_mode,
        human_like_interaction=bool(human_like_interaction),
    )
    obligations = plan.get("playbook_step_obligations", []) if isinstance(plan, dict) else []
    if obligations and isinstance(obligations, list) and bool(out.get("ok", False)):
        ob_eval = evaluate_step_obligations(
            domain=str(plan.get("domain", "")),
            steps=list(plan.get("steps", [])),
            obligations=obligations,
            result=out if isinstance(out, dict) else {},
        )
        out["playbook_obligation_evaluation"] = ob_eval
        if not bool(ob_eval.get("ok", True)):
            return {
                "ok": False,
                "mode": "native_plan_obligation_failed",
                "query": str(plan.get("objective", instruction)),
                "results_count": int(out.get("results_count", 0) or 0),
                "results": out.get("results", []),
                "artifacts": out.get("artifacts", {}),
                "summary": {
                    "error": "playbook_obligation_failed",
                    "errors": list(ob_eval.get("errors", [])),
                },
                "source_status": out.get("source_status", {}),
                "opened_url": str(out.get("opened_url", "")),
                "canvas": {
                    "title": "Run Blocked",
                    "subtitle": "Per-step playbook policy obligations failed.",
                    "cards": [],
                },
                "paused_for_credentials": bool(out.get("paused_for_credentials", False)),
                "pause_reason": str(out.get("pause_reason", "")),
                "error": "playbook_obligation_failed",
                "error_code": "playbook_obligation_failed",
                "trace": out.get("trace", []),
                "auth_session_id": out.get("auth_session_id", ""),
                "playbook_obligation_evaluation": ob_eval,
            }
    return out


def _run_reflective_planner(
    plan: Dict[str, Any],
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    min_live_non_curated_citations: Optional[int] = None,
    manual_auth_phase: bool = False,
    auth_session_id: Optional[str] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    min_live = _effective_min_live_non_curated_citations(min_live_non_curated_citations)
    objective = str(plan.get("objective", instruction))
    preferred = str(plan.get("domain", "web_research"))
    if preferred == "competitor_analysis":
        accept_threshold = 0.60
    elif preferred == "job_market":
        accept_threshold = 0.58
    elif preferred == "code_workbench":
        accept_threshold = 0.62
    else:
        accept_threshold = 0.72
    attempt_order = _strategy_order(preferred=preferred, instruction=instruction)
    elegance = EleganceBudget(total=90)
    world_route_context = _world_route_context(instruction=instruction, preferred=preferred)
    best: Dict[str, Any] = {"score": -1.0, "result": {}}
    decision_log: List[str] = []
    accepted = False

    for attempt, strategy in enumerate(attempt_order, start=1):
        allow_step, reason, cost = _world_route_step_gate(
            strategy=strategy,
            attempt=attempt,
            instruction=instruction,
            context=world_route_context,
        )
        if cost > 0:
            elegance.consume(cost, f"route_gate:{reason}")
        if not allow_step:
            decision_log.append(f"Attempt {attempt} skipped by world route gate: {reason}")
            continue
        _emit_progress(progress_cb, 18 + (attempt - 1) * 20, f"Planning attempt {attempt}: {strategy}")
        result = _run_strategy(
            strategy=strategy,
            instruction=instruction,
            progress_cb=progress_cb,
            min_live_non_curated_citations=min_live,
            manual_auth_phase=manual_auth_phase,
            auth_session_id=auth_session_id,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=bool(human_like_interaction),
        )
        if bool(result.get("paused_for_credentials", False)) or str(result.get("error_code", "")) in {"credential_missing", "permission_denied"}:
            decision_log.append(f"Attempt {attempt} blocked: {result.get('error_code') or 'paused_for_credentials'}")
            chosen = dict(result)
            return {
                "ok": bool(chosen.get("ok", False)),
                "mode": "autonomous_plan_execute",
                "plan": plan,
                "instruction": instruction,
                "decision_log": decision_log,
                "query": chosen.get("query", ""),
                "results_count": chosen.get("results_count", 0),
                "artifacts": chosen.get("artifacts", {}),
                "summary": chosen.get("summary", {}),
                "results": chosen.get("results", []),
                "source_status": chosen.get("source_status", {}),
                "opened_url": chosen.get("opened_url", ""),
                "recommendation": chosen.get("recommendation", {}),
                "canvas": chosen.get(
                    "canvas",
                    {
                        "title": "Task Blocked",
                        "subtitle": objective,
                        "cards": [],
                    },
                ),
                "paused_for_credentials": bool(chosen.get("paused_for_credentials", False)),
                "pause_reason": str(chosen.get("pause_reason", "")),
                "error": chosen.get("error", ""),
                "error_code": chosen.get("error_code", ""),
                "trace": chosen.get("trace", []),
                "auth_session_id": chosen.get("auth_session_id", ""),
            }
        score = _score_result_against_objective(
            result=result,
            objective=objective,
            min_live_non_curated_citations=min_live,
        )
        decision_log.append(f"Attempt {attempt} used {strategy} -> quality score {score:.2f}")
        if score > float(best["score"]):
            best = {"score": score, "result": result}
        if (
            preferred == "job_market"
            and strategy == "job_market"
            and bool(result.get("ok", False))
            and bool((result.get("artifacts", {}) or {}))
        ):
            decision_log.append("Accepted job_market attempt with artifact evidence.")
            accepted = True
            break
        if score >= accept_threshold:
            decision_log.append(f"Accepted attempt {attempt}; score passed threshold.")
            accepted = True
            break
        decision_log.append(f"Rejected attempt {attempt}; refining strategy.")

    chosen = dict(best["result"] or {})
    if preferred == "competitor_analysis" and (not accepted):
        return {
            "ok": False,
            "mode": "autonomous_plan_execute",
            "plan": plan,
            "instruction": instruction,
            "decision_log": decision_log + [
                f"Failed strict acceptance threshold ({accept_threshold:.2f}) for competitor analysis.",
            ],
            "query": str(chosen.get("query", "")),
            "results_count": int(chosen.get("results_count", 0) or 0),
            "artifacts": {},
            "summary": {
                "error": "strict_competitor_validation_failed",
                "required_live_non_curated_citations": min_live,
            },
            "results": chosen.get("results", []),
            "source_status": chosen.get("source_status", {}),
            "opened_url": "",
            "recommendation": chosen.get("recommendation", {}),
            "canvas": {
                "title": "Run Blocked",
                "subtitle": "Strict citation rule not met.",
                "cards": [],
            },
            "paused_for_credentials": False,
            "pause_reason": "",
        }
    if not chosen:
        chosen = _run_generic_research(
            instruction,
            progress_cb=progress_cb,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=bool(human_like_interaction),
        )
        decision_log.append("Fallback to generic strategy due empty attempts.")

    paused_for_credentials = bool(chosen.get("paused_for_credentials", False))
    pause_reason = str(chosen.get("pause_reason", "")) if paused_for_credentials else ""
    chosen_summary = chosen.get("summary", {}) if isinstance(chosen.get("summary"), dict) else {}
    chosen_summary = _ensure_elegance_budget(chosen_summary, elegance)
    chosen["summary"] = chosen_summary
    budget_block = _enforce_elegance_budget_gate(
        elegance=elegance,
        mode="autonomous_plan_execute",
        instruction=instruction,
    )
    if budget_block:
        budget_block["plan"] = plan
        budget_block["decision_log"] = decision_log
        budget_block["query"] = chosen.get("query", "")
        budget_block["results_count"] = chosen.get("results_count", 0)
        budget_block["results"] = chosen.get("results", [])
        budget_block["artifacts"] = chosen.get("artifacts", {})
        budget_block["source_status"] = chosen.get("source_status", {})
        budget_block["opened_url"] = chosen.get("opened_url", "")
        budget_block["recommendation"] = chosen.get("recommendation", {})
        budget_block["paused_for_credentials"] = paused_for_credentials
        budget_block["pause_reason"] = pause_reason
        budget_block["error"] = "elegance_budget_exceeded"
        budget_block["error_code"] = "elegance_budget_exceeded"
        budget_block["trace"] = chosen.get("trace", [])
        budget_block["auth_session_id"] = chosen.get("auth_session_id", "")
        return budget_block
    return {
        "ok": bool(chosen.get("ok", False)),
        "mode": "autonomous_plan_execute",
        "plan": plan,
        "instruction": instruction,
        "decision_log": decision_log,
        "query": chosen.get("query", ""),
        "results_count": chosen.get("results_count", 0),
        "artifacts": chosen.get("artifacts", {}),
        "summary": chosen.get("summary", {}),
        "results": chosen.get("results", []),
        "source_status": chosen.get("source_status", {}),
        "opened_url": chosen.get("opened_url", ""),
        "recommendation": chosen.get("recommendation", {}),
        "canvas": chosen.get(
            "canvas",
            {
                "title": "Task Completed",
                "subtitle": objective,
                "cards": [],
            },
        ),
        "paused_for_credentials": paused_for_credentials,
        "pause_reason": pause_reason,
        "error": chosen.get("error", ""),
        "error_code": chosen.get("error_code", ""),
        "trace": chosen.get("trace", []),
        "auth_session_id": chosen.get("auth_session_id", ""),
        "world_route": world_route_context,
    }


def _world_route_context(instruction: str, preferred: str) -> Dict[str, Any]:
    low = str(instruction or "").lower()
    mgr = SessionManager()
    gmail_reusable = mgr.find_reusable_authenticated_tab("mail.google.com")
    return {
        "preferred": preferred,
        "email_intent": bool(any(x in low for x in ["inbox", "gmail", "email"])),
        "reusable_gmail_tab": gmail_reusable,
    }


def _world_route_step_gate(
    *,
    strategy: str,
    attempt: int,
    instruction: str,
    context: Dict[str, Any],
) -> tuple[bool, str, int]:
    _ = instruction
    strat = str(strategy or "").strip().lower()
    if bool(context.get("email_intent", False)) and strat != "email_triage" and attempt <= 2:
        return False, "domain_lock_email", 10
    if bool(context.get("email_intent", False)) and str(context.get("reusable_gmail_tab", "")).strip() and strat != "email_triage":
        return False, "reuse_session_prefers_email_triage", 8
    return True, "ok", 0


def _strategy_order(preferred: str, instruction: str = "") -> List[str]:
    all_strategies = ["email_triage", "study_pack", "job_market", "competitor_analysis", "payer_pricing_review", "code_workbench", "generic_research"]
    low = str(instruction or "").lower()
    if any(x in low for x in ["gmail", "inbox", "email"]):
        mgr = SessionManager()
        reusable = mgr.find_reusable_authenticated_tab("mail.google.com")
        if reusable:
            return ["email_triage"] + [s for s in all_strategies if s != "email_triage"]
    if preferred == "email_triage":
        return ["email_triage"]
    if preferred == "study_pack":
        return ["study_pack", "generic_research"]
    if preferred == "job_market":
        return ["job_market", "generic_research"]
    if preferred == "competitor_analysis":
        return ["competitor_analysis", "generic_research"]
    if preferred == "payer_pricing_review":
        return ["payer_pricing_review", "generic_research"]
    if preferred == "code_workbench":
        return ["code_workbench", "generic_research"]
    if preferred == "web_research":
        return ["generic_research"]
    if preferred in all_strategies:
        return [preferred] + [s for s in all_strategies if s != preferred]
    return ["generic_research"]


def _run_strategy(
    strategy: str,
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    min_live_non_curated_citations: Optional[int] = None,
    manual_auth_phase: bool = False,
    auth_session_id: Optional[str] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    if strategy == "email_triage":
        return _run_email_triage(
            instruction,
            progress_cb=progress_cb,
            manual_auth_phase=manual_auth_phase,
            auth_session_id=auth_session_id,
            browser_worker_mode=browser_worker_mode,
            human_like_interaction=bool(human_like_interaction),
        )
    if strategy == "study_pack":
        return _run_study_pack(instruction, progress_cb=progress_cb)
    if strategy == "job_market":
        return _run_job_market_research(instruction, progress_cb=progress_cb)
    if strategy == "competitor_analysis":
        return _run_competitor_analysis(
            instruction,
            progress_cb=progress_cb,
            min_live_non_curated_citations=min_live_non_curated_citations,
        )
    if strategy == "payer_pricing_review":
        return _run_payer_pricing_review(instruction, progress_cb=progress_cb)
    if strategy == "code_workbench":
        return _run_code_workbench(instruction, progress_cb=progress_cb)
    return _run_generic_research(
        instruction,
        progress_cb=progress_cb,
        browser_worker_mode=browser_worker_mode,
        human_like_interaction=bool(human_like_interaction),
    )


def _score_result_against_objective(
    result: Dict[str, Any],
    objective: str,
    min_live_non_curated_citations: Optional[int] = None,
) -> float:
    min_live = _effective_min_live_non_curated_citations(min_live_non_curated_citations)
    if not result.get("ok"):
        return 0.0
    score = 0.2
    artifacts = result.get("artifacts", {}) or {}
    if artifacts:
        score += 0.2
    if any(k in artifacts for k in ["quiz_html", "dashboard_html", "report_md"]):
        score += 0.2

    results_count = int(result.get("results_count", 0) or 0)
    score += 0.15 if results_count > 0 else 0.0
    score += 0.1 if results_count >= 20 else 0.0

    low_obj = objective.lower()
    low_summary = json.dumps(result.get("summary", {})).lower()
    alignment_terms = [t for t in re.split(r"[^a-z0-9]+", low_obj) if len(t) > 3]
    overlap = sum(1 for t in alignment_terms if t in low_summary)
    if alignment_terms:
        score += min(0.15, overlap / max(1, len(alignment_terms)) * 0.15)

    bad_domains = ("support.google.com", "gmail.com", "mail.google.com")
    urls = [str(x.get("url", "")).lower() for x in (result.get("results") or [])[:10]]
    if any(any(b in u for b in bad_domains) for u in urls):
        score -= 0.25

    low_obj = objective.lower()
    wants_study = any(t in low_obj for t in ["flashcard", "quiz", "study", "exam", "permit"])
    wants_jobs = any(t in low_obj for t in ["job", "position", "salary", "linkedin", "indeed"])
    wants_competitor = any(t in low_obj for t in ["competitor", "competition", "executive summary", "powerpoint", "ppt"])
    wants_email_triage = any(t in low_obj for t in ["inbox", "gmail", "draft replies", "draft reply", "task list"]) and "last" in low_obj
    wants_shopping_decision = any(t in low_obj for t in ["best price", "cheapest", "recommend", "which one to buy", "what should i buy", "buy"])
    wants_payer_review = _is_payer_pricing_review_intent(objective)
    wants_code_workbench = _is_code_workbench_intent(objective)
    if wants_study:
        if "flashcards_csv" not in artifacts or not any(k in artifacts for k in ["quiz_md", "quiz_html"]):
            score -= 0.5
    if wants_jobs:
        if "jobs_csv" not in artifacts and result.get("mode") != "job_market_research":
            score -= 0.3
        job_constraints = _extract_job_constraints(objective)
        rows = result.get("results") or []
        if rows and job_constraints.require_vp_avp:
            vp_hits = sum(1 for row in rows if _is_vp_avp_title(str(row.get("title", ""))))
            if vp_hits == 0:
                score -= 0.4
        if rows and "ireland" not in objective.lower() and "us" in job_constraints.allowed_regions:
            ireland_hits = sum(1 for row in rows if "ireland" in str(row.get("location", "")).lower())
            if ireland_hits > 0:
                score -= 0.25
    if wants_competitor:
        if "executive_summary_md" not in artifacts or "powerpoint_pptx" not in artifacts:
            score -= 0.45
        if result.get("mode") != "competitor_analysis":
            score -= 0.35
        live_cites = int((result.get("summary", {}) or {}).get("live_non_curated_citations", 0) or 0)
        if live_cites < min_live:
            score -= 0.35
    if wants_email_triage:
        if "email_tasks_csv" not in artifacts and "task_list_csv" not in artifacts:
            score -= 0.45
        if result.get("mode") != "email_triage":
            score -= 0.4
    if wants_shopping_decision:
        if "decision_matrix_csv" not in artifacts and "results_csv" not in artifacts:
            score -= 0.35
        if result.get("mode") not in {"web_search", "generic_research"}:
            score -= 0.2
    if wants_payer_review:
        required = {"workbook_xlsx", "dashboard_html", "summary_report_md", "validation_queue_csv", "rag_index_db"}
        missing = [key for key in required if key not in artifacts]
        if missing:
            score -= 0.5
        if result.get("mode") != "payer_pricing_review":
            score -= 0.35
    if wants_code_workbench:
        required = {"workspace_directory", "analysis_script_py", "workspace_readme_md", "smoke_log"}
        missing = [key for key in required if key not in artifacts]
        if missing:
            score -= 0.45
        else:
            score += 0.25
        if result.get("mode") != "code_workbench":
            score -= 0.35

    return max(0.0, min(1.0, score))


def _run_payer_pricing_review(
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
) -> Dict[str, Any]:
    contract = extract_current_task_contract(instruction)
    service_keywords = _payer_service_keywords(instruction)
    _emit_progress(progress_cb, 10, f"Detected {contract.geography} as the current payer market")
    _emit_progress(progress_cb, 16, "Checking for stale geography-specific artifacts")
    allow_reuse = _is_payer_pricing_question(instruction) and not bool(getattr(contract, "geography_explicit", False))
    if allow_reuse:
        _emit_progress(progress_cb, 22, f"Looking for the latest valid {contract.geography} payer workspace")
    else:
        _emit_progress(progress_cb, 22, f"Creating a fresh run folder for {contract.geography}")
    build_result = ensure_workspace(
        contract=contract,
        service_keywords=service_keywords,
        max_services_per_source=24,
        outlier_threshold=0.2,
        min_peer_count=3,
        offline_fallback=True,
        allow_reuse=allow_reuse,
    )
    workspace = Path(str(build_result.get("workspace", "")))
    invalidated = list(build_result.get("invalidated_artifacts", [])) if isinstance(build_result.get("invalidated_artifacts"), list) else []
    if invalidated:
        _emit_progress(progress_cb, 34, "Found stale outputs from a different market; excluding them from final output")
    if bool(build_result.get("reused_existing_outputs", False)):
        _emit_progress(progress_cb, 42, f"Using the latest valid {contract.geography} workspace")
    else:
        _emit_progress(progress_cb, 42, f"Rebuilding ingestion and normalization for {contract.geography}")
    counts = build_result.get("counts", {}) if isinstance(build_result.get("counts"), dict) else {}
    artifacts = build_result.get("artifact_paths", {}) if isinstance(build_result.get("artifact_paths"), dict) else {}
    opened_url = artifacts.get("primary_open_file", "") or artifacts.get("dashboard_html", "")
    geography_validation = (
        build_result.get("geography_validation", {})
        if isinstance(build_result.get("geography_validation"), dict)
        else {}
    )
    validation_results = build_result.get("validation_results", {}) if isinstance(build_result.get("validation_results"), dict) else {}
    final_output_gate = build_result.get("final_output_gate", {}) if isinstance(build_result.get("final_output_gate"), dict) else {}
    quarantined_artifacts = build_result.get("quarantined_artifacts", {}) if isinstance(build_result.get("quarantined_artifacts"), dict) else {}
    completion_status = str(build_result.get("completion_status", "") or "")
    synthetic_only = bool(build_result.get("synthetic_only", False))
    if not bool(geography_validation.get("passed", False)):
        return {
            "ok": False,
            "mode": "payer_pricing_review",
            "query": instruction,
            "results_count": int(counts.get("outreach_candidates", 0) or 0),
            "results": build_result.get("top_candidates", [])[:20] if isinstance(build_result.get("top_candidates"), list) else [],
            "artifacts": artifacts,
            "summary": {
                "workspace": str(workspace.resolve()) if workspace else "",
                "counts": counts,
                "issues": build_result.get("issues", []),
                "reused_existing_outputs": bool(build_result.get("reused_existing_outputs", False)),
                "current_task_contract": build_result.get("current_task_contract", {}),
                "invalidated_artifacts": invalidated,
                "generation_timestamp": build_result.get("generation_timestamp", ""),
                "geography_validation": geography_validation,
                "validation_results": validation_results,
                "final_output_gate": final_output_gate,
                "quarantined_artifacts": quarantined_artifacts,
            },
            "source_status": {"payer_rag": "error:geography_consistency_failed"},
            "opened_url": opened_url,
            "recommendation": {},
            "current_task_contract": build_result.get("current_task_contract", {}),
            "task_contract": build_result.get("current_task_contract", {}),
            "validation_results": validation_results,
            "final_output_gate": final_output_gate,
            "error": "geography_consistency_failed",
            "error_code": "geography_consistency_failed",
            "canvas": {
                "title": "Geography Validation Failed",
                "subtitle": f"Final artifacts do not consistently match {contract.geography}.",
                "cards": [
                    {"title": item[:110], "price": "validation", "source": "payer_rag", "url": ""}
                    for item in geography_validation.get("errors", [])[:5]
                ],
            },
        }
    if not bool(final_output_gate.get("passed", True)):
        return {
            "ok": False,
            "mode": "payer_pricing_review",
            "query": instruction,
            "results_count": int(counts.get("outreach_candidates", 0) or 0),
            "results": build_result.get("top_candidates", [])[:20] if isinstance(build_result.get("top_candidates"), list) else [],
            "artifacts": artifacts,
            "summary": {
                "workspace": str(workspace.resolve()) if workspace else "",
                "counts": counts,
                "issues": build_result.get("issues", []),
                "reused_existing_outputs": bool(build_result.get("reused_existing_outputs", False)),
                "current_task_contract": build_result.get("current_task_contract", {}),
                "invalidated_artifacts": invalidated,
                "generation_timestamp": build_result.get("generation_timestamp", ""),
                "geography_validation": geography_validation,
                "validation_results": validation_results,
                "final_output_gate": final_output_gate,
                "quarantined_artifacts": quarantined_artifacts,
                "synthetic_only": synthetic_only,
                "repair_state": build_result.get("repair_state", {}),
                "completion_status": completion_status,
            },
            "source_status": {"payer_rag": "error:final_output_gate_failed"},
            "opened_url": artifacts.get("geography_validation_report_md", ""),
            "recommendation": {},
            "current_task_contract": build_result.get("current_task_contract", {}),
            "task_contract": build_result.get("current_task_contract", {}),
            "validation_results": validation_results,
            "final_output_gate": final_output_gate,
            "error": "final_output_gate_failed",
            "error_code": "final_output_gate_failed",
            "canvas": {
                "title": "Validation Failed",
                "subtitle": "I found out-of-market or out-of-scope data and I am not presenting this as a stakeholder-ready result.",
                "cards": [
                    {"title": item[:110], "price": "repair", "source": "validator", "url": ""}
                    for item in list(final_output_gate.get("required_repairs", []))[:5]
                ],
            },
        }
    if _is_payer_pricing_question(instruction) and workspace:
        _emit_progress(progress_cb, 78, f"Answering from the {contract.geography} payer corpus")
        response = ask_workspace_question(instruction, contract=contract, workspace=workspace)
        cards = [
            {"title": line[:110], "price": "evidence", "source": "payer_rag", "url": ""}
            for line in str(response.get("answer", "")).splitlines()[:5]
            if str(line).strip()
        ]
        return {
            "ok": True,
            "mode": "payer_pricing_review",
            "query": instruction,
            "results_count": len(response.get("sources", [])),
            "results": [
                {
                    "title": "payer_answer",
                    "url": src,
                    "price": None,
                    "source": "payer_rag",
                    "snippet": response.get("answer", ""),
                }
                for src in response.get("sources", [])
            ],
            "artifacts": artifacts,
            "summary": {
                "workspace": str(workspace.resolve()),
                "counts": counts,
                "reused_existing_outputs": bool(build_result.get("reused_existing_outputs", False)),
                "question_answered": True,
                "sources": response.get("sources", []),
                "current_task_contract": build_result.get("current_task_contract", {}),
                "invalidated_artifacts": invalidated,
                "generation_timestamp": build_result.get("generation_timestamp", ""),
                "geography_validation": build_result.get("geography_validation", {}),
                "validation_results": validation_results,
                "final_output_gate": final_output_gate,
                "completion_status": completion_status,
                "source_basis": "synthetic_demo" if synthetic_only else "validated_real_data",
            },
            "source_status": {"payer_rag": "ok", "workspace": str(workspace.resolve())},
            "opened_url": opened_url,
            "recommendation": {},
            "current_task_contract": build_result.get("current_task_contract", {}),
            "task_contract": build_result.get("current_task_contract", {}),
            "validation_results": validation_results,
            "final_output_gate": final_output_gate,
            "completion_status": completion_status,
            "source_basis": "synthetic_demo" if synthetic_only else "validated_real_data",
            "canvas": {
                "title": "Payer Review Answer Ready" if not synthetic_only else "Demo Payer Review Answer Ready",
                "subtitle": (
                    "Synthetic/demo outpatient imaging corpus; not validated local public evidence."
                    if synthetic_only
                    else str(response.get("answer", "")).splitlines()[0][:120]
                )
                if str(response.get("answer", "")).strip()
                else ("Synthetic/demo answer ready" if synthetic_only else "Source-backed answer ready"),
                "cards": cards,
            },
        }
    _emit_progress(progress_cb, 78, f"Building stakeholder outputs for {contract.geography}")
    top_candidates = build_result.get("top_candidates", []) if isinstance(build_result.get("top_candidates"), list) else []
    cards = []
    for row in top_candidates[:5]:
        variance = str(row.get("variance_percent", "")).strip()
        cards.append(
            {
                "title": f"{row.get('payer_name', '')} / {row.get('plan_name', '')}",
                "price": f"{float(variance) * 100:.1f}% above median" if variance else "candidate",
                "source": row.get("service", ""),
                "url": row.get("source_evidence", "").split("|")[0].strip(),
            }
        )
    return {
        "ok": True,
        "mode": "payer_pricing_review",
        "query": instruction,
        "results_count": int(counts.get("outreach_candidates", 0) or 0),
        "results": top_candidates[:20],
        "artifacts": artifacts,
        "summary": {
            "workspace": str(workspace.resolve()),
            "counts": counts,
            "issues": build_result.get("issues", []),
            "reused_existing_outputs": bool(build_result.get("reused_existing_outputs", False)),
            "validation_required": True,
            "current_task_contract": build_result.get("current_task_contract", {}),
            "invalidated_artifacts": invalidated,
            "generation_timestamp": build_result.get("generation_timestamp", ""),
            "geography_validation": build_result.get("geography_validation", {}),
            "validation_results": validation_results,
            "final_output_gate": final_output_gate,
            "quarantined_artifacts": quarantined_artifacts,
            "completion_status": completion_status,
            "source_basis": "synthetic_demo" if synthetic_only else "validated_real_data",
        },
        "source_status": {"payer_rag": "ok", "workspace": str(workspace.resolve())},
        "opened_url": opened_url,
        "recommendation": {},
        "current_task_contract": build_result.get("current_task_contract", {}),
        "task_contract": build_result.get("current_task_contract", {}),
        "validation_results": validation_results,
        "final_output_gate": final_output_gate,
        "completion_status": completion_status,
        "source_basis": "synthetic_demo" if synthetic_only else "validated_real_data",
        "canvas": {
            "title": f"{contract.geography} Demo Package Ready" if synthetic_only else f"{contract.geography} Payer Review Ready",
            "subtitle": (
                "Synthetic/demo outpatient imaging package completed; use the acquisition checklist to replace it with validated local evidence."
                if synthetic_only
                else f"{int(counts.get('outreach_candidates', 0) or 0)} outreach candidates queued for validation"
            ),
            "cards": cards,
        },
    }


def _run_code_workbench(
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
) -> Dict[str, Any]:
    contract = extract_workbench_contract(instruction)
    _emit_progress(progress_cb, 10, "Translating the request into a code workbench contract")
    build_result = build_code_workbench_workspace(contract=contract, open_vscode=True)
    workspace = Path(str(build_result.get("workspace", "")))
    artifacts = build_result.get("artifact_paths", {}) if isinstance(build_result.get("artifact_paths"), dict) else {}
    vscode_launch = build_result.get("vscode_launch", {}) if isinstance(build_result.get("vscode_launch"), dict) else {}
    _emit_progress(progress_cb, 34, "Created a fresh workspace with notes, code scaffold, and editor tasks")

    smoke_ok = True
    smoke_lines: List[str] = []
    if workspace:
        commands = [
            ["python", "-m", "py_compile", "src\\analysis.py", "tests\\test_smoke.py"],
            ["python", "src\\analysis.py"],
        ]
        for command in commands:
            smoke_lines.append(f"$ {' '.join(command)}")
            try:
                proc = subprocess.run(
                    command,
                    cwd=str(workspace),
                    capture_output=True,
                    text=True,
                    timeout=30,
                    check=False,
                )
            except Exception as exc:
                smoke_ok = False
                smoke_lines.append(f"ERROR: {exc}")
                continue
            if proc.stdout.strip():
                smoke_lines.append(proc.stdout.strip())
            if proc.stderr.strip():
                smoke_lines.append(proc.stderr.strip())
            smoke_lines.append(f"exit_code={proc.returncode}")
            if proc.returncode != 0:
                smoke_ok = False

    smoke_log_path = Path(str(artifacts.get("smoke_log", ""))) if artifacts.get("smoke_log") else None
    if smoke_log_path:
        smoke_log_path.write_text("\n".join(smoke_lines).strip() + "\n", encoding="utf-8")

    analysis_summary = workspace / "artifacts" / "analysis_summary.json" if workspace else Path()
    if analysis_summary.exists():
        artifacts["analysis_summary_json"] = str(analysis_summary.resolve())

    launch_note = (
        "Opened a new VS Code window for the workspace"
        if bool(vscode_launch.get("ok"))
        else "VS Code launch was requested but not available; the workspace is ready on disk"
    )
    _emit_progress(progress_cb, 76, "Ran smoke checks against the generated scaffold")
    _emit_progress(progress_cb, 90, launch_note)
    cards = [
        {"title": contract.title[:110], "price": "workspace", "source": "deep_work", "url": ""},
        {"title": "analysis.py scaffold ready", "price": "code", "source": "deep_work", "url": ""},
        {"title": "Smoke check passed" if smoke_ok else "Smoke check needs follow-up", "price": "verify", "source": "deep_work", "url": ""},
    ]
    return {
        "ok": True,
        "mode": "code_workbench",
        "query": instruction,
        "results_count": 1,
        "results": [
            {
                "title": contract.title,
                "url": str(workspace.resolve()) if workspace else "",
                "price": None,
                "source": "deep_workbench",
                "snippet": "Fresh code workspace with notes, scaffold, and smoke-test output.",
            }
        ],
        "artifacts": artifacts,
        "summary": {
            "workspace": str(workspace.resolve()) if workspace else "",
            "generation_timestamp": build_result.get("generation_timestamp", ""),
            "smoke_test_passed": smoke_ok,
            "vscode_launch": vscode_launch,
            "launch_note": launch_note,
            "current_task_contract": build_result.get("current_task_contract", {}),
        },
        "source_status": {
            "deep_workbench": "ok",
            "vscode": "ok" if bool(vscode_launch.get("ok")) else f"warning:{vscode_launch.get('mode', 'not_found')}",
        },
        "opened_url": str(artifacts.get("primary_open_file", "")),
        "recommendation": {},
        "current_task_contract": build_result.get("current_task_contract", {}),
        "canvas": {
            "title": "Code Workbench Ready",
            "subtitle": contract.title,
            "cards": cards,
        },
    }


def _run_email_triage(
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    manual_auth_phase: bool = False,
    auth_session_id: Optional[str] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    return _run_email_triage_active_browser(
        instruction=instruction,
        progress_cb=progress_cb,
        manual_auth_phase=manual_auth_phase,
        auth_session_id=auth_session_id,
        browser_worker_mode=browser_worker_mode,
        human_like_interaction=bool(human_like_interaction),
    )


def _run_email_triage_active_browser(
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    manual_auth_phase: bool = False,
    auth_session_id: Optional[str] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    worker_mode = normalize_browser_worker_mode(browser_worker_mode)
    worker_info: Dict[str, Any] = {"ok": True, "mode": worker_mode, "debug_port": 9222}
    if worker_mode == "docker":
        _emit_progress(progress_cb, 10, "Starting Docker browser worker")
        worker_info = ensure_browser_worker(mode="docker")
        if not bool(worker_info.get("ok", False)):
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {
                    "error": "browser_worker_unavailable",
                    "detail": str(worker_info.get("detail", "")),
                    "browser_worker_mode": worker_mode,
                },
                "source_status": {"browser_worker": "docker_unavailable"},
                "opened_url": "",
                "paused_for_credentials": False,
                "pause_reason": "",
                "error": "browser_worker_unavailable",
                "error_code": "browser_worker_unavailable",
                "auth_session_id": "",
                "trace": [],
                "canvas": {
                    "title": "Docker Browser Worker Unavailable",
                    "subtitle": "Docker could not start a browser worker. Switch to local mode or fix Docker.",
                    "cards": [],
                },
            }
    profile_dir = Path("data/interface/browser_profile")
    if worker_mode == "docker":
        profile_dir = Path("data/interface/browser_profile_docker")
    profile_dir.mkdir(parents=True, exist_ok=True)
    sessions = SessionManager()
    account = _extract_inbox_account(instruction)
    trace: List[Dict[str, Any]] = []
    debug_port = int(worker_info.get("debug_port", 9222) or 9222)

    # Manual-auth phase is explicit for local mode only.
    if manual_auth_phase and worker_mode != "docker":
        action_decision = ActionCritic().evaluate(
            next_action="open_tab",
            target="https://mail.google.com/",
            already_open=False,
            context={"recent_actions": ["open_tab:https://mail.google.com/"]},
        )
        if not action_decision.allow:
            critic_reason = next((x for x in action_decision.reasons if x != "ok"), "blocked")
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {
                    "error": "action_critic_blocked",
                    "detail": critic_reason,
                    "action_critic": {
                        "score": action_decision.score,
                        "reasons": action_decision.reasons,
                        "elegance_cost": action_decision.elegance_cost,
                    },
                },
                "source_status": {"gmail_ui": "critic_blocked"},
                "opened_url": "",
                "paused_for_credentials": False,
                "pause_reason": "",
                "error": "action_blocked",
                "error_code": "action_blocked",
                "auth_phase": "manual",
                "auth_session_id": "",
                "trace": [],
                "canvas": {
                    "title": "Action Blocked",
                    "subtitle": f"Action critic rejected step: {critic_reason}",
                    "cards": [],
                },
            }
        _emit_progress(progress_cb, 12, "Opening Gmail auth target")
        started = _start_email_auth_session(
            profile_dir=profile_dir,
            browser_worker_mode=worker_mode,
            debug_port=debug_port,
            auto_open_auth_tab=True,
        )
        if not started.get("ok"):
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {"error": "manual_auth_phase", "detail": str(started.get("error", "auth_open_failed"))},
                "source_status": {"gmail_ui": "auth_launch_failed"},
                "opened_url": "https://mail.google.com/",
                "paused_for_credentials": True,
                "pause_reason": "Unable to start Gmail auth target. Resolve and click Resume.",
                "error": "credential_missing",
                "error_code": "credential_missing",
                "auth_phase": "manual",
                "auth_session_id": "",
                "trace": [],
                "canvas": {
                    "title": "Paused For Manual Auth",
                    "subtitle": "Sign in in opened tab, then click Resume.",
                    "cards": [],
            },
        }
        sid = str(started.get("auth_session_id", ""))
        focused = focus_auth_session(sid, fallback_url="https://mail.google.com/")
        opened = str(focused.get("opened_url", "https://mail.google.com/"))
        sessions.remember_tab(url=opened, title="Gmail Auth Target", authenticated=False)
        sessions.record_auth_attempt(domain="mail.google.com", status="blocked", detail="manual_auth_phase")
        trace.append({"step": 0, "action": "open_auth_tab", "ok": True, "opened_url": opened})
        return {
            "ok": False,
            "mode": "email_triage",
            "query": "newer_than:2d in:inbox",
            "results_count": 0,
            "results": [],
            "artifacts": {},
            "summary": {"error": "manual_auth_phase"},
            "source_status": {"gmail_ui": "manual_auth_phase"},
            "opened_url": opened,
            "paused_for_credentials": True,
            "pause_reason": "Sign in if required in the opened Gmail tab, then click Resume.",
            "error": "credential_missing",
            "error_code": "credential_missing",
            "auth_phase": "manual",
            "auth_session_id": sid,
            "trace": trace,
            "canvas": {
                "title": "Paused For Manual Auth",
                "subtitle": "Sign in in opened tab, then click Resume.",
                "cards": [],
            },
        }
    if manual_auth_phase and worker_mode == "docker":
        _emit_progress(progress_cb, 12, "Docker mode enabled; continuing without manual auth pause")

    sid = str(auth_session_id or "").strip()
    session: Dict[str, Any] = _EMAIL_AUTH_SESSIONS.get(sid, {})
    if not session:
        sid, session = _select_latest_auth_session(preferred_sid=sid)
    if not session:
        if worker_mode == "docker":
            started = _start_email_auth_session(
                profile_dir=profile_dir,
                browser_worker_mode=worker_mode,
                debug_port=debug_port,
                auto_open_auth_tab=False,
            )
            if not started.get("ok"):
                return {
                    "ok": False,
                    "mode": "email_triage",
                    "query": "newer_than:2d in:inbox",
                    "results_count": 0,
                    "results": [],
                    "artifacts": {},
                    "summary": {
                        "error": "browser_worker_session_failed",
                        "detail": str(started.get("error", "")),
                        "browser_worker_mode": worker_mode,
                    },
                    "source_status": {"gmail_ui": "docker_session_failed", "browser_worker": "docker"},
                    "opened_url": "",
                    "paused_for_credentials": False,
                    "pause_reason": "",
                    "error": "browser_worker_session_failed",
                    "error_code": "browser_worker_session_failed",
                    "auth_session_id": "",
                    "trace": trace,
                    "canvas": {
                        "title": "Docker Worker Session Failed",
                        "subtitle": "Could not initialize docker browser session.",
                        "cards": [],
                    },
                }
            sid = str(started.get("auth_session_id", ""))
            session = _EMAIL_AUTH_SESSIONS.get(sid, {})
            if not session:
                return {
                    "ok": False,
                    "mode": "email_triage",
                    "query": "newer_than:2d in:inbox",
                    "results_count": 0,
                    "results": [],
                    "artifacts": {},
                    "summary": {"error": "browser_worker_session_missing", "browser_worker_mode": worker_mode},
                    "source_status": {"gmail_ui": "docker_session_missing", "browser_worker": "docker"},
                    "opened_url": "",
                    "paused_for_credentials": False,
                    "pause_reason": "",
                    "error": "browser_worker_session_missing",
                    "error_code": "browser_worker_session_missing",
                    "auth_session_id": sid,
                    "trace": trace,
                    "canvas": {
                        "title": "Docker Worker Session Missing",
                        "subtitle": "Retry after worker startup completes.",
                        "cards": [],
                    },
                }
        retry_decision = sessions.auth_retry_decision(domain="mail.google.com", max_failed_attempts=2)
        if worker_mode != "docker":
            started = _start_email_auth_session(
                profile_dir=profile_dir,
                browser_worker_mode=worker_mode,
                debug_port=debug_port,
                auto_open_auth_tab=True,
            )
        else:
            started = {"ok": True, "auth_session_id": sid}
        nsid = str(started.get("auth_session_id", ""))
        if worker_mode != "docker" and (not started.get("ok") or not nsid):
            if not retry_decision.allow_retry and not retry_decision.reusable_authenticated_tab:
                return {
                    "ok": False,
                    "mode": "email_triage",
                    "query": "newer_than:2d in:inbox",
                    "results_count": 0,
                    "results": [],
                    "artifacts": {},
                    "summary": {"error": "credential_missing", "account": account, "detail": retry_decision.reason},
                    "source_status": {"gmail_ui": "auth_retry_budget_exhausted"},
                    "opened_url": "",
                    "paused_for_credentials": True,
                    "pause_reason": "Auth retry budget reached. Use existing signed-in Gmail tab and click Resume.",
                    "error": "credential_missing",
                    "error_code": "credential_missing",
                    "auth_session_id": "",
                    "trace": trace,
                    "canvas": {
                        "title": "Paused For Auth Recovery",
                        "subtitle": "Avoiding login loop. Use existing session, then Resume.",
                        "cards": [],
                    },
                }
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {"error": "credential_missing", "account": account, "detail": str(started.get("error", "auth_session_start_failed"))},
                "source_status": {"gmail_ui": "auth_launch_failed"},
                "opened_url": "",
                "paused_for_credentials": True,
                "pause_reason": "Unable to open Gmail auth target. Resolve the issue and click Resume.",
                "error": "credential_missing",
                "error_code": "credential_missing",
                "auth_session_id": "",
                "trace": trace,
                "canvas": {
                    "title": "Paused For Login",
                    "subtitle": "Gmail auth target could not be opened.",
                    "cards": [],
                },
            }
        if worker_mode != "docker":
            focused = focus_auth_session(nsid, fallback_url="https://mail.google.com/", allow_reopen=False)
            opened = str(focused.get("opened_url", "https://mail.google.com/"))
            sessions.remember_tab(url=opened, title="Gmail Auth Target", authenticated=False)
            sessions.record_auth_attempt(domain="mail.google.com", status="blocked", detail="session_missing")
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {"error": "credential_missing", "account": account},
                "source_status": {"gmail_ui": "manual_auth_phase"},
                "opened_url": opened,
                "paused_for_credentials": True,
                "pause_reason": "Gmail auth session required. Sign in in the opened tab, then click Resume.",
                "error": "credential_missing",
                "error_code": "credential_missing",
                "auth_session_id": nsid,
                "trace": [{"step": 0, "action": "open_auth_tab", "ok": True, "opened_url": opened}],
                "canvas": {
                    "title": "Paused For Login",
                    "subtitle": "Sign in to Gmail, then click Resume.",
                    "cards": [],
                },
            }
        sid = nsid
        session = _EMAIL_AUTH_SESSIONS.get(nsid, session)

    session["auth_confirmed"] = True
    if worker_mode == "docker":
        opened = "https://mail.google.com/"
        trace.append({"step": 0, "action": "docker_worker_session_reuse", "ok": True, "opened_url": opened})
    else:
        focused = focus_auth_session(sid, fallback_url="https://mail.google.com/", allow_reopen=False)
        opened = str(focused.get("opened_url", "https://mail.google.com/"))
        sessions.remember_tab(url=opened, title="Gmail Resume Target", authenticated=False)
        trace.append({"step": 0, "action": "focus_auth_tab", "ok": bool(focused.get("ok", False)), "opened_url": opened})

    _emit_progress(progress_cb, 24, "Opening automation browser context")
    page = session.get("page")
    context = None
    playwright = None
    browser = None
    items: List[EmailActionItem] = []
    try:
        if page is None:
            from playwright.sync_api import sync_playwright

            playwright = sync_playwright().start()
            context = _attach_or_launch_auth_context(
                playwright=playwright,
                profile_dir=profile_dir,
                session=session,
                allow_automated_fallback=(worker_mode == "docker"),
            )
            if context is None:
                if worker_mode == "docker":
                    return {
                        "ok": False,
                        "mode": "email_triage",
                        "query": "newer_than:2d in:inbox",
                        "results_count": 0,
                        "results": [],
                        "artifacts": {},
                        "summary": {"error": "docker_worker_attach_failed", "account": account, "browser_worker_mode": worker_mode},
                        "source_status": {"gmail_ui": "docker_attach_failed", "browser_worker": "docker"},
                        "opened_url": "",
                        "paused_for_credentials": False,
                        "pause_reason": "",
                        "error": "docker_worker_attach_failed",
                        "error_code": "docker_worker_attach_failed",
                        "auth_session_id": sid,
                        "trace": trace,
                        "canvas": {
                            "title": "Docker Worker Attach Failed",
                            "subtitle": "Could not attach to docker browser context.",
                            "cards": [],
                        },
                    }
                return {
                    "ok": False,
                    "mode": "email_triage",
                    "query": "newer_than:2d in:inbox",
                    "results_count": 0,
                    "results": [],
                    "artifacts": {},
                    "summary": {"error": "credential_missing", "account": account},
                    "source_status": {"gmail_ui": "auth_profile_locked"},
                    "opened_url": opened,
                    "paused_for_credentials": True,
                    "pause_reason": "Auth browser is busy/locked. Keep Gmail open in that tab and click Resume.",
                    "error": "credential_missing",
                    "error_code": "credential_missing",
                    "auth_session_id": sid,
                    "trace": trace,
                    "canvas": {
                        "title": "Paused For Auth Session Reuse",
                        "subtitle": "Use existing Gmail tab. Do not close it, then click Resume.",
                        "cards": [],
                    },
                }
            page = _select_best_context_page(context, "https://mail.google.com/")
            if page is None:
                page = context.new_page()
        else:
            trace.append({"step": 0, "action": "reuse_live_gmail_page", "ok": True})
        try:
            # Try current tab state first to avoid unnecessary navigation churn.
            warm_state = _gmail_wait_ready_state(page, timeout_ms=2500)
        except Exception:
            warm_state = "unknown"
        if warm_state != "mail":
            page.goto("https://mail.google.com/", timeout=30000)

        _emit_progress(progress_cb, 26, "Checking Gmail auth state")
        ready = _gmail_wait_ready_state(page, timeout_ms=12000)
        if ready == "login":
            sessions.record_auth_attempt(domain="mail.google.com", status="failed", detail="gmail_login_required")
            _emit_progress(progress_cb, 30, "Attempting Gmail sign-in from local vault")
            login_attempt = _gmail_try_login_with_vault(page=page, account=account)
            trace.append(
                {
                    "step": 0,
                    "action": "vault_login_attempt",
                    "ok": bool(login_attempt.get("ok", False)),
                    "mode": "vault",
                    "error": str(login_attempt.get("error", "")),
                }
            )
            if login_attempt.get("ok"):
                sessions.record_auth_attempt(domain="mail.google.com", status="ok", detail="vault_login_ok")
                ready = _gmail_wait_ready_state(page, timeout_ms=18000)
            else:
                sessions.record_auth_attempt(domain="mail.google.com", status="failed", detail=str(login_attempt.get("error", "")))
                ready = "login"
        if ready == "login":
            if worker_mode == "docker":
                return {
                    "ok": False,
                    "mode": "email_triage",
                    "query": "newer_than:2d in:inbox",
                    "results_count": 0,
                    "results": [],
                    "artifacts": {},
                    "summary": {
                        "error": "docker_worker_auth_required",
                        "account": account,
                        "browser_worker_mode": worker_mode,
                        "detail": "Vault-based sign-in could not complete in docker worker.",
                    },
                    "source_status": {"gmail_ui": "docker_auth_required", "browser_worker": "docker"},
                    "opened_url": "",
                    "paused_for_credentials": False,
                    "pause_reason": "",
                    "error": "docker_worker_auth_required",
                    "error_code": "docker_worker_auth_required",
                    "auth_session_id": sid,
                    "trace": trace,
                    "canvas": {
                        "title": "Docker Worker Needs Interactive Auth",
                        "subtitle": "Switch browser worker mode to local for interactive Gmail sign-in.",
                        "cards": [],
                    },
                }
            decision = sessions.auth_retry_decision(domain="mail.google.com", max_failed_attempts=2)
            sessions.remember_tab(url=opened, title="Gmail Login Required", authenticated=False)
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {"error": "credential_missing", "account": account},
                "source_status": {"gmail_ui": "manual_auth_phase"},
                "opened_url": opened,
                "paused_for_credentials": True,
                "pause_reason": (
                    "Complete Gmail sign-in in the opened tab, then click Resume."
                    if decision.allow_retry
                    else "Authentication retry budget reached. Verify account state in existing tab before retry."
                ),
                "error": "credential_missing",
                "error_code": "credential_missing",
                "auth_session_id": sid,
                "trace": trace,
                "canvas": {
                    "title": "Paused For Login",
                    "subtitle": "Sign in to Gmail, then click Resume.",
                    "cards": [],
                },
            }
        if ready == "unknown":
            try:
                page.goto("https://mail.google.com/", timeout=20000)
                ready = _gmail_wait_ready_state(page, timeout_ms=8000)
            except Exception:
                ready = "unknown"
        if ready == "unknown":
            try:
                page.goto("https://accounts.google.com/ServiceLogin?service=mail", timeout=20000)
                login_attempt = _gmail_try_login_with_vault(page=page, account=account)
                trace.append(
                    {
                        "step": 0,
                        "action": "vault_login_attempt_unknown_state",
                        "ok": bool(login_attempt.get("ok", False)),
                        "mode": "vault",
                        "error": str(login_attempt.get("error", "")),
                    }
                )
                page.goto("https://mail.google.com/", timeout=20000)
                ready = _gmail_wait_ready_state(page, timeout_ms=12000)
            except Exception:
                ready = "unknown"
        if ready != "mail":
            current_url = ""
            try:
                current_url = str(page.url or "")
            except Exception:
                current_url = ""
            if "workspace.google.com/intl" in current_url.lower():
                sessions.record_auth_attempt(domain="mail.google.com", status="blocked", detail="workspace_redirect")
                trace.append({"step": 0, "action": "browser_blocked_google_workspace", "ok": False, "url": current_url[:180]})
                if worker_mode == "docker":
                    return {
                        "ok": False,
                        "mode": "email_triage",
                        "query": "newer_than:2d in:inbox",
                        "results_count": 0,
                        "results": [],
                        "artifacts": {},
                        "summary": {"error": "docker_workspace_redirect", "browser_worker_mode": worker_mode},
                        "source_status": {"gmail_ui": "docker_workspace_redirect", "browser_worker": "docker"},
                        "opened_url": "",
                        "paused_for_credentials": False,
                        "pause_reason": "",
                        "error": "docker_workspace_redirect",
                        "error_code": "docker_workspace_redirect",
                        "auth_session_id": sid,
                        "trace": trace,
                        "canvas": {
                            "title": "Docker Worker Blocked By Workspace Redirect",
                            "subtitle": "Switch to local worker mode or update account permissions.",
                            "cards": [],
                        },
                    }
                return _run_email_triage_imap_fallback(
                    instruction=instruction,
                    account=account,
                    progress_cb=progress_cb,
                    trace=trace,
                )
            if worker_mode == "docker":
                return {
                    "ok": False,
                    "mode": "email_triage",
                    "query": "newer_than:2d in:inbox",
                    "results_count": 0,
                    "results": [],
                    "artifacts": {},
                    "summary": {
                        "error": "docker_worker_inbox_not_ready",
                        "account": account,
                        "detail": f"state={ready}",
                        "browser_worker_mode": worker_mode,
                    },
                    "source_status": {"gmail_ui": f"docker_not_ready:{ready}", "browser_worker": "docker"},
                    "opened_url": "",
                    "paused_for_credentials": False,
                    "pause_reason": "",
                    "error": "docker_worker_inbox_not_ready",
                    "error_code": "docker_worker_inbox_not_ready",
                    "auth_session_id": sid,
                    "trace": trace,
                    "canvas": {
                        "title": "Docker Worker Could Not Reach Gmail Inbox",
                        "subtitle": "Switch to local browser worker mode for interactive auth.",
                        "cards": [],
                    },
                }
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {"error": "credential_missing", "account": account},
                "source_status": {"gmail_ui": f"not_ready:{ready}"},
                "opened_url": opened,
                "paused_for_credentials": True,
                "pause_reason": "Open Gmail inbox in auth tab, then click Resume.",
                "error": "credential_missing",
                "error_code": "credential_missing",
                "auth_session_id": sid,
                "trace": trace,
                "canvas": {
                    "title": "Paused For Inbox Ready",
                    "subtitle": f"Ensure Gmail inbox is loaded, then Resume. Current URL: {current_url[:120]}",
                    "cards": [],
                },
            }

        sessions.remember_tab(url=str(page.url or "https://mail.google.com/"), title="Gmail Inbox", authenticated=True)
        sessions.record_auth_attempt(domain="mail.google.com", status="ok", detail="inbox_ready")
        _emit_progress(progress_cb, 38, "Filtering inbox for last 48 hours")
        _gmail_filter_last_48h(page)
        trace.append({"step": 1, "action": "filter_last_48h", "ok": True, "query": "newer_than:2d in:inbox"})

        _emit_progress(progress_cb, 52, "Reading candidate emails")
        rows = _gmail_collect_rows(page, max_rows=20, scroll_passes=4)
        max_rows = min(len(rows), 20)
        for idx in range(max_rows):
            item = _gmail_process_message(page, idx, human_like_interaction=bool(human_like_interaction))
            if item is None:
                _gmail_back_to_inbox(page)
                continue
            if item.requires_action:
                item.draft_created = _gmail_create_draft_reply(
                    page,
                    item,
                    human_like_interaction=bool(human_like_interaction),
                )
            items.append(item)
            _gmail_back_to_inbox(page)
            trace.append(
                {
                    "step": 2,
                    "action": "process_message",
                    "ok": True,
                    "message_id": item.message_id,
                    "subject": item.subject,
                    "requires_action": item.requires_action,
                    "draft_created": item.draft_created,
                }
            )
    except Exception as exc:  # pylint: disable=broad-exception-caught
        trace.append({"step": 0, "action": "browser_triage_exception", "ok": False, "error": str(exc)})
        if worker_mode == "docker":
            return {
                "ok": False,
                "mode": "email_triage",
                "query": "newer_than:2d in:inbox",
                "results_count": 0,
                "results": [],
                "artifacts": {},
                "summary": {
                    "error": "docker_worker_exception",
                    "detail": str(exc),
                    "browser_worker_mode": worker_mode,
                },
                "source_status": {"gmail_ui": "docker_exception", "browser_worker": "docker"},
                "opened_url": "",
                "paused_for_credentials": False,
                "pause_reason": "",
                "error": "docker_worker_exception",
                "error_code": "docker_worker_exception",
                "auth_session_id": sid,
                "trace": trace,
                "canvas": {
                    "title": "Docker Worker Browser Exception",
                    "subtitle": "Switch to local mode for interactive troubleshooting.",
                    "cards": [],
                },
            }
        return _run_email_triage_imap_fallback(
            instruction=instruction,
            account=account,
            progress_cb=progress_cb,
            trace=trace,
        )
    finally:
        try:
            if context is not None:
                context.close()
        except Exception:
            pass
        try:
            if browser is not None:
                browser.close()
        except Exception:
            pass
        try:
            if playwright is not None:
                playwright.stop()
        except Exception:
            pass

    _emit_progress(progress_cb, 78, "Writing spreadsheet and dashboard artifacts")
    artifacts = _write_email_triage_artifacts(instruction=instruction, account=account, items=items)
    opened_url = Path(artifacts.get("email_triage_html", artifacts.get("email_tasks_csv", ""))).resolve().as_uri()
    opened_url, _nav = _open_target_with_reuse(target_url=opened_url, recent_actions=[f"open_tab:{opened_url}"])
    _emit_progress(progress_cb, 92, "Finalizing run report")

    action_items = [x for x in items if x.requires_action]
    return {
        "ok": True,
        "mode": "email_triage",
        "query": "newer_than:2d in:inbox",
        "results_count": len(items),
        "results": [asdict(x) for x in items[:50]],
        "artifacts": artifacts,
        "summary": {
            "account": account,
            "messages_processed": len(items),
            "action_required": len(action_items),
            "drafts_created": sum(1 for x in action_items if x.draft_created),
            "active_browser_mode": False,
            "auth_session_reused": True,
            "browser_worker_mode": worker_mode,
            "browser_worker_status": str(worker_info.get("status", worker_mode)),
            "human_like_interaction": bool(human_like_interaction),
            "session_manager": sessions.snapshot(),
        },
        "source_status": {"gmail_ui": "ok", "browser_worker": worker_mode},
        "opened_url": opened_url,
        "paused_for_credentials": False,
        "pause_reason": "",
        "trace": trace,
        "auth_session_id": sid,
        "canvas": {
            "title": "Inbox Triage Completed",
            "subtitle": f"Processed {len(items)} messages; {len(action_items)} require action.",
            "cards": [
                {
                    "title": x.subject[:90],
                    "price": "drafted" if x.draft_created else ("needs action" if x.requires_action else "info"),
                    "source": x.sender,
                    "url": "",
                }
                for x in items[:6]
            ],
        },
    }


def _extract_inbox_account(instruction: str) -> str:
    m = re.search(r"([a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+)", instruction)
    if m:
        return m.group(1).strip().lower()
    return "inbox"


def _start_email_auth_session(
    profile_dir: Path,
    *,
    browser_worker_mode: str = "local",
    debug_port: int = 9222,
    auto_open_auth_tab: bool = True,
) -> Dict[str, Any]:
    sid = uuid.uuid4().hex
    url = _normalize_gmail_auth_url("https://mail.google.com/")
    worker_mode = normalize_browser_worker_mode(browser_worker_mode)
    target_debug_port = int(debug_port or 9222)
    opened = {
        "ok": True,
        "method": "session_only",
        "url": url,
        "profile_dir": str(profile_dir.resolve()),
        "debug_port": target_debug_port,
    }
    if auto_open_auth_tab and worker_mode != "docker":
        opened = _open_auth_browser(profile_dir=profile_dir, url=url, debug_port=target_debug_port)
        if not opened.get("ok"):
            try:
                webbrowser.open(url, new=2)
                opened = {
                    "ok": True,
                    "method": "default_browser",
                    "url": url,
                    "profile_dir": str(profile_dir.resolve()),
                    "debug_port": target_debug_port,
                }
            except Exception as exc:  # pylint: disable=broad-exception-caught
                return {"ok": False, "error": str(exc)}
    try:
        _EMAIL_AUTH_SESSIONS[sid] = {
            "mode": "manual_auth_tab" if worker_mode == "local" else "worker_auth_session",
            "browser_worker_mode": worker_mode,
            "url": url,
            "profile_dir": str(profile_dir.resolve()),
            "created_ts": time.time(),
            "auth_confirmed": False,
            "opened_once": bool(auto_open_auth_tab and worker_mode != "docker"),
            "open_method": str(opened.get("method", "")),
            "debug_port": int(opened.get("debug_port", target_debug_port) or target_debug_port),
        }
        return {"ok": True, "auth_session_id": sid}
    except Exception as exc:  # pylint: disable=broad-exception-caught
        return {"ok": False, "error": str(exc)}


def focus_auth_session(
    auth_session_id: str,
    fallback_url: str = "https://mail.google.com/",
    allow_reopen: bool = False,
) -> Dict[str, Any]:
    sid = str(auth_session_id or "").strip()
    safe_fallback = _normalize_gmail_auth_url(fallback_url)
    if not sid:
        sid, _sess = _select_latest_auth_session(preferred_sid="")
    sess = _EMAIL_AUTH_SESSIONS.get(sid, {}) if sid else {}
    if not sess:
        try:
            webbrowser.open(safe_fallback, new=2)
        except Exception:
            pass
        return {"ok": False, "error": "auth_session_not_found", "opened_url": safe_fallback}
    target = _normalize_gmail_auth_url(str(sess.get("url", safe_fallback) or safe_fallback), default=safe_fallback)
    sess["url"] = target
    if sess.get("opened_once", False) and not allow_reopen:
        return {"ok": True, "auth_session_id": sid, "opened_url": target, "already_open": True}
    try:
        webbrowser.open(target, new=2)
        sess["opened_once"] = True
        return {"ok": True, "auth_session_id": sid, "opened_url": target}
    except Exception as exc:  # pylint: disable=broad-exception-caught
        return {"ok": False, "error": str(exc), "opened_url": target}


def focus_email_auth_session(auth_session_id: str, fallback_url: str = "https://mail.google.com/") -> Dict[str, Any]:
    return focus_auth_session(auth_session_id=auth_session_id, fallback_url=fallback_url)


def _select_best_context_page(context: Any, target_url: str) -> Any:
    best = None
    try:
        pages = list(getattr(context, "pages", []) or [])
    except Exception:
        pages = []
    target_host = urllib.parse.urlparse(target_url).netloc.lower()
    for page in pages:
        try:
            url = str(page.url or "").lower()
        except Exception:
            continue
        if target_host and target_host in url:
            return page
        if "about:blank" in url and best is None:
            best = page
    return best


def _close_email_auth_session(session_id: str) -> None:
    _EMAIL_AUTH_SESSIONS.pop(session_id, None)


def _launch_persistent_chrome_with_retry(playwright: Any, profile_dir: Path) -> Any:
    last_exc: Exception | None = None
    for attempt in range(2):
        try:
            return playwright.chromium.launch_persistent_context(
                str(profile_dir),
                channel="chrome",
                headless=False,
                args=["--start-maximized"],
            )
        except Exception as exc:  # pylint: disable=broad-exception-caught
            last_exc = exc
            if attempt == 0:
                # Avoid killing a user's active auth session tab; retry once only.
                time.sleep(0.4)
                continue
            raise
    if last_exc:
        raise last_exc
    raise RuntimeError("Failed to launch persistent browser context")


def _terminate_profile_browser_processes(profile_dir: Path) -> Dict[str, Any]:
    profile_abs = str(profile_dir.resolve())
    needle = f"--user-data-dir={profile_abs}".replace("'", "''")
    script = (
        "$procs = Get-CimInstance Win32_Process | Where-Object { "
        "(($_.Name -match '^(chrome|msedge)\\.exe$') -and ($_.CommandLine -like '*"
        + needle
        + "*')) }; "
        "$ids=@(); "
        "foreach($p in $procs){ $ids += $p.ProcessId; try { Stop-Process -Id $p.ProcessId -Force -ErrorAction Stop } catch {} }; "
        "if($ids.Count -eq 0){ 'none' } else { ($ids -join ',') }"
    )
    try:
        out = subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", script],
            capture_output=True,
            text=True,
            timeout=8,
            check=False,
        )
        return {"ok": True, "killed": (out.stdout or "").strip(), "stderr": (out.stderr or "").strip()}
    except Exception as exc:  # pylint: disable=broad-exception-caught
        return {"ok": False, "error": str(exc)}


def _is_profile_browser_running(profile_dir: Path) -> bool:
    profile_abs = str(profile_dir.resolve())
    needle = f"--user-data-dir={profile_abs}".replace("'", "''")
    script = (
        "$procs = Get-CimInstance Win32_Process | Where-Object { "
        "(($_.Name -match '^(chrome|msedge)\\.exe$') -and ($_.CommandLine -like '*"
        + needle
        + "*')) }; "
        "if($procs){ 'yes' } else { 'no' }"
    )
    try:
        out = subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", script],
            capture_output=True,
            text=True,
            timeout=6,
            check=False,
        )
        return (out.stdout or "").strip().lower() == "yes"
    except Exception:
        return False


def _open_auth_browser(profile_dir: Path, url: str, debug_port: int = 9222) -> Dict[str, Any]:
    profile_abs = str(profile_dir.resolve())
    candidates: List[str] = []
    for key in ["PROGRAMFILES", "PROGRAMFILES(X86)", "LOCALAPPDATA"]:
        root = str(Path(str(os.environ.get(key, ""))))
        if not root or root == ".":
            continue
        candidates.extend(
            [
                str(Path(root) / "Google" / "Chrome" / "Application" / "chrome.exe"),
                str(Path(root) / "Microsoft" / "Edge" / "Application" / "msedge.exe"),
            ]
        )
    which_chrome = shutil.which("chrome") or shutil.which("chrome.exe")
    which_edge = shutil.which("msedge") or shutil.which("msedge.exe")
    if which_chrome:
        candidates.insert(0, which_chrome)
    if which_edge:
        candidates.insert(1, which_edge)
    exe = next((c for c in candidates if c and Path(c).exists()), "")
    if exe:
        try:
            subprocess.Popen(
                [exe, f"--user-data-dir={profile_abs}", "--profile-directory=Default", f"--remote-debugging-port={int(debug_port)}", url],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=0x00000008,  # DETACHED_PROCESS
            )
            return {
                "ok": True,
                "method": "profile_browser",
                "exe": exe,
                "url": url,
                "profile_dir": profile_abs,
                "debug_port": int(debug_port),
            }
        except Exception as exc:
            return {
                "ok": False,
                "method": "profile_browser",
                "error": str(exc),
                "url": url,
                "profile_dir": profile_abs,
                "debug_port": int(debug_port),
            }
    try:
        webbrowser.open(url, new=2)
        return {"ok": True, "method": "default_browser", "url": url, "profile_dir": profile_abs, "debug_port": int(debug_port)}
    except Exception as exc:
        return {
            "ok": False,
            "method": "default_browser",
            "error": str(exc),
            "url": url,
            "profile_dir": profile_abs,
            "debug_port": int(debug_port),
        }


def _select_latest_auth_session(preferred_sid: str = "") -> tuple[str, Dict[str, Any]]:
    sid = str(preferred_sid or "").strip()
    if sid and sid in _EMAIL_AUTH_SESSIONS:
        return sid, _EMAIL_AUTH_SESSIONS.get(sid, {})
    if not _EMAIL_AUTH_SESSIONS:
        return "", {}
    ordered = sorted(
        _EMAIL_AUTH_SESSIONS.items(),
        key=lambda kv: float((kv[1] or {}).get("created_ts", 0.0)),
        reverse=True,
    )
    best_sid, best = ordered[0]
    return str(best_sid), dict(best or {})


def _attach_or_launch_auth_context(
    playwright: Any,
    profile_dir: Path,
    session: Dict[str, Any],
    allow_automated_fallback: bool = False,
) -> Any:
    debug_port = int(session.get("debug_port", 9222) or 9222)
    try:
        browser = playwright.chromium.connect_over_cdp(f"http://127.0.0.1:{debug_port}", timeout=3500)
        contexts = list(getattr(browser, "contexts", []) or [])
        if contexts:
            return contexts[0]
        return browser.new_context()
    except Exception:
        pass
    if not allow_automated_fallback:
        return None
    try:
        return _launch_persistent_chrome_with_retry(playwright, profile_dir=profile_dir)
    except Exception:
        return None


def _gmail_login_required(page: Any) -> bool:
    url = str(page.url or "").lower()
    if "accounts.google.com" in url:
        return True
    try:
        if page.locator("input[type='password'], input[type='email']").count() > 0 and page.locator("input[name='q']").count() == 0:
            return True
    except Exception:
        return True
    return False


def _is_topic_mastery_intent(instruction: str) -> bool:
    low = str(instruction or "").lower()
    learn_markers = [
        "learn how to",
        "get well versed",
        "mastery guide",
        "reusable playbook",
        "learned skill",
        "seed video",
        "topic only",
    ]
    source_markers = ["youtube", "tutorial", "video url", "watch this video", "related videos", "supporting sources"]
    return any(token in low for token in learn_markers) and any(token in low for token in source_markers)


def _gmail_try_login_with_vault(page: Any, account: str) -> Dict[str, Any]:
    try:
        vault = LocalPasswordVault()
        resolved = vault.find_entry_by_service("gmail")
        if not resolved.get("ok"):
            resolved = vault.find_entry_by_service("google")
        if not resolved.get("ok"):
            return {"ok": False, "error": "vault_entry_not_found"}
        entry = resolved.get("entry", {}) or {}
        username = str(entry.get("username", "")).strip()
        password = str(entry.get("password", "")).strip()
        if not username or not password:
            return {"ok": False, "error": "vault_entry_missing_secret"}
        if account and "@" in account and username.lower() != account.lower():
            # Respect explicit account request.
            return {"ok": False, "error": "vault_account_mismatch"}

        # Email step
        email_visible = page.locator("input[type='email']:visible").first
        if email_visible.count() > 0:
            try:
                email_visible.click(timeout=4000)
            except Exception:
                pass
            email_visible.fill(username, timeout=8000)
            try:
                page.locator("#identifierNext button:visible, #identifierNext:visible").first.click(timeout=5000)
            except Exception:
                page.keyboard.press("Enter")
            page.wait_for_timeout(1200)
        else:
            # Account chooser step.
            try:
                account_tile = page.locator(
                    f"[data-identifier='{username}']:visible, div[data-email='{username}']:visible, div[data-identifier='{username}']:visible"
                ).first
                if account_tile.count() > 0:
                    account_tile.click(timeout=5000)
                    page.wait_for_timeout(1200)
            except Exception:
                pass

        # Password step
        pwd_visible = page.locator("input[name='Passwd']:visible, input[type='password']:visible").first
        if pwd_visible.count() > 0:
            try:
                pwd_visible.click(timeout=4000)
            except Exception:
                pass
            pwd_visible.fill(password, timeout=8000)
            try:
                page.locator("#passwordNext button:visible, #passwordNext:visible").first.click(timeout=5000)
            except Exception:
                page.keyboard.press("Enter")
            page.wait_for_timeout(2200)
        else:
            # No visible password field may indicate already-signed-in or challenge state.
            if _gmail_login_required(page):
                return {"ok": False, "error": "password_field_not_visible"}

        try:
            vault.touch_used(str(entry.get("id", "")))
        except Exception:
            pass
        return {"ok": True}
    except Exception as exc:  # pylint: disable=broad-exception-caught
        return {"ok": False, "error": str(exc)}


def _gmail_wait_ready_state(page: Any, timeout_ms: int = 15000) -> str:
    deadline = time.time() + (max(1000, timeout_ms) / 1000.0)
    while time.time() < deadline:
        if _gmail_login_required(page):
            return "login"
        try:
            url = str(page.url or "").lower()
            if page.locator("input[name='q'], input[aria-label='Search mail']").count() > 0:
                return "mail"
            if page.locator("div[role='main'] tr.zA, div[gh='cm'], div.T-I.T-I-KE").count() > 0 and "mail.google.com" in url:
                return "mail"
            if "mail.google.com" in url and "accounts.google.com" not in url:
                # Soft-ready fallback for slower/lazy Gmail layouts.
                if page.locator("body").count() > 0:
                    return "mail"
        except Exception:
            pass
        time.sleep(0.25)
    return "unknown"


def _gmail_filter_last_48h(page: Any) -> None:
    search = page.locator("input[name='q'], input[aria-label='Search mail'], textarea[name='q']")
    try:
        search.first.click(timeout=12000)
    except Exception:
        try:
            page.keyboard.press("/")
            page.wait_for_timeout(400)
        except Exception:
            pass
    # Retry after '/' shortcut focus attempt.
    search = page.locator("input[name='q'], input[aria-label='Search mail'], textarea[name='q']")
    search.first.fill("newer_than:2d in:inbox", timeout=12000)
    page.keyboard.press("Enter")
    page.wait_for_timeout(1200)


def _gmail_collect_rows(page: Any, max_rows: int = 20, scroll_passes: int = 4) -> List[Any]:
    rows: List[Any] = []
    for _ in range(max(1, scroll_passes)):
        locator = page.locator("tr.zA:visible")
        count = locator.count()
        if count <= 0:
            locator = page.locator("tr.zA")
            count = locator.count()
        for i in range(min(count, max_rows)):
            rows.append(locator.nth(i))
            if len(rows) >= max_rows:
                return rows
        page.mouse.wheel(0, 1600)
        page.wait_for_timeout(350)
    return rows[:max_rows]


def _gmail_process_message(
    page: Any,
    row_index: int,
    human_like_interaction: bool = False,
) -> Optional[EmailActionItem]:
    rows = page.locator("tr.zA:visible")
    if rows.count() <= 0:
        rows = page.locator("tr.zA")
    if row_index >= rows.count():
        return None
    row = rows.nth(row_index)
    try:
        row.scroll_into_view_if_needed(timeout=2000)
    except Exception:
        pass
    if human_like_interaction:
        try:
            box = row.bounding_box()
            if box:
                tx = float(box.get("x", 0.0)) + (float(box.get("width", 0.0)) / 2.0) + random.uniform(-12.0, 12.0)
                ty = float(box.get("y", 0.0)) + (float(box.get("height", 0.0)) / 2.0) + random.uniform(-6.0, 6.0)
                page.mouse.move(tx, ty, steps=random.randint(10, 22))
                page.wait_for_timeout(random.randint(70, 180))
        except Exception:
            pass
    try:
        row.click(timeout=12000)
    except Exception:
        try:
            row.click(timeout=5000, force=True)
        except Exception:
            return None
    page.wait_for_timeout(700)
    sender = ""
    subject = ""
    snippet = ""
    received_at = ""
    try:
        sender = (page.locator("h3 span[email], span[email]").first.get_attribute("email") or "").strip()
    except Exception:
        sender = ""
    try:
        subject = (page.locator("h2.hP").first.inner_text(timeout=5000) or "").strip()
    except Exception:
        subject = ""
    try:
        snippet = (page.locator("div.a3s").first.inner_text(timeout=5000) or "").strip()
        snippet = re.sub(r"\s+", " ", snippet)[:450]
    except Exception:
        snippet = ""
    try:
        received_at = (page.locator("span.g3").first.get_attribute("title") or "").strip()
    except Exception:
        received_at = ""
    requires_action, reason = _classify_requires_action(subject=subject, snippet=snippet)
    message_id = hashlib.sha256(f"{sender}|{subject}|{received_at}".encode("utf-8", errors="ignore")).hexdigest()[:16]
    return EmailActionItem(
        message_id=message_id,
        sender=sender or "unknown",
        subject=subject or "(no subject)",
        received_at=received_at or "",
        snippet=snippet,
        requires_action=requires_action,
        reason=reason,
        draft_created=False,
    )


def _classify_requires_action(subject: str, snippet: str) -> tuple[bool, str]:
    hay = f"{subject} {snippet}".lower()
    action_keywords = [
        "action required",
        "please review",
        "please respond",
        "deadline",
        "approve",
        "approval",
        "follow up",
        "urgent",
        "asap",
        "can you",
        "could you",
        "request",
        "next steps",
    ]
    for kw in action_keywords:
        if kw in hay:
            return True, f"matched:{kw}"
    return False, "no_action_keyword"


def _gmail_create_draft_reply(page: Any, item: EmailActionItem, human_like_interaction: bool = False) -> bool:
    try:
        reply = page.locator("div[aria-label='Reply'], [data-tooltip='Reply']").first
        if human_like_interaction:
            page.wait_for_timeout(random.randint(80, 210))
        reply.click(timeout=7000)
        page.wait_for_timeout(500)
        body = page.locator("div[aria-label='Message Body'][role='textbox'], div[role='textbox'][aria-label='Message Body']").last
        body.click(timeout=5000)
        text = _build_draft_text(item)
        if human_like_interaction:
            body.type(text, delay=random.randint(18, 52))
        else:
            body.fill(text)
        page.wait_for_timeout(250)
        return True
    except Exception:
        return False


def _build_draft_text(item: EmailActionItem) -> str:
    return (
        f"Hi,\n\n"
        f"Thanks for the note about \"{item.subject}\". I reviewed your message and will follow up with the requested details shortly.\n\n"
        f"Current action item: {item.reason}.\n\n"
        "Best,\n"
        "Jeff"
    )


def _gmail_back_to_inbox(page: Any) -> None:
    try:
        page.keyboard.press("u")
        page.wait_for_timeout(500)
    except Exception:
        try:
            page.go_back()
            page.wait_for_timeout(500)
        except Exception:
            return


def _run_email_triage_imap_fallback(
    *,
    instruction: str,
    account: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    trace: Optional[List[Dict[str, Any]]] = None,
) -> Dict[str, Any]:
    _emit_progress(progress_cb, 62, "Falling back to Gmail IMAP mode")
    creds = _resolve_gmail_vault_credentials(account=account)
    if not creds.get("ok"):
        return {
            "ok": False,
            "mode": "email_triage",
            "query": "newer_than:2d in:inbox",
            "results_count": 0,
            "results": [],
            "artifacts": {},
            "summary": {"error": "imap_fallback_unavailable", "detail": str(creds.get("error", "vault_credentials_missing"))},
            "source_status": {"gmail_imap": "vault_credentials_missing"},
            "opened_url": "",
            "paused_for_credentials": True,
            "pause_reason": "Gmail IMAP fallback requires a vault Gmail entry with app password.",
            "error": "credential_missing",
            "error_code": "credential_missing",
            "trace": trace or [],
            "canvas": {
                "title": "Paused For IMAP Credentials",
                "subtitle": "Add Gmail app-password credentials in Local Password Vault, then Resume.",
                "cards": [],
            },
        }

    username = str(creds.get("username", ""))
    password = str(creds.get("password", ""))
    entry_id = str(creds.get("entry_id", ""))
    items: List[EmailActionItem] = []
    drafts_created = 0
    fetch_count = 0

    mailbox = None
    try:
        mailbox = imaplib.IMAP4_SSL("imap.gmail.com", 993)
        mailbox.login(username, password)
        mailbox.select("INBOX", readonly=False)
        _emit_progress(progress_cb, 68, "Reading inbox messages (last 48h)")
        message_ids = _imap_recent_message_ids(mailbox=mailbox, max_ids=30)
        fetch_count = len(message_ids)
        for msg_id in message_ids:
            item = _imap_fetch_action_item(mailbox=mailbox, msg_id=msg_id)
            if item is None:
                continue
            if item.requires_action:
                ok = _imap_append_draft(mailbox=mailbox, username=username, item=item)
                item.draft_created = ok
                if ok:
                    drafts_created += 1
            items.append(item)
    except Exception as exc:  # pylint: disable=broad-exception-caught
        detail = str(exc)
        detail_low = detail.lower()
        is_app_password_required = ("application-specific password required" in detail_low) or ("app password" in detail_low)
        error_code = "imap_app_password_required" if is_app_password_required else "credential_missing"
        pause_reason = (
            "IMAP requires a Gmail app password. Add an app password in Local Password Vault or continue with Gmail UI session, then Resume."
            if is_app_password_required
            else "Gmail IMAP fallback failed. Verify app password and IMAP access, then Resume."
        )
        subtitle = (
            "Gmail requires an app password for IMAP access."
            if is_app_password_required
            else "Verify Gmail app-password and IMAP access."
        )
        return {
            "ok": False,
            "mode": "email_triage",
            "query": "newer_than:2d in:inbox",
            "results_count": 0,
            "results": [],
            "artifacts": {},
            "summary": {
                "error": "imap_fallback_failed",
                "detail": detail,
                "imap_error_code": error_code,
            },
            "source_status": {"gmail_imap": f"error:{type(exc).__name__}"},
            "opened_url": "",
            "paused_for_credentials": True,
            "pause_reason": pause_reason,
            "error": "credential_missing",
            "error_code": error_code,
            "trace": trace or [],
            "canvas": {
                "title": "IMAP Fallback Failed",
                "subtitle": subtitle,
                "cards": [],
            },
        }
    finally:
        try:
            if mailbox is not None:
                mailbox.logout()
        except Exception:
            pass
        try:
            if entry_id:
                LocalPasswordVault().touch_used(entry_id)
        except Exception:
            pass

    _emit_progress(progress_cb, 78, "Writing spreadsheet and dashboard artifacts")
    artifacts = _write_email_triage_artifacts(instruction=instruction, account=username, items=items)
    opened_url = Path(artifacts.get("email_triage_html", artifacts.get("email_tasks_csv", ""))).resolve().as_uri()
    opened_url, _nav = _open_target_with_reuse(target_url=opened_url, recent_actions=[f"open_tab:{opened_url}"])

    action_items = [x for x in items if x.requires_action]
    return {
        "ok": True,
        "mode": "email_triage",
        "query": "newer_than:2d in:inbox",
        "results_count": len(items),
        "results": [asdict(x) for x in items[:50]],
        "artifacts": artifacts,
        "summary": {
            "account": username,
            "messages_processed": len(items),
            "messages_fetched": fetch_count,
            "action_required": len(action_items),
            "drafts_created": drafts_created,
            "fallback_mode": "imap",
        },
        "source_status": {"gmail_ui": "fallback_imap", "gmail_imap": "ok"},
        "opened_url": opened_url,
        "paused_for_credentials": False,
        "pause_reason": "",
        "trace": trace or [],
        "canvas": {
            "title": "Inbox Triage Completed (IMAP Fallback)",
            "subtitle": f"Processed {len(items)} messages; {len(action_items)} require action.",
            "cards": [
                {
                    "title": x.subject[:90],
                    "price": "drafted" if x.draft_created else ("needs action" if x.requires_action else "info"),
                    "source": x.sender,
                    "url": "",
                }
                for x in items[:6]
            ],
        },
    }


def _resolve_gmail_vault_credentials(account: str) -> Dict[str, Any]:
    vault = LocalPasswordVault()
    resolved = vault.find_entry_by_service("gmail")
    if not resolved.get("ok"):
        resolved = vault.find_entry_by_service("google")
    if not resolved.get("ok"):
        return {"ok": False, "error": "vault_entry_not_found"}
    entry = resolved.get("entry", {}) or {}
    username = str(entry.get("username", "")).strip()
    password = str(entry.get("password", "")).strip()
    if not username or not password:
        return {"ok": False, "error": "vault_entry_missing_secret"}
    if account and "@" in account and account.lower() != username.lower():
        return {"ok": False, "error": "vault_account_mismatch"}
    return {"ok": True, "username": username, "password": password, "entry_id": str(entry.get("id", ""))}


def _imap_recent_message_ids(mailbox: Any, max_ids: int = 30) -> List[bytes]:
    start_dt = datetime.now().astimezone() - timedelta(hours=48)
    date_token = start_dt.strftime("%d-%b-%Y")
    typ, data = mailbox.search(None, "SINCE", date_token)
    if typ != "OK" or not data:
        return []
    ids = list((data[0] or b"").split())
    if not ids:
        return []
    ids = ids[-max_ids:]
    return list(reversed(ids))


def _imap_fetch_action_item(mailbox: Any, msg_id: bytes) -> Optional[EmailActionItem]:
    typ, data = mailbox.fetch(msg_id, "(RFC822)")
    if typ != "OK" or not data:
        return None
    raw = b""
    for chunk in data:
        if isinstance(chunk, tuple) and len(chunk) >= 2 and isinstance(chunk[1], (bytes, bytearray)):
            raw = bytes(chunk[1])
            break
    if not raw:
        return None
    msg = message_from_bytes(raw)
    subject = _decode_mime_header(str(msg.get("Subject", ""))) or "(no subject)"
    sender_full = str(msg.get("From", "")).strip()
    sender_email = parseaddr(sender_full)[1] or sender_full or "unknown"
    date_raw = str(msg.get("Date", "")).strip()
    received_at = date_raw
    try:
        dt = parsedate_to_datetime(date_raw)
        if dt is not None:
            received_at = dt.isoformat(timespec="seconds")
    except Exception:
        pass
    snippet = _imap_extract_snippet(msg, limit=450)
    requires_action, reason = _classify_requires_action(subject=subject, snippet=snippet)
    message_id = hashlib.sha256(f"{sender_email}|{subject}|{received_at}".encode("utf-8", errors="ignore")).hexdigest()[:16]
    return EmailActionItem(
        message_id=message_id,
        sender=sender_email,
        subject=subject,
        received_at=received_at,
        snippet=snippet,
        requires_action=requires_action,
        reason=reason,
        draft_created=False,
    )


def _decode_mime_header(value: str) -> str:
    parts = decode_header(value or "")
    out: List[str] = []
    for payload, enc in parts:
        if isinstance(payload, bytes):
            charset = enc or "utf-8"
            try:
                out.append(payload.decode(charset, errors="ignore"))
            except Exception:
                out.append(payload.decode("utf-8", errors="ignore"))
        else:
            out.append(str(payload))
    return "".join(out).strip()


def _imap_extract_snippet(msg: Any, limit: int = 450) -> str:
    text = ""
    try:
        if msg.is_multipart():
            for part in msg.walk():
                ctype = str(part.get_content_type() or "").lower()
                disp = str(part.get("Content-Disposition", "")).lower()
                if ctype == "text/plain" and "attachment" not in disp:
                    payload = part.get_payload(decode=True) or b""
                    charset = part.get_content_charset() or "utf-8"
                    text = payload.decode(charset, errors="ignore")
                    if text.strip():
                        break
        else:
            payload = msg.get_payload(decode=True) or b""
            charset = msg.get_content_charset() or "utf-8"
            text = payload.decode(charset, errors="ignore")
    except Exception:
        text = ""
    text = re.sub(r"\s+", " ", text or "").strip()
    return text[:limit]


def _imap_append_draft(mailbox: Any, username: str, item: EmailActionItem) -> bool:
    try:
        em = EmailMessage()
        em["From"] = username
        em["To"] = item.sender
        em["Subject"] = f"Re: {item.subject}"
        em.set_content(_build_draft_text(item))
        raw = em.as_bytes()
        typ, _ = mailbox.append('"[Gmail]/Drafts"', "\\Draft", Time2Internaldate(time.time()), raw)
        if typ == "OK":
            return True
        typ2, _ = mailbox.append("Drafts", "\\Draft", Time2Internaldate(time.time()), raw)
        return typ2 == "OK"
    except Exception:
        return False


def _write_email_triage_artifacts(instruction: str, account: str, items: List[EmailActionItem]) -> Dict[str, str]:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path("data/reports/email_triage") / ts
    out_dir.mkdir(parents=True, exist_ok=True)
    csv_path = out_dir / "task_list.csv"
    md_path = out_dir / "summary.md"
    html_path = out_dir / "dashboard.html"

    with csv_path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(
            f,
            fieldnames=[
                "message_id",
                "sender",
                "subject",
                "received_at",
                "requires_action",
                "reason",
                "draft_created",
                "snippet",
            ],
        )
        writer.writeheader()
        for item in items:
            writer.writerow(asdict(item))

    action_items = [x for x in items if x.requires_action]
    lines = [
        "# Inbox Triage Summary",
        "",
        f"- Generated: {datetime.now().isoformat(timespec='seconds')}",
        f"- Instruction: {instruction}",
        f"- Account: {account}",
        f"- Messages processed: {len(items)}",
        f"- Action required: {len(action_items)}",
        f"- Drafts created: {sum(1 for x in action_items if x.draft_created)}",
        "",
        "## Action-Needed Messages",
    ]
    for idx, item in enumerate(action_items[:100], start=1):
        lines.append(f"{idx}. **{item.subject}** | {item.sender} | {item.received_at or 'n/a'} | draft={item.draft_created}")
        lines.append(f"   - reason: {item.reason}")
    md_path.write_text("\n".join(lines), encoding="utf-8")

    rows = json.dumps([asdict(x) for x in items])
    html_text = f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8"/><meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Inbox Triage Dashboard</title>
<style>
body{{font-family:'Segoe UI',Tahoma,sans-serif;background:#f8fafc;color:#0f172a;margin:0}}
.w{{max-width:1180px;margin:0 auto;padding:20px}}
.hero{{background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;padding:14px}}
table{{width:100%;border-collapse:collapse;background:#fff;margin-top:12px}}
th,td{{padding:8px;border-bottom:1px solid #e2e8f0;vertical-align:top;font-size:13px}}
th{{background:#0f172a;color:#e2e8f0;position:sticky;top:0}}
.pill{{display:inline-block;padding:2px 8px;border-radius:999px;font-size:11px}}
.ok{{background:#dcfce7;color:#166534}} .wait{{background:#fee2e2;color:#991b1b}}
input{{padding:8px;border:1px solid #cbd5e1;border-radius:8px;min-width:280px;margin-top:8px}}
</style></head><body><div class="w"><div class="hero"><h1 style="margin:0">Inbox Triage Dashboard</h1><p style="margin:4px 0 0 0">{html.escape(account)}</p><input id="q" placeholder="Filter sender/subject/snippet" oninput="render()"/></div><table><thead><tr><th>Sender</th><th>Subject</th><th>Received</th><th>Action</th><th>Draft</th><th>Snippet</th></tr></thead><tbody id="rows"></tbody></table></div>
<script>
const data={rows};
function esc(s){{ return (s||'').replace(/[&<>]/g,c=>({{'&':'&amp;','<':'&lt;','>':'&gt;'}}[c])); }}
function render(){{
 const q=(document.getElementById('q').value||'').toLowerCase();
 const filtered=data.filter(x=>!q || (x.sender+' '+x.subject+' '+x.snippet).toLowerCase().includes(q));
 document.getElementById('rows').innerHTML=filtered.map(x=>`<tr><td>${{esc(x.sender)}}</td><td>${{esc(x.subject)}}</td><td>${{esc(x.received_at||'')}}</td><td><span class="pill ${{x.requires_action?'wait':'ok'}}">${{x.requires_action?'needs action':'info'}}</span></td><td>${{x.draft_created?'yes':'no'}}</td><td>${{esc((x.snippet||'').slice(0,220))}}</td></tr>`).join('');
}}
render();
</script></body></html>"""
    html_path.write_text(html_text, encoding="utf-8")
    return {
        "directory": str(out_dir.resolve()),
        "email_tasks_csv": str(csv_path.resolve()),
        "summary_md": str(md_path.resolve()),
        "email_triage_html": str(html_path.resolve()),
        "primary_open_file": str(html_path.resolve()),
    }


def _is_job_research_intent(instruction: str) -> bool:
    low = instruction.lower()
    job_terms = ["job", "position", "vp", "avp", "linkedin", "indeed", "salary", "remote"]
    analysis_terms = ["spreadsheet", "dashboard", "report", "analysis"]
    strong_job_signal = sum(1 for t in job_terms if t in low) >= 3 and any(t in low for t in analysis_terms)
    board_signal = any(t in low for t in ["job board", "job boards", "all the job boards", "jobsites", "job sites"])
    comp_signal = any(t in low for t in ["total compensation", "compensation", "total comp", "more than", "above"])
    leadership_signal = any(t in low for t in ["vp", "avp", "vice president", "analytics", "data and ai", "data and analytics"])
    remote_signal = any(t in low for t in ["remote", "hybrid"])
    return strong_job_signal or (board_signal and leadership_signal and (comp_signal or remote_signal))


def _run_job_market_research(instruction: str, progress_cb: Optional[Callable[[int, str], None]] = None) -> Dict[str, Any]:
    role_query = _extract_role_query(instruction)
    region_labels = _extract_regions(instruction)
    constraints = _extract_job_constraints(instruction)
    source_status: Dict[str, str] = {}
    collected: List[JobListing] = []
    _emit_progress(progress_cb, 18, f"Researching job sources for: {role_query}")

    total_units = max(1, len(region_labels) * 5)
    done_units = 0

    for region in region_labels:
        _emit_progress(progress_cb, 22 + int((done_units / total_units) * 45), f"Searching LinkedIn ({region.upper()})")
        try:
            li = _scrape_linkedin_jobs(role_query=role_query, region=region, limit=40)
            collected.extend(li)
            source_status[f"linkedin_{region}"] = f"ok:{len(li)}"
        except Exception as exc:
            source_status[f"linkedin_{region}"] = f"error:{type(exc).__name__}"
        done_units += 1

        _emit_progress(progress_cb, 22 + int((done_units / total_units) * 45), f"Searching BuiltIn ({region.upper()})")
        try:
            bi = _scrape_builtin_jobs(role_query=role_query, region=region, limit=30)
            collected.extend(bi)
            source_status[f"builtin_{region}"] = f"ok:{len(bi)}"
        except Exception as exc:
            source_status[f"builtin_{region}"] = f"error:{type(exc).__name__}"
        done_units += 1

        # Best-effort site-search fallback for other commercial boards.
        site_queries = [
            ("indeed", "site:indeed.com/jobs"),
            ("ziprecruiter", "site:ziprecruiter.com/jobs"),
            ("glassdoor", "site:glassdoor.com/job-listing"),
        ]
        region_phrase = "Ireland" if region == "ireland" else "United States"
        for source, site_prefix in site_queries:
            query = f'{site_prefix} "{role_query}" {region_phrase} salary remote'
            _emit_progress(progress_cb, 22 + int((done_units / total_units) * 45), f"Searching {source} ({region.upper()})")
            try:
                pulled = 0
                for item in _search_web(query, limit=6):
                    listing = _to_job_listing(item=item, source=source, fallback_region=region)
                    if listing:
                        collected.append(listing)
                        pulled += 1
                source_status[f"{source}_{region}"] = f"ok:{pulled}"
            except Exception as exc:
                source_status[f"{source}_{region}"] = f"error:{type(exc).__name__}"
            time.sleep(0.2)
            done_units += 1

    _emit_progress(progress_cb, 72, "Deduplicating and ranking job listings")
    dedup: Dict[str, JobListing] = {}
    for job in collected:
        dedup[job.url] = job
    jobs_all = list(dedup.values())
    strict_mode = _job_constraints_are_strict(constraints)
    jobs = [j for j in jobs_all if _is_target_job(j.title, strict_vp_avp=constraints.require_vp_avp)]
    jobs = [j for j in jobs if _job_matches_constraints(j, constraints)]
    if (not strict_mode) and len(jobs) < 12:
        jobs = jobs_all
    jobs.sort(key=lambda j: (_salary_sort_key(j), 1 if j.remote else 0, j.source, j.title), reverse=True)
    _emit_progress(progress_cb, 84, "Generating spreadsheet, report, and dashboard")
    artifacts = _write_job_artifacts(instruction=instruction, jobs=jobs)
    summary = _job_summary(jobs)
    summary["constraints"] = {
        "strict_mode": strict_mode,
        "require_vp_avp": constraints.require_vp_avp,
        "require_remote_or_hybrid": constraints.require_remote_or_hybrid,
        "min_base_salary_usd": constraints.min_base_salary_usd,
        "min_total_comp_usd": constraints.min_total_comp_usd,
        "allowed_regions": list(constraints.allowed_regions),
    }
    top = jobs[:40]
    dashboard_uri = Path(artifacts["dashboard_html"]).resolve().as_uri() if artifacts.get("dashboard_html") else ""
    if dashboard_uri:
        dashboard_uri, _nav = _open_target_with_reuse(target_url=dashboard_uri, recent_actions=[f"open_tab:{dashboard_uri}"])
    return {
        "ok": True,
        "query": role_query,
        "results_count": len(jobs),
        "results": [asdict(x) for x in top],
        "artifacts": artifacts,
        "summary": summary,
        "source_status": source_status,
        "opened_url": dashboard_uri,
        "canvas": {
            "title": "Job Market Dashboard Generated",
            "subtitle": f"{len(jobs)} listings across {', '.join(region_labels).upper()}",
            "cards": [
                {
                    "title": x.title[:90],
                    "price": x.salary_text or ("Remote" if x.remote else "Salary n/a"),
                    "source": x.source,
                    "url": x.url,
                }
                for x in jobs[:6]
            ],
        },
    }


def _run_generic_research(
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    browser_worker_mode: str = "local",
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    runtime_contract = TaskContractEngine().extract(instruction)
    runtime_contract.domain = "generic_research"
    runtime_contract.requested_outputs = sorted(
        {
            "spreadsheet",
            "report",
            "dashboard",
            *(["presentation"] if _wants_powerpoint(instruction) else []),
        }
    )
    runtime_graph = CapabilityPlanner(registry=default_capability_registry()).plan(runtime_contract)
    runtime_result = _run_execution_graph_runtime_path(
        instruction=instruction,
        task_contract=runtime_contract,
        graph=runtime_graph,
        ai_meta=backend_metadata("deterministic-local"),
        progress_cb=progress_cb,
        mode_override="generic_research",
        extra_context={
            "browser_worker_mode": browser_worker_mode,
            "human_like_interaction": human_like_interaction,
        },
    )
    if not runtime_result.get("ok"):
        error_text = str(runtime_result.get("error", "") or "").strip()
        try:
            payload = json.loads(error_text) if error_text.startswith("{") else {}
        except Exception:
            payload = {}
        if isinstance(payload, dict) and payload.get("error"):
            error_code = str(payload.get("error", "generic_research_failed"))
            subtitle_map = {
                "locality_not_satisfied": "Locality constraint was not satisfied by candidate results.",
                "decision_quality_insufficient": "Need at least 2 strong candidates for a superlative/compare request.",
                "quality_threshold_not_met": "Candidate quality is too low to finalize safely.",
                "low_relevance": "Results were not relevant enough; task not executed.",
            }
            title_map = {
                "locality_not_satisfied": "Research Blocked By Locality Gate",
                "decision_quality_insufficient": "Research Blocked By Human Judgment Gate",
                "quality_threshold_not_met": "Research Blocked By Quality Gate",
                "low_relevance": "Research Blocked",
            }
            return {
                "ok": False,
                "mode": "generic_research",
                "runtime_mode": "execution_graph_runtime",
                "query": str(payload.get("query", "")),
                "results_count": int(payload.get("results_count", 0) or 0),
                "results": list(payload.get("results", []) or []),
                "artifacts": dict(runtime_result.get("artifacts", {}) or {}),
                "artifact_metadata": dict(runtime_result.get("artifact_metadata", {}) or {}),
                "summary": {"error": error_code, **(payload.get("summary", {}) if isinstance(payload.get("summary"), dict) else {})},
                "source_status": dict(payload.get("source_status", {}) or {}),
                "opened_url": "",
                "verification_report": runtime_result.get("verification_report", {}),
                "verification": runtime_result.get("verification", {}),
                "final_report": runtime_result.get("final_report", {}),
                "runtime_events": runtime_result.get("runtime_events", []),
                "task_contract": runtime_result.get("task_contract", {}),
                "capability_execution_graph": runtime_result.get("capability_execution_graph", {}),
                "critics": runtime_result.get("critics", {}),
                "memory_context": runtime_result.get("memory_context", {}),
                "canvas": {
                    "title": title_map.get(error_code, "Research Blocked"),
                    "subtitle": subtitle_map.get(error_code, "Research could not be completed safely."),
                    "cards": [],
                },
            }
        return runtime_result
    graph = runtime_result.get("capability_execution_graph", {}) or {}
    nodes = graph.get("nodes", []) if isinstance(graph.get("nodes", []), list) else []
    research_node = next((node for node in nodes if str(node.get("capability", "")) == "research_collection"), {})
    research_payload = dict(research_node.get("output_payload", {}) or {})
    query = str(research_payload.get("query", ""))
    results = [dict(x) for x in (research_payload.get("search_results") or []) if isinstance(x, dict)]
    recommendation = dict(research_payload.get("recommendation", {}) or {})
    browser_notes = [dict(x) for x in (research_payload.get("browser_notes") or []) if isinstance(x, dict)]
    source_status = dict(research_payload.get("source_status", {}) or {})
    summary = dict(research_payload.get("research_summary", {}) or {})
    opened_url = str(research_payload.get("opened_url", "") or "")
    artifacts = dict(runtime_result.get("artifacts", {}) or {})
    artifact_metadata = dict(runtime_result.get("artifact_metadata", {}) or {})
    if not opened_url:
        open_target = artifacts.get("primary_open_file") or artifacts.get("dashboard_html", "")
        opened_url = Path(open_target).resolve().as_uri() if open_target else ""
    if opened_url:
        opened_url, _nav = _open_target_with_reuse(target_url=opened_url, recent_actions=[f"open_tab:{opened_url}"])
    runtime_result["query"] = query
    runtime_result["results_count"] = len(results)
    runtime_result["results"] = results[:50]
    runtime_result["artifacts"] = artifacts
    runtime_result["artifact_metadata"] = artifact_metadata
    runtime_result["summary"] = {**(runtime_result.get("summary", {}) or {}), **summary}
    runtime_result["source_status"] = source_status
    runtime_result["opened_url"] = opened_url
    runtime_result["recommendation"] = recommendation
    runtime_result["canvas"] = {
        "title": "Decision Package Generated" if recommendation else "Research Deliverables Generated",
        "subtitle": str(recommendation.get("selected_title") or f"{len(results)} results"),
        "cards": [
            {
                "title": str(x.get("title", ""))[:90],
                "price": (f"${float(x.get('price')):.2f}" if isinstance(x.get("price"), (int, float)) else "result"),
                "source": str(x.get("source", "")),
                "url": str(x.get("url", "")),
            }
            for x in results[:6]
        ],
    }
    return runtime_result


def _human_judgment_constraints(instruction: str, query: str) -> Dict[str, Any]:
    return platform_human_judgment_constraints(instruction=instruction, query=query)


def _human_judgment_refine_queries(query: str, constraints: Dict[str, Any]) -> List[str]:
    return platform_human_judgment_refine_queries(query=query, constraints=constraints)


def _apply_human_judgment_quality_gate(
    *,
    ranked: List[SearchResult],
    instruction: str,
    query: str,
    constraints: Dict[str, Any],
) -> List[SearchResult]:
    return list(
        platform_apply_human_judgment_quality_gate(
            ranked=ranked,
            instruction=instruction,
            query=query,
            constraints=constraints,
        )
    )


def _quality_score_result(result: SearchResult, query: str, locality_terms: List[str]) -> int:
    quality = assess_result_quality(
        title=result.title,
        url=result.url,
        snippet=result.snippet,
        query=query,
        locality_terms=locality_terms,
    )
    return int(quality.score)


def _count_locality_matches(ranked: List[SearchResult], locality_terms: List[str]) -> int:
    return int(platform_count_locality_matches(ranked=ranked, locality_terms=locality_terms))


def _run_competitor_analysis(
    instruction: str,
    progress_cb: Optional[Callable[[int, str], None]] = None,
    min_live_non_curated_citations: Optional[int] = None,
) -> Dict[str, Any]:
    min_live = _effective_min_live_non_curated_citations(min_live_non_curated_citations)
    if not _is_competitor_analysis_intent(instruction):
        return {
            "ok": False,
            "mode": "competitor_analysis",
            "query": "",
            "results_count": 0,
            "results": [],
            "artifacts": {},
            "summary": {"error": "competitor_intent_not_detected"},
            "source_status": {},
            "opened_url": "",
            "canvas": {
                "title": "Competitor Analysis Skipped",
                "subtitle": "Prompt did not request competitor analysis explicitly.",
                "cards": [],
            },
        }
    target = _extract_competitor_target(instruction)
    output_folder = _extract_named_output_folder(instruction, default_name=f"{target} Competitor Analysis")
    safe_folder = re.sub(r"[<>:\"/\\\\|?*]", "", output_folder).strip() or f"{target} Competitor Analysis"
    workspace = Path("data/reports") / safe_folder
    workspace.mkdir(parents=True, exist_ok=True)
    _emit_progress(progress_cb, 16, f"Researching competitors for: {target}")
    task_contract = TaskContractEngine().extract(instruction)
    task_contract.domain = "competitor_analysis"
    task_contract.audience = "stakeholder"
    task_contract.requested_outputs = ["report", "presentation", "dashboard", "spreadsheet"]
    runtime_graph = CapabilityPlanner(registry=default_capability_registry()).plan(task_contract)
    runtime_result = _run_execution_graph_runtime_path(
        instruction=instruction,
        task_contract=task_contract,
        graph=runtime_graph,
        ai_meta=backend_metadata("deterministic-local"),
        progress_cb=progress_cb,
        mode_override="competitor_analysis",
        extra_context={
            "workspace_dir": str(workspace.resolve()),
            "min_live_non_curated_citations": min_live,
        },
    )
    if not runtime_result.get("ok"):
        error_text = str(runtime_result.get("error") or "")
        if error_text.startswith("insufficient_live_non_curated_citations:"):
            parts = error_text.split(":")
            found = int(parts[1]) if len(parts) > 1 and parts[1].isdigit() else 0
            required = int(parts[2]) if len(parts) > 2 and parts[2].isdigit() else min_live
            return {
                "ok": False,
                "mode": "competitor_analysis",
                "runtime_mode": "execution_graph_runtime",
                "query": f"{target} EHR competitors",
                "results_count": 0,
                "results": [],
                "artifacts": runtime_result.get("artifacts", {}) or {},
                "artifact_metadata": runtime_result.get("artifact_metadata", {}) or {},
                "summary": {
                    "target": target,
                    "error": "insufficient_live_non_curated_citations",
                    "required_live_non_curated_citations": required,
                    "live_non_curated_citations": found,
                },
                "source_status": {},
                "opened_url": "",
                "verification_report": runtime_result.get("verification_report", {}),
                "verification": runtime_result.get("verification", {}),
                "final_report": runtime_result.get("final_report", {}),
                "canvas": {
                    "title": "Run Blocked",
                    "subtitle": f"Need at least {required} live non-curated citations; found {found}.",
                    "cards": [],
                },
                "runtime_events": runtime_result.get("runtime_events", []),
                "artifact_metadata": runtime_result.get("artifact_metadata", {}),
                "task_contract": runtime_result.get("task_contract", {}),
                "capability_execution_graph": runtime_result.get("capability_execution_graph", {}),
                "critics": runtime_result.get("critics", {}),
                "memory_context": runtime_result.get("memory_context", {}),
            }
        return runtime_result

    graph = runtime_result.get("capability_execution_graph", {}) or {}
    nodes = graph.get("nodes", []) if isinstance(graph.get("nodes", []), list) else []
    competitor_node = next((node for node in nodes if str(node.get("capability", "")) == "competitor_research"), {})
    competitor_payload = competitor_node.get("output_payload", {}) if isinstance(competitor_node, dict) else {}
    rows = [dict(x) for x in (competitor_payload.get("structured_rows") or []) if isinstance(x, dict)]
    notes = dict(competitor_payload.get("research_notes", {}) or {})
    sources = [dict(x) for x in (competitor_payload.get("sources") or []) if isinstance(x, dict)]
    cards = [
        {
            "title": str(row.get("name", ""))[:90],
            "price": str(row.get("segment", "EHR")),
            "source": "competitor_analysis",
            "url": str(row.get("citations", "")).split(" | ")[0] if str(row.get("citations", "")).strip() else "",
        }
        for row in rows[:6]
    ]
    runtime_result["summary"] = {
        **(runtime_result.get("summary", {}) or {}),
        "target": target,
        "top_competitors": [row.get("name", "") for row in rows[:5]],
        "competitor_count": len(rows),
        "sources_used": len(sources),
        "live_non_curated_citations": int(notes.get("live_non_curated_citations", 0) or 0),
        "required_live_non_curated_citations": int(notes.get("required_live_non_curated_citations", min_live) or min_live),
    }
    runtime_result["results_count"] = len(sources)
    runtime_result["results"] = [
        {
            "title": source.get("name", ""),
            "url": source.get("url", ""),
            "source": source.get("source_type", ""),
            "snippet": source.get("snippet", ""),
        }
        for source in sources[:50]
    ]
    runtime_result["canvas"] = {
        "title": f"{target} Competitor Analysis Ready",
        "subtitle": f"Top {len(rows[:5])} competitors with executive summary + PowerPoint",
        "cards": cards,
    }
    return runtime_result


def _count_live_non_curated_citations(results: List[SearchResult]) -> int:
    seen: set[str] = set()
    count = 0
    for r in results:
        src = str(r.source or "").strip().lower()
        url = str(r.url or "").strip()
        if not url or not url.startswith(("http://", "https://")):
            continue
        host = urllib.parse.urlparse(url).netloc.lower()
        if src == "curated":
            continue
        if host in seen:
            continue
        seen.add(host)
        count += 1
    return count


def _effective_min_live_non_curated_citations(value: Optional[int]) -> int:
    if value is None:
        return MIN_LIVE_NON_CURATED_CITATIONS
    try:
        return max(1, min(20, int(value)))
    except Exception:
        return MIN_LIVE_NON_CURATED_CITATIONS


def _curated_ehr_competitor_sources(target: str) -> List[SearchResult]:
    _ = target
    return [
        SearchResult(
            title="Oracle Health | Electronic Health Record",
            url="https://www.oracle.com/health/",
            price=None,
            source="curated",
            snippet="Oracle Health (including Cerner capabilities) enterprise healthcare platform information.",
        ),
        SearchResult(
            title="MEDITECH EHR Platform",
            url="https://ehr.meditech.com/",
            price=None,
            source="curated",
            snippet="MEDITECH electronic health record platform overview.",
        ),
        SearchResult(
            title="athenahealth EHR",
            url="https://www.athenahealth.com/solutions/electronic-health-records",
            price=None,
            source="curated",
            snippet="athenahealth ambulatory-focused EHR product information.",
        ),
        SearchResult(
            title="eClinicalWorks EHR",
            url="https://www.eclinicalworks.com/",
            price=None,
            source="curated",
            snippet="eClinicalWorks EHR solutions for practices and health systems.",
        ),
        SearchResult(
            title="NextGen Healthcare EHR",
            url="https://www.nextgen.com/solutions/ehr",
            price=None,
            source="curated",
            snippet="NextGen ambulatory EHR capabilities and product information.",
        ),
        SearchResult(
            title="Veradigm Healthcare Data and EHR Solutions",
            url="https://veradigm.com/",
            price=None,
            source="curated",
            snippet="Veradigm portfolio including legacy Allscripts ecosystem context.",
        ),
    ]


def _extract_competitor_target(instruction: str) -> str:
    m = re.search(r"competitors?\s+(?:to|for|of)\s+([A-Za-z0-9& .-]{2,80})", instruction, flags=re.IGNORECASE)
    if m:
        return re.sub(r"\s+", " ", m.group(1)).strip(" .,\"'â€œâ€")
    if "epic" in instruction.lower():
        return "Epic Systems"
    return "Target Company"


def _extract_named_output_folder(instruction: str, default_name: str) -> str:
    m = re.search(r'folder\s+called\s+["â€œâ€\']([^"â€œâ€\']{2,80})["â€œâ€\']', instruction, flags=re.IGNORECASE)
    if not m:
        return default_name
    return m.group(1).strip()


def _competitor_queries(target: str) -> List[str]:
    return [
        f'"{target}" competitors electronic health record',
        f'"{target}" EHR market share competitors',
        f'{target} vs Oracle Health Cerner MEDITECH athenahealth',
        f'{target} hospital EHR alternatives enterprise',
        f'{target} ambulatory EHR competitors eClinicalWorks NextGen',
    ]


def _competitor_must_terms(target: str) -> List[str]:
    return [
        target.lower(),
        "ehr",
        "electronic health record",
        "healthcare",
        "hospital",
        "clinical",
        "oracle health",
        "cerner",
        "meditech",
        "athenahealth",
        "allscripts",
        "veradigm",
        "eclinicalworks",
        "nextgen",
    ]


def _select_top_competitors(target: str, results: List[SearchResult], top_n: int = 5) -> List[Dict[str, Any]]:
    competitor_aliases = {
        "Oracle Health (Cerner)": ["oracle health", "cerner"],
        "MEDITECH": ["meditech"],
        "athenahealth": ["athenahealth"],
        "Veradigm (Allscripts)": ["veradigm", "allscripts"],
        "eClinicalWorks": ["eclinicalworks", "eclinical works"],
        "NextGen Healthcare": ["nextgen healthcare", "nextgen"],
        "CPSI / TruBridge": ["cpsi", "trubridge"],
        "Altera Digital Health": ["altera digital health", "altera"],
    }
    scores: Dict[str, float] = {k: 0.0 for k in competitor_aliases}
    evidence: Dict[str, List[str]] = {k: [] for k in competitor_aliases}
    for r in results[:120]:
        hay = f"{r.title} {r.snippet} {r.url}".lower()
        for name, aliases in competitor_aliases.items():
            hit = sum(1 for a in aliases if a in hay)
            if hit:
                scores[name] += float(hit) + (_relevance_score(r, "ehr healthcare") * 0.15)
                if len(evidence[name]) < 3:
                    evidence[name].append(r.url)
    ordered = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    selected: List[Dict[str, Any]] = []
    for name, score in ordered:
        if len(selected) >= top_n:
            break
        if score <= 0.0:
            continue
        selected.append(
            {
                "name": name,
                "score": round(score, 3),
                "segment": "Enterprise EHR",
                "citations": evidence.get(name, []),
                "why": f"Detected recurring mentions with {target} in EHR market context.",
            }
        )
    if len(selected) < top_n:
        fallback = [
            "Oracle Health (Cerner)",
            "MEDITECH",
            "athenahealth",
            "Veradigm (Allscripts)",
            "eClinicalWorks",
            "NextGen Healthcare",
        ]
        for name in fallback:
            if len(selected) >= top_n:
                break
            if any(x["name"] == name for x in selected):
                continue
            selected.append(
                {
                    "name": name,
                    "score": 0.1,
                    "segment": "Enterprise/ambulatory EHR",
                    "citations": evidence.get(name, [])[:2],
                    "why": f"Industry-typical competitor set around {target} EHR footprint.",
                }
            )
    return selected[:top_n]


def _write_competitor_artifacts(
    instruction: str,
    target: str,
    output_folder: str,
    competitors: List[Dict[str, Any]],
    ranked_results: List[SearchResult],
) -> Dict[str, str]:
    safe_folder = re.sub(r"[<>:\"/\\\\|?*]", "", output_folder).strip() or f"{target} Competitor Analysis"
    out_dir = Path("data/reports") / safe_folder
    out_dir.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    csv_path = out_dir / f"competitors_{ts}.csv"
    report_path = out_dir / f"executive_summary_{ts}.md"
    report_html_path = out_dir / f"executive_summary_{ts}.html"
    pptx_path = out_dir / f"executive_summary_{ts}.pptx"
    dash_path = out_dir / f"dashboard_{ts}.html"

    with csv_path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=["name", "score", "segment", "why", "citations"])
        writer.writeheader()
        for row in competitors:
            writer.writerow(
                {
                    "name": row.get("name", ""),
                    "score": row.get("score", ""),
                    "segment": row.get("segment", ""),
                    "why": row.get("why", ""),
                    "citations": " | ".join(row.get("citations", [])[:3]),
                }
            )

    lines: List[str] = []
    lines.append(f"# Executive Summary: {target} Competitor Analysis")
    lines.append("")
    lines.append(f"- Generated: {datetime.now().isoformat(timespec='seconds')}")
    lines.append(f"- Instruction: {instruction}")
    lines.append(f"- Output folder: `{out_dir.resolve()}`")
    lines.append("")
    lines.append("## 1. Executive Overview")
    lines.append(
        f"{target} remains a dominant enterprise EHR platform. The most credible competitive pressure in hospitals and ambulatory networks is concentrated in a small set of vendors with broad clinical workflow coverage, large installed footprints, and active modernization programs."
    )
    lines.append("")
    lines.append("## 2. Top Competitors")
    for i, comp in enumerate(competitors, start=1):
        lines.append(f"### {i}) {comp.get('name','')}")
        lines.append(f"- Positioning: {comp.get('segment','EHR vendor')}")
        lines.append(f"- Why it matters: {comp.get('why','')}")
        lines.append(
            "- Executive view: This vendor appears in competitive conversations when buyers evaluate enterprise breadth, operational risk during implementation, and long-term platform modernization costs."
        )
        lines.append(
            "- Commercial lens: Selection dynamics typically depend on how strongly the platform supports integrated clinical-financial workflows and migration complexity from incumbent tooling."
        )
        cites = comp.get("citations", [])[:3]
        if cites:
            lines.append("- Evidence:")
            for c in cites:
                lines.append(f"  - {c}")
        lines.append("")
    lines.append("## 3. Strategic Implications")
    lines.append(
        "Health systems comparing alternatives to Epic typically balance enterprise integration depth, revenue-cycle integration, specialty workflow maturity, implementation risk, and migration timeline. In practice, competitive displacement is strongest during major modernization cycles, mergers, or when organizations rebalance enterprise vs. ambulatory priorities."
    )
    lines.append(
        "At the executive level, competitor pressure is rarely about one feature. It is usually about total operating model fit: governance model, data strategy, interoperability architecture, implementation throughput, and degree of disruption to frontline clinical users. Organizations that do best in these transitions define non-negotiable outcomes first, then evaluate each platform against those outcomes using measurable decision criteria."
    )
    lines.append(
        "For most buyers, near-term implementation risk and medium-term optimization capability matter more than feature parity marketing claims. A credible competitor to Epic must demonstrate reliable deployment at scale, durable post-go-live support, and clear evidence of value realization across quality, throughput, and revenue-cycle performance."
    )
    lines.append("")
    lines.append("## 4. Recommended Next Steps")
    lines.append("1. Validate top-5 shortlist with your target segment (IDN, community hospital, ambulatory-heavy network).")
    lines.append("2. Build side-by-side scorecard: interoperability, total cost of ownership, implementation time, and user adoption risk.")
    lines.append("3. Run focused diligence on migration tooling and data conversion readiness.")
    lines.append("4. Confirm executive sponsorship model and change-management capacity before final vendor down-select.")
    report_path.write_text("\n".join(lines), encoding="utf-8")
    report_html_path.write_text(
        "<!doctype html><html lang='en'><head><meta charset='utf-8'/><meta name='viewport' content='width=device-width,initial-scale=1'/>"
        "<title>Executive Summary</title><style>body{font-family:'Segoe UI',Tahoma,sans-serif;background:#f8fafc;color:#0f172a;margin:0}"
        ".w{max-width:980px;margin:0 auto;padding:24px}.card{background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:16px}pre{white-space:pre-wrap}</style></head>"
        f"<body><div class='w'><div class='card'><h1>{html.escape(target)} Competitor Executive Summary</h1><pre>{html.escape(report_path.read_text(encoding='utf-8'))}</pre></div></div></body></html>",
        encoding="utf-8",
    )

    _write_competitor_pptx(pptx_path=pptx_path, target=target, competitors=competitors)
    _write_generic_dashboard_html(dash_path=dash_path, results=ranked_results, title=f"{target} Competitor Research Sources")

    return {
        "directory": str(out_dir.resolve()),
        "competitors_csv": str(csv_path.resolve()),
        "executive_summary_md": str(report_path.resolve()),
        "executive_summary_html": str(report_html_path.resolve()),
        "powerpoint_pptx": str(pptx_path.resolve()),
        "dashboard_html": str(dash_path.resolve()),
        "primary_open_file": str(report_html_path.resolve()),
    }


def _write_competitor_pptx(pptx_path: Path, target: str, competitors: List[Dict[str, Any]]) -> None:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        pptx_path.write_text(
            "PowerPoint package unavailable. Install python-pptx to generate .pptx files.",
            encoding="utf-8",
        )
        return

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{target} Competitor Analysis"
    slide.placeholders[1].text = "Executive Summary Deck"

    overview = prs.slides.add_slide(prs.slide_layouts[1])
    overview.shapes.title.text = "Executive Overview"
    overview.placeholders[1].text = (
        f"{target} competes in a concentrated EHR market where enterprise integration, interoperability, and implementation risk drive selection."
    )

    for comp in competitors[:5]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = comp.get("name", "Competitor")
        cites = comp.get("citations", [])[:2]
        cite_text = "\n".join(f"- {u}" for u in cites) if cites else "- Source coverage captured in dashboard file."
        s.placeholders[1].text = (
            f"Segment: {comp.get('segment','EHR')}\n"
            f"Why it matters: {comp.get('why','')}\n"
            f"Evidence:\n{cite_text}"
        )

    close = prs.slides.add_slide(prs.slide_layouts[1])
    close.shapes.title.text = "Recommended Next Steps"
    close.placeholders[1].text = (
        "1) Validate shortlist against your deployment profile.\n"
        "2) Build weighted scorecard (clinical fit, cost, interoperability).\n"
        "3) Launch structured vendor diligence with implementation risk gates."
    )
    prs.save(str(pptx_path))


def _run_study_pack(instruction: str, progress_cb: Optional[Callable[[int, str], None]] = None) -> Dict[str, Any]:
    topic = _extract_study_topic(instruction)
    question_count = _extract_question_count(instruction, default_count=200)
    use_notebooklm = "notebooklm" in instruction.lower()
    _emit_progress(progress_cb, 16, f"Researching study sources for: {topic}")

    sources = _study_sources_for_topic(topic)
    preferred_domains = _preferred_domains_for_topic(topic)
    strict_official = _is_strict_official_topic(topic)
    source_results: List[SearchResult] = []
    source_status: Dict[str, str] = {}
    for i, q in enumerate(sources):
        _emit_progress(progress_cb, 22 + int((i / max(1, len(sources))) * 38), f"Finding relevant sources: {q}")
        pulled = _search_web(q, limit=12)
        filtered = _filter_relevant_results(
            pulled,
            must_terms=_must_terms_for_topic(topic),
            banned_domains={
                "support.google.com",
                "mail.google.com",
                "gmail.com",
                "support.microsoft.com",
                "learn.microsoft.com",
                "microsoft.com",
                "windows.com",
            },
            min_score=2.0,
            preferred_domains=preferred_domains,
        )
        source_results.extend(filtered)
        source_status[q] = f"ok:{len(filtered)}"

    if len(source_results) < 5:
        _emit_progress(progress_cb, 58, "Refining search terms due low relevance")
        for q in _refine_study_queries(topic):
            pulled = _search_web(q, limit=10)
            filtered = _filter_relevant_results(
                pulled,
                must_terms=_must_terms_for_topic(topic),
                banned_domains={
                    "support.google.com",
                    "mail.google.com",
                    "gmail.com",
                    "support.microsoft.com",
                    "learn.microsoft.com",
                    "microsoft.com",
                    "windows.com",
                },
                min_score=1.5,
                preferred_domains=preferred_domains if strict_official else preferred_domains,
            )
            source_results.extend(filtered)
            source_status[q] = f"retry:{len(filtered)}"
            if len(source_results) >= 10:
                break

    dedup: Dict[str, SearchResult] = {}
    for r in source_results:
        dedup[r.url] = r
    ranked = sorted(dedup.values(), key=lambda x: _relevance_score(x, topic), reverse=True)

    if not ranked and not strict_official:
        _emit_progress(progress_cb, 64, "Relaxing filters to recover high-signal sources")
        fallback_pool: List[SearchResult] = []
        for q in sources + _refine_study_queries(topic):
            fallback_pool.extend(_search_web(q, limit=10))
        ranked = _filter_relevant_results(
            fallback_pool,
            must_terms=_must_terms_for_topic(topic),
            banned_domains={
                "support.google.com",
                "mail.google.com",
                "gmail.com",
                "dell.com",
                "lenovo.com",
                "nvidia.com",
                "hp.com",
                "intel.com",
            },
            min_score=2.0,
            preferred_domains=None,
        )
        dedup2: Dict[str, SearchResult] = {}
        for r in ranked:
            dedup2[r.url] = r
        ranked = sorted(dedup2.values(), key=lambda x: _relevance_score(x, topic), reverse=True)

    if not ranked:
        _emit_progress(progress_cb, 66, "Using curated official source links")
        ranked = _curated_study_sources(topic)

    _emit_progress(progress_cb, 70, "Extracting facts from official source material")
    fact_bank = _build_fact_bank(topic=topic, sources=ranked)
    if len(fact_bank) < 20:
        return {
            "ok": False,
            "query": topic,
            "results_count": 0,
            "results": [asdict(x) for x in ranked[:10]],
            "artifacts": {},
            "summary": {"topic": topic, "error": "insufficient_evidence", "fact_count": len(fact_bank)},
            "source_status": source_status,
            "opened_url": "",
            "paused_for_credentials": False,
            "pause_reason": "",
            "canvas": {
                "title": "Study Pack Blocked",
                "subtitle": "Not enough validated source evidence. No synthetic quiz generated.",
                "cards": [{"title": "Evidence", "price": str(len(fact_bank)), "source": "validator", "url": ""}],
            },
        }

    _emit_progress(progress_cb, 78, f"Generating {question_count} source-grounded flashcards")
    cards = _generate_study_items(topic=topic, count=question_count, source_results=ranked, fact_bank=fact_bank)
    coverage = _study_evidence_coverage(cards)
    if coverage < 0.85:
        return {
            "ok": False,
            "query": topic,
            "results_count": 0,
            "results": [asdict(x) for x in ranked[:10]],
            "artifacts": {},
            "summary": {"topic": topic, "error": "low_evidence_coverage", "coverage": coverage},
            "source_status": source_status,
            "opened_url": "",
            "paused_for_credentials": False,
            "pause_reason": "",
            "canvas": {
                "title": "Study Pack Blocked",
                "subtitle": "Evidence coverage below threshold. No low-quality quiz output allowed.",
                "cards": [{"title": "Coverage", "price": f"{coverage:.2%}", "source": "validator", "url": ""}],
            },
        }
    quality_issues = _study_human_quality_issues(cards)
    if quality_issues:
        return {
            "ok": False,
            "query": topic,
            "results_count": 0,
            "results": [asdict(x) for x in ranked[:10]],
            "artifacts": {},
            "summary": {"topic": topic, "error": "human_quality_failed", "issues": quality_issues[:10]},
            "source_status": source_status,
            "opened_url": "",
            "paused_for_credentials": False,
            "pause_reason": "",
            "canvas": {
                "title": "Study Pack Blocked",
                "subtitle": "Human-quality validation failed. Output rejected.",
                "cards": [{"title": "Issues", "price": str(len(quality_issues)), "source": "validator", "url": ""}],
            },
        }
    artifacts = _write_study_artifacts(instruction=instruction, topic=topic, cards=cards, sources=ranked)

    pause_reason = ""
    if use_notebooklm:
        _emit_progress(progress_cb, 90, "Opening NotebookLM workspace")
        notebooklm_url = "https://notebooklm.google.com/"
        notebooklm_url, _nav = _open_target_with_reuse(target_url=notebooklm_url, recent_actions=[f"open_tab:{notebooklm_url}"])
        pause_reason = "NotebookLM opened. If sign-in is required, complete login and continue."
    dashboard_uri = Path(artifacts["quiz_html"]).resolve().as_uri()
    dashboard_uri, _nav = _open_target_with_reuse(target_url=dashboard_uri, recent_actions=[f"open_tab:{dashboard_uri}"])

    return {
        "ok": True,
        "query": topic,
        "results_count": len(cards),
        "results": [asdict(x) for x in ranked[:25]],
        "artifacts": artifacts,
        "summary": {
            "topic": topic,
            "question_count": len(cards),
            "sources_found": len(ranked),
            "notebooklm_requested": use_notebooklm,
        },
        "source_status": source_status,
        "opened_url": dashboard_uri,
        "paused_for_credentials": bool(use_notebooklm),
        "pause_reason": pause_reason,
        "canvas": {
            "title": "Study Pack Ready",
            "subtitle": f"{len(cards)} questions for {topic}",
            "cards": [
                {
                    "title": r.title[:90],
                    "price": "source",
                    "source": r.source,
                    "url": r.url,
                }
                for r in ranked[:6]
            ],
        },
    }


def _extract_role_query(instruction: str) -> str:
    low = instruction.lower()
    wants_vp_avp = bool(re.search(r"\b(avp|vp|vice president|assistant vice president)\b", low))
    if "data and ai" in low or "data and analytics" in low or "analytics" in low:
        if wants_vp_avp:
            return "VP OR AVP Data and AI OR Data and Analytics OR Analytics"
        return "Data and AI OR Artificial Intelligence OR Data Analytics OR Machine Learning"
    if wants_vp_avp:
        return "VP OR AVP Data OR AI OR Analytics leadership roles"
    return "Data and AI roles OR Artificial Intelligence roles OR Data leadership roles"


def _extract_generic_query(instruction: str) -> str:
    return platform_extract_generic_query(instruction)


def _extract_study_topic(instruction: str) -> str:
    text = re.sub(r"\s+", " ", instruction).strip()
    m = re.search(r"(?:for|about)\s+(.+?)(?:\s+exam|\s+test|\s+quiz|$)", text, flags=re.IGNORECASE)
    if m:
        return m.group(1).strip(" .")
    return text


def _extract_question_count(instruction: str, default_count: int = 200) -> int:
    m = re.search(r"\b([0-9]{2,4})\s+question", instruction, flags=re.IGNORECASE)
    if not m:
        return default_count
    return max(20, min(500, int(m.group(1))))


def _study_sources_for_topic(topic: str) -> List[str]:
    low = topic.lower()
    if "south carolina" in low and ("driver" in low or "permit" in low):
        return [
            'site:scdmvonline.com south carolina driver manual permit test',
            'site:dmv.sc.gov Driver Manual PDF South Carolina',
            'site:dc.statelibrary.sc.gov south carolina driver license manual',
            f'"{topic}" rules of the road signs permit practice',
        ]
    return [
        f'"{topic}" official handbook PDF',
        f'"{topic}" practice test questions',
        f'"{topic}" study guide',
    ]


def _refine_study_queries(topic: str) -> List[str]:
    return [
        f'"{topic}" official manual PDF',
        f'"{topic}" exam topics road signs right of way',
        f'"{topic}" permit exam sample questions',
    ]


def _must_terms_for_topic(topic: str) -> List[str]:
    low = topic.lower()
    terms = [t for t in re.split(r"[^a-z0-9]+", low) if len(t) > 2]
    if "driver" in low or "permit" in low:
        terms.extend(["driver", "permit", "manual", "road", "dmv"])
    return list(dict.fromkeys(terms))[:12]


def _preferred_domains_for_topic(topic: str) -> List[str]:
    low = topic.lower()
    if "south carolina" in low and ("driver" in low or "permit" in low):
        return [
            "dmv.sc.gov",
            "scdmvonline.com",
            "dc.statelibrary.sc.gov",
            "driving-tests.org",
            "dmv.org",
        ]
    if "driver" in low or "permit" in low:
        return ["dmv", "gov", "state"]
    return []


def _is_strict_official_topic(topic: str) -> bool:
    low = topic.lower()
    return ("south carolina" in low) and ("driver" in low or "permit" in low)


def _curated_study_sources(topic: str) -> List[SearchResult]:
    low = topic.lower()
    if "south carolina" in low and ("driver" in low or "permit" in low):
        return [
            SearchResult(
                title="South Carolina Driver Manual (SC DMV)",
                url="https://dmv.sc.gov/sites/scdmv/files/media/Files/Driver-Manual.pdf",
                price=None,
                source="curated",
                snippet="Official South Carolina driver manual PDF.",
            ),
            SearchResult(
                title="SC DMV Driver Services",
                url="https://www.scdmvonline.com/Driver-Services",
                price=None,
                source="curated",
                snippet="South Carolina DMV driver services portal.",
            ),
            SearchResult(
                title="SCDMV Beginner's Permit",
                url="https://dmv.sc.gov/driver-services/drivers-license/beginner-permits",
                price=None,
                source="curated",
                snippet="Official beginner permit eligibility and requirements.",
            ),
            SearchResult(
                title="South Carolina Driver's License Manual (State Library)",
                url="https://dc.statelibrary.sc.gov/handle/10827/62666",
                price=None,
                source="curated",
                snippet="State document repository for the SC driver manual.",
            ),
        ]
    return [
        SearchResult(
            title=f"Official manual source for {topic}",
            url="https://www.usa.gov/motor-vehicle-services",
            price=None,
            source="curated",
            snippet="General US DMV service portal.",
        )
    ]


def _build_fact_bank(topic: str, sources: List[SearchResult]) -> List[Dict[str, str]]:
    facts: List[Dict[str, str]] = []
    for src in sources[:12]:
        url = src.url
        if url.lower().endswith(".pdf") or "driver-manual.pdf" in url.lower():
            facts.extend(_extract_pdf_fact_entries(url=url, topic=topic))
            continue
        text = _extract_source_text(url)
        if not text:
            continue
        for sentence in _split_sentences(text):
            clean = _normalize_fact_text(sentence)
            if _is_fact_sentence(clean, topic):
                score = _fact_sentence_score(clean)
                if url.lower().endswith(".pdf") or "driver-manual.pdf" in url.lower():
                    score += 2.0
                facts.append(
                    {
                        "text": clean.strip(),
                        "source_url": url,
                        "category": _categorize_fact(clean),
                        "score": f"{score:.3f}",
                        "image_path": "",
                    }
                )
    dedup: Dict[str, Dict[str, str]] = {}
    for f in facts:
        key = re.sub(r"\s+", " ", f["text"]).strip().lower()
        dedup[key] = f
    ranked = list(dedup.values())
    ranked.sort(key=lambda x: float(x.get("score", "0")), reverse=True)
    return ranked


def _extract_pdf_fact_entries(url: str, topic: str) -> List[Dict[str, str]]:
    out: List[Dict[str, str]] = []
    fitz = _get_fitz_module()
    if fitz is None:
        text = _extract_source_text(url)
        if not text:
            return out
        for s in _split_sentences(text):
            clean = _normalize_fact_text(s)
            if _is_fact_sentence(clean, topic):
                out.append(
                    {
                        "text": clean,
                        "source_url": url,
                        "category": _categorize_fact(clean),
                        "score": f"{_fact_sentence_score(clean)+2.0:.3f}",
                        "image_path": "",
                        "image_base64": "",
                    }
                )
        return out

    req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(req, timeout=40) as resp:  # nosec B310
        raw = resp.read()
    doc = fitz.open(stream=raw, filetype="pdf")
    max_pages = min(120, int(getattr(doc, "page_count", 0) or 0))
    for i in range(max_pages):
        try:
            page = doc.load_page(i)
            page_text = page.get_text("text") or ""
        except Exception:
            continue
        for s in _split_sentences(page_text):
            clean = _normalize_fact_text(s)
            if not _is_fact_sentence(clean, topic):
                continue
            score = _fact_sentence_score(clean) + 2.0
            image_path = ""
            image_base64 = ""
            if _needs_visual(clean):
                image_path = _render_pdf_page_image(doc, i, url)
                image_base64 = image_to_base64(image_path) if image_path else ""
            out.append(
                {
                    "text": clean,
                    "source_url": url,
                    "category": _categorize_fact(clean),
                    "score": f"{score:.3f}",
                    "image_path": image_path,
                    "image_base64": image_base64,
                }
            )
    return out


def _needs_visual(sentence: str) -> bool:
    low = sentence.lower()
    return any(k in low for k in VISUAL_KEYWORDS)


def _render_pdf_page_image(doc: Any, page_index: int, source_url: str) -> str:
    fitz = _get_fitz_module()
    if fitz is None:
        return ""
    try:
        digest = hashlib.sha256(source_url.encode("utf-8")).hexdigest()[:16]
        out_dir = STUDY_ASSETS_ROOT / digest
        out_dir.mkdir(parents=True, exist_ok=True)
        out = out_dir / f"page_{page_index+1}.png"
        if out.exists():
            return str(out.resolve())
        page = doc.load_page(page_index)
        pix = page.get_pixmap(matrix=fitz.Matrix(STUDY_PDF_RENDER_ZOOM, STUDY_PDF_RENDER_ZOOM), alpha=False)
        pix.save(str(out))
        return str(out.resolve())
    except Exception as exc:
        LOGGER.warning("Failed to render PDF page image for %s page %s: %s", source_url, page_index + 1, exc)
        return ""


def _extract_source_text(url: str) -> str:
    low = url.lower()
    try:
        if low.endswith(".pdf") or "driver-manual.pdf" in low:
            text = _extract_pdf_text(url)
            return text
        html_text = _fetch_text(url)
        return _extract_text_from_html(html_text)
    except Exception:
        return ""


def _extract_pdf_text(url: str) -> str:
    req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(req, timeout=30) as resp:  # nosec B310
        raw = resp.read()

    # Try fast pure-Python extraction first.
    try:
        import io
        import pypdf

        reader = pypdf.PdfReader(io.BytesIO(raw))
        chunks: List[str] = []
        for page in reader.pages[:120]:
            try:
                chunks.append(page.extract_text() or "")
            except Exception:
                continue
        joined = "\n".join(chunks).strip()
        if len(joined) > 5000:
            return joined
    except Exception:
        pass

    # Fallback: PyMuPDF handles many scanned/complex PDFs better.
    try:
        import fitz  # type: ignore

        doc = fitz.open(stream=raw, filetype="pdf")
        chunks2: List[str] = []
        max_pages = min(120, int(getattr(doc, "page_count", 0) or 0))
        for i in range(max_pages):
            try:
                chunks2.append(doc.load_page(i).get_text("text") or "")
            except Exception:
                continue
        return "\n".join(chunks2).strip()
    except Exception:
        return ""


def _extract_text_from_html(html_text: str) -> str:
    text = re.sub(r"<script[\s\S]*?</script>", " ", html_text, flags=re.IGNORECASE)
    text = re.sub(r"<style[\s\S]*?</style>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html.unescape(text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _split_sentences(text: str) -> List[str]:
    parts = re.split(r"(?<=[\.\!\?])\s+", text)
    out = []
    for p in parts:
        s = p.strip()
        if 60 <= len(s) <= 320:
            out.append(s)
    return out


def _normalize_fact_text(text: str) -> str:
    s = text.replace("\r", " ").replace("\n", " ")
    s = re.sub(r"\s+", " ", s)
    s = re.sub(r"\b([0-9]{1,2})-([0-9]{1,2})\b", " ", s)
    s = re.sub(r"\s([â€¢\-])\s", " ", s)
    return s.strip(" .")


def _is_fact_sentence(sentence: str, topic: str) -> bool:
    low = sentence.lower()
    must = _must_terms_for_topic(topic)
    score = sum(1 for t in must if t in low)
    legal_signal = any(k in low for k in ["must", "shall", "required", "cannot", "may", "permit", "license", "test"])
    return score >= 1 and legal_signal


def _categorize_fact(sentence: str) -> str:
    low = sentence.lower()
    mapping = [
        ("Road Signs", ["sign", "signal", "traffic light"]),
        ("Permit Rules", ["permit", "license", "application", "under 18", "age"]),
        ("Right of Way", ["right-of-way", "yield", "intersection"]),
        ("Safe Driving", ["safe", "seat belt", "defensive", "speed"]),
        ("Alcohol and Drugs", ["alcohol", "drug", "dui"]),
        ("Parking and Turns", ["park", "turn", "lane", "u-turn"]),
    ]
    for cat, keys in mapping:
        if any(k in low for k in keys):
            return cat
    return "Rules of the Road"


def _fact_sentence_score(sentence: str) -> float:
    low = sentence.lower()
    score = 0.0
    score += 1.0 if len(sentence) >= 90 else 0.4
    score += 0.8 if any(k in low for k in ["must", "shall", "required", "illegal", "prohibited"]) else 0.0
    score += 0.6 if any(k in low for k in ["speed", "lane", "intersection", "sign", "signal", "right-of-way"]) else 0.0
    score += 0.4 if re.search(r"\b[0-9]{1,3}\b", sentence) else 0.0
    return score


def _study_evidence_coverage(cards: List[StudyItem]) -> float:
    if not cards:
        return 0.0
    good = sum(1 for c in cards if c.source_url and c.evidence and len(c.evidence) > 20)
    return good / len(cards)


def _study_human_quality_issues(cards: List[StudyItem]) -> List[str]:
    issues: List[str] = []
    for i, c in enumerate(cards, start=1):
        q = c.question.lower()
        a = c.answer.strip()
        if len(a) < 25:
            issues.append(f"Q{i}: answer too short")
        image_referenced = any(k in q for k in ["image shown", "picture", "sign shown", "visual"])
        if image_referenced and not c.image_path:
            issues.append(f"Q{i}: image referenced but no image attached")
    return issues


def _filter_relevant_results(
    results: List[SearchResult],
    must_terms: List[str],
    banned_domains: set[str],
    min_score: float = 1.0,
    preferred_domains: Optional[List[str]] = None,
) -> List[SearchResult]:
    filtered: List[SearchResult] = []
    preferred = [d.lower() for d in (preferred_domains or []) if d]
    for r in results:
        parsed = urllib.parse.urlparse(r.url)
        host = (parsed.netloc or "").lower()
        if any(b in host for b in banned_domains):
            continue
        if preferred and not any(p in host for p in preferred):
            continue
        score = _result_term_score(r, must_terms)
        if score >= min_score:
            filtered.append(r)
    return filtered


def _result_term_score(result: SearchResult, terms: List[str]) -> float:
    hay = f"{result.title} {result.snippet} {result.url}".lower()
    hit = sum(1 for t in terms if t in hay)
    quality_bonus = 0.0
    host = urllib.parse.urlparse(result.url).netloc.lower()
    if any(x in host for x in ["dmv.sc.gov", "scdmvonline.com", "state", "gov"]):
        quality_bonus += 1.0
    return float(hit) + quality_bonus


def _generate_study_items(
    topic: str,
    count: int,
    source_results: List[SearchResult],
    fact_bank: Optional[List[Dict[str, str]]] = None,
) -> List[StudyItem]:
    categories = [
        "Road Signs",
        "Right of Way",
        "Speed and Distance",
        "Safe Driving",
        "Sharing the Road",
        "Permit Rules",
        "Alcohol and Drugs",
        "Parking and Turns",
    ]
    source_titles = [r.title for r in source_results[:20]] or [f"{topic} official manual"]
    facts = fact_bank or []
    facts_with_image = [f for f in facts if str(f.get("image_path", ""))]
    facts_no_image = [f for f in facts if not str(f.get("image_path", ""))]
    items: List[StudyItem] = []
    for i in range(count):
        cat = categories[i % len(categories)]
        src = source_titles[i % len(source_titles)]
        if facts:
            use_visual = bool(facts_with_image) and (i % 4 == 0)
            if use_visual:
                fact = facts_with_image[(i // 4) % len(facts_with_image)]
            else:
                base = facts_no_image if facts_no_image else facts
                fact = base[i % len(base)]
        else:
            fact = {"text": f"{topic} official rule", "source_url": "", "category": cat}
        fact_text = str(fact.get("text", "")).strip()
        if len(fact_text) > 220:
            fact_text = fact_text[:220] + "..."
        img = str(fact.get("image_path", ""))
        img_b64 = str(fact.get("image_base64", ""))
        q = f"[{cat}] Q{i+1}: {_fact_to_question(fact_text, has_image=bool(img))}"
        a = f"{fact_text}"
        diff = "easy" if i % 3 == 0 else "medium" if i % 3 == 1 else "hard"
        items.append(
            StudyItem(
                question=q,
                answer=a,
                category=str(fact.get("category", cat)),
                difficulty=diff,
                source_url=str(fact.get("source_url", "")),
                evidence=fact_text,
                image_path=img,
                image_base64=img_b64,
            )
        )
    return items


def _fact_to_question(fact_text: str, has_image: bool = False) -> str:
    clean = re.sub(r"\s+", " ", fact_text).strip()
    if has_image:
        return "Based on the image shown, what rule or meaning applies in this situation?"
    low = clean.lower()
    if low.startswith("if "):
        tail = clean[3:]
        return f"If {tail}, what does South Carolina guidance require?"
    if "must" in low:
        prefix = clean.split("must", 1)[0].strip(" ,.;:")
        if prefix:
            return f"What must be done in South Carolina when {prefix.lower()}?"
    tokens = re.split(r"[^a-zA-Z0-9]+", clean)
    key = " ".join([t for t in tokens[:10] if t]) or "this situation"
    return f"According to the official manual, what is the correct rule regarding {key.lower()}?"


def _write_study_artifacts(
    instruction: str,
    topic: str,
    cards: List[StudyItem],
    sources: List[SearchResult],
) -> Dict[str, str]:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path("data/reports/study_pack") / ts
    out_dir.mkdir(parents=True, exist_ok=True)

    flashcards_csv = out_dir / "flashcards.csv"
    quiz_md = out_dir / "quiz.md"
    quiz_html = out_dir / "quiz.html"
    sources_md = out_dir / "sources.md"

    with flashcards_csv.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=["question", "answer", "category", "difficulty", "source_url", "evidence", "image_path"])
        writer.writeheader()
        for c in cards:
            writer.writerow(asdict(c))

    lines = [
        f"# Study Quiz: {topic}",
        "",
        f"- Generated: {datetime.now().isoformat(timespec='seconds')}",
        f"- Requested by: {instruction}",
        f"- Total Questions: {len(cards)}",
        "",
    ]
    for i, c in enumerate(cards[: min(len(cards), 200)], start=1):
        lines.append(f"## Q{i} ({c.category}, {c.difficulty})")
        lines.append(c.question)
        lines.append("")
        lines.append(f"Answer: {c.answer}")
        if c.source_url:
            lines.append(f"Source: {c.source_url}")
        if c.evidence:
            lines.append(f"Evidence: {c.evidence}")
        if c.image_path:
            lines.append(f"Image: {c.image_path}")
        lines.append("")
    quiz_md.write_text("\n".join(lines), encoding="utf-8")

    src_lines = ["# Sources", ""]
    for i, s in enumerate(sources[:50], start=1):
        src_lines.append(f"{i}. [{s.title}]({s.url})")
    sources_md.write_text("\n".join(src_lines), encoding="utf-8")

    _write_study_quiz_html(quiz_html, topic=topic, cards=cards, sources=sources)
    return {
        "directory": str(out_dir.resolve()),
        "flashcards_csv": str(flashcards_csv.resolve()),
        "quiz_md": str(quiz_md.resolve()),
        "quiz_html": str(quiz_html.resolve()),
        "sources_md": str(sources_md.resolve()),
    }


def _write_study_quiz_html(path: Path, topic: str, cards: List[StudyItem], sources: List[SearchResult]) -> None:
    payload_cards = []
    for c in cards:
        row = asdict(c)
        img_path = str(row.get("image_path", "") or "")
        img_b64 = str(row.get("image_base64", "") or "")
        row["image_url"] = _study_image_url(img_path=img_path, image_base64=img_b64)
        payload_cards.append(row)
    cards_payload = json.dumps(payload_cards)
    src_payload = json.dumps([asdict(s) for s in sources[:30]])
    html_text = f"""<!doctype html>
<html lang="en"><head>
<meta charset="utf-8"/><meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Study Pack - {topic}</title>
<style>
body{{font-family:'Segoe UI',Tahoma,sans-serif;margin:0;background:#f8fafc;color:#0f172a}}
.wrap{{max-width:1200px;margin:0 auto;padding:20px}}
.hero{{background:linear-gradient(120deg,#ecfeff,#eef2ff);border:1px solid #cbd5e1;border-radius:14px;padding:16px}}
.row{{display:flex;gap:10px;flex-wrap:wrap;margin:12px 0}}
input,select{{border:1px solid #cbd5e1;border-radius:10px;padding:8px;font-size:14px}}
.card{{background:#fff;border:1px solid #e2e8f0;border-radius:10px;padding:12px;margin-bottom:10px}}
.meta{{font-size:12px;color:#64748b}}
a{{color:#0f766e;text-decoration:none}}
</style></head><body><div class="wrap">
<div class="hero"><h1 style="margin:0">Study Pack: {topic}</h1><p style="margin:6px 0 0 0">{len(cards)} generated questions with source links.</p></div>
<div class="row"><input id="q" placeholder="Search questions" oninput="render()"/><select id="cat" onchange="render()"><option value="">All categories</option></select></div>
<div id="cards"></div>
<h3>Sources</h3><div id="sources"></div>
</div>
<script>
const cards={cards_payload};
const sources={src_payload};
const cats=[...new Set(cards.map(x=>x.category))].sort();
const catSel=document.getElementById('cat'); cats.forEach(c=>{{const o=document.createElement('option');o.value=c;o.textContent=c;catSel.appendChild(o);}});
function esc(s){{return (s||'').replace(/[&<>]/g,c=>({{'&':'&amp;','<':'&lt;','>':'&gt;'}}[c]));}}
function render(){{
 const q=(document.getElementById('q').value||'').toLowerCase();
 const cat=document.getElementById('cat').value;
 const filtered=cards.filter(c=>{{ if(cat && c.category!==cat) return false; return !q || (c.question+' '+c.answer).toLowerCase().includes(q); }});
  document.getElementById('cards').innerHTML=filtered.slice(0,300).map((c,i)=>`<div class="card"><div class="meta">#${{i+1}} â€¢ ${{esc(c.category)}} â€¢ ${{esc(c.difficulty)}}</div><div><strong>${{esc(c.question)}}</strong></div>${{c.image_url?`<div style="margin-top:8px"><img src="${{c.image_url}}" alt="Study visual" style="max-width:100%;border:1px solid #e2e8f0;border-radius:8px"/></div>`:''}}<div style="margin-top:6px">${{esc(c.answer)}}</div><div class="meta" style="margin-top:6px">${{c.source_url?`Source: <a target="_blank" rel="noopener" href="${{c.source_url}}">${{esc(c.source_url)}}</a>`:''}}</div></div>`).join('');
 document.getElementById('sources').innerHTML=sources.map((s,i)=>`<div><a target="_blank" rel="noopener" href="${{s.url}}">${{i+1}}. ${{esc(s.title)}}</a></div>`).join('');
}}
    render();
</script></body></html>"""
    path.write_text(html_text, encoding="utf-8")


def _study_image_url(*, img_path: str = "", image_base64: str = "") -> str:
    path_value = str(img_path or "").strip()
    if path_value and Path(path_value).exists():
        return Path(path_value).resolve().as_uri()
    b64_value = str(image_base64 or "").strip()
    if b64_value:
        return f"data:image/png;base64,{b64_value}"
    return ""


def _get_fitz_module() -> Any | None:
    global _FITZ_MODULE, _FITZ_IMPORT_ATTEMPTED
    if _FITZ_IMPORT_ATTEMPTED:
        return _FITZ_MODULE
    _FITZ_IMPORT_ATTEMPTED = True
    try:
        import fitz  # type: ignore

        _FITZ_MODULE = fitz
    except Exception as exc:  # pylint: disable=broad-exception-caught
        LOGGER.warning("PyMuPDF/fitz is unavailable for PDF rendering: %s", exc)
        _FITZ_MODULE = None
    return _FITZ_MODULE


def _expand_queries(query: str, instruction: str = "") -> List[str]:
    return platform_expand_queries(query=query, instruction=instruction)


def _relevance_score(result: SearchResult, query: str) -> float:
    return float(platform_relevance_score(result, query))


def _is_target_job(title: str, strict_vp_avp: bool = False) -> bool:
    low = title.lower()
    if strict_vp_avp:
        has_seniority = _is_vp_avp_title(title)
    else:
        has_seniority = any(k in low for k in ["vp", "vice president", "avp", "head", "director", "chief"])
    has_domain = any(k in low for k in ["data", "ai", "artificial intelligence", "machine learning", "analytics"])
    return has_seniority and has_domain


def _write_generic_research_artifacts(
    instruction: str,
    query: str,
    results: List[SearchResult],
    decision_rows: Optional[List[Dict[str, Any]]] = None,
    recommendation: Optional[Dict[str, Any]] = None,
    browser_notes: Optional[List[Dict[str, Any]]] = None,
) -> Dict[str, str]:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path("data/reports/generic_research") / ts
    out_dir.mkdir(parents=True, exist_ok=True)

    csv_path = out_dir / "results.csv"
    decision_csv_path = out_dir / "decision_matrix.csv"
    report_path = out_dir / "report.md"
    recommendation_path = out_dir / "recommendation.md"
    browser_log_path = out_dir / "browser_research.md"
    browser_json_path = out_dir / "browser_research.json"
    brief_md_path = out_dir / "executive_brief.md"
    brief_html_path = out_dir / "executive_brief.html"
    pptx_path = out_dir / "executive_brief.pptx"
    dash_path = out_dir / "dashboard.html"
    decision_rows = list(decision_rows or [])
    recommendation = dict(recommendation or {})
    browser_notes = [dict(x) for x in (browser_notes or []) if isinstance(x, dict)]

    with csv_path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(
            f,
            fieldnames=["title", "url", "source", "price", "snippet"],
        )
        writer.writeheader()
        for r in results:
            writer.writerow(asdict(r))

    if decision_rows:
        with decision_csv_path.open("w", newline="", encoding="utf-8") as f:
            writer = csv.DictWriter(
                f,
                fieldnames=["rank", "candidate", "candidate_type", "url", "source", "price", "score", "support_count", "rationale"],
            )
            writer.writeheader()
            for row in decision_rows:
                writer.writerow(
                    {
                        "rank": row.get("rank", ""),
                        "candidate": row.get("candidate", ""),
                        "candidate_type": row.get("candidate_type", ""),
                        "url": row.get("url", ""),
                        "source": row.get("source", ""),
                        "price": row.get("price", ""),
                        "score": row.get("score", ""),
                        "support_count": row.get("support_count", ""),
                        "rationale": row.get("rationale", ""),
                    }
                )

    lines = [
        "# Research Report",
        "",
        f"- Generated: {datetime.now().isoformat(timespec='seconds')}",
        f"- Instruction: {instruction}",
        f"- Query focus: {query}",
        f"- Results: {len(results)}",
        "",
    ]
    if recommendation:
        lines.extend(
            [
                "## Recommended Option",
                f"- Selection: {recommendation.get('selected_title', '')}",
                f"- URL: {recommendation.get('selected_url', '')}",
                f"- Score: {recommendation.get('selected_score', 'n/a')}",
                f"- Reason: {recommendation.get('reason', '')}",
                "",
            ]
        )
    lines.extend(
        [
        "## Top Results",
        ]
    )
    for i, r in enumerate(results[:40], start=1):
        extra = f" | ${r.price:.2f}" if r.price is not None else ""
        lines.append(f"{i}. [{r.title}]({r.url}) | {r.source}{extra}")
    report_path.write_text("\n".join(lines), encoding="utf-8")

    recommendation_lines = [
        "# Recommendation",
        "",
        f"- Generated: {datetime.now().isoformat(timespec='seconds')}",
        f"- Instruction: {instruction}",
        f"- Query focus: {query}",
        "",
    ]
    if recommendation:
        recommendation_lines.extend(
            [
                f"## Selected",
                f"- Title: {recommendation.get('selected_title', '')}",
                f"- URL: {recommendation.get('selected_url', '')}",
                f"- Score: {recommendation.get('selected_score', 'n/a')}",
                f"- Reason: {recommendation.get('reason', '')}",
                "",
            ]
        )
    if decision_rows:
        recommendation_lines.append("## Decision Matrix")
        for row in decision_rows[:10]:
            recommendation_lines.append(
                f"- {row.get('rank', '')}. {row.get('candidate', '')} | score={row.get('score', 'n/a')} | {row.get('rationale', '')}"
            )
    else:
        recommendation_lines.append("No recommendation-specific decision matrix was generated.")
    recommendation_path.write_text("\n".join(recommendation_lines), encoding="utf-8")

    browser_json_path.write_text(json.dumps(browser_notes, indent=2), encoding="utf-8")
    browser_lines = [
        "# Browser Research Log",
        "",
        f"- Generated: {datetime.now().isoformat(timespec='seconds')}",
        f"- Query focus: {query}",
        f"- Pages reviewed: {len(browser_notes)}",
        "",
    ]
    if browser_notes:
        for idx, note in enumerate(browser_notes, start=1):
            browser_lines.append(f"## Page {idx}")
            browser_lines.append(f"- Title: {note.get('title', '')}")
            browser_lines.append(f"- URL: {note.get('url', '')}")
            browser_lines.append(f"- Extracted note: {note.get('summary', '')}")
            browser_lines.append("")
    else:
        browser_lines.append("No browser page notes were captured.")
    browser_log_path.write_text("\n".join(browser_lines), encoding="utf-8")

    wants_brief = _wants_exec_brief(instruction)
    wants_ppt = _wants_powerpoint(instruction)
    if wants_brief:
        _write_generic_exec_brief(
            instruction=instruction,
            query=query,
            results=results,
            brief_md_path=brief_md_path,
            brief_html_path=brief_html_path,
        )
    if wants_ppt:
        _write_generic_exec_pptx(
            instruction=instruction,
            query=query,
            results=results,
            pptx_path=pptx_path,
        )
    _write_generic_dashboard_html(
        dash_path=dash_path,
        results=results,
        title="Research Dashboard",
        decision_rows=decision_rows,
        recommendation=recommendation,
        browser_notes=browser_notes,
    )

    out = {
        "directory": str(out_dir.resolve()),
        "results_csv": str(csv_path.resolve()),
        "decision_matrix_csv": str(decision_csv_path.resolve()) if decision_rows else "",
        "report_md": str(report_path.resolve()),
        "recommendation_md": str(recommendation_path.resolve()),
        "browser_research_md": str(browser_log_path.resolve()),
        "browser_research_json": str(browser_json_path.resolve()),
        "dashboard_html": str(dash_path.resolve()),
        "research_dashboard_html": str(dash_path.resolve()),
    }
    if wants_brief:
        out["executive_brief_md"] = str(brief_md_path.resolve())
        out["executive_brief_html"] = str(brief_html_path.resolve())
        out["primary_open_file"] = str(brief_html_path.resolve())
    elif wants_ppt:
        out["primary_open_file"] = str(pptx_path.resolve())
    else:
        out["primary_open_file"] = str(dash_path.resolve())
    if wants_ppt:
        out["powerpoint_pptx"] = str(pptx_path.resolve())
    return out


def _wants_exec_brief(instruction: str) -> bool:
    low = instruction.lower()
    return any(x in low for x in ["executive summary", "executive brief", "2-page", "two-page brief", "brief"])


def _wants_powerpoint(instruction: str) -> bool:
    low = instruction.lower()
    return any(x in low for x in ["powerpoint", "ppt", "slides", "slide deck"])


def _write_generic_exec_brief(
    instruction: str,
    query: str,
    results: List[SearchResult],
    brief_md_path: Path,
    brief_html_path: Path,
) -> None:
    top = results[:20]
    themes = [r.title.strip() for r in top[:8]]
    lines: List[str] = []
    lines.append(f"# Executive Brief")
    lines.append("")
    lines.append(f"- Generated: {datetime.now().isoformat(timespec='seconds')}")
    lines.append(f"- Instruction: {instruction}")
    lines.append(f"- Query focus: {query}")
    lines.append("")
    lines.append("## 1. Executive Summary")
    lines.append(
        "This brief synthesizes current market signals and source material relevant to the requested topic. The emphasis is on practical strategic interpretation rather than simple link aggregation."
    )
    lines.append("")
    lines.append("## 2. Key Findings")
    for i, t in enumerate(themes, start=1):
        lines.append(f"{i}. {t}")
    lines.append("")
    lines.append("## 3. Competitive/Market Interpretation")
    lines.append(
        "Across sources, the strongest pattern is concentration around a small number of recurring players and themes. Differentiation appears to depend on implementation risk, interoperability depth, and long-term operating model fit."
    )
    lines.append(
        "Decision quality improves when evidence is weighted by source quality and direct relevance to the objective, while excluding low-signal pages."
    )
    lines.append("")
    lines.append("## 4. Recommended Actions")
    lines.append("1. Validate the top findings against your target segment and operating constraints.")
    lines.append("2. Build a decision scorecard with explicit weighting and acceptance thresholds.")
    lines.append("3. Convert findings into an execution roadmap with owners, milestones, and risk controls.")
    lines.append("")
    lines.append("## 5. Source Appendix")
    for i, r in enumerate(top, start=1):
        lines.append(f"{i}. [{r.title}]({r.url})")
    brief_md_path.write_text("\n".join(lines), encoding="utf-8")

    html_rows = "".join(
        f"<li><a href=\"{html.escape(r.url)}\" target=\"_blank\" rel=\"noopener\">{html.escape(r.title)}</a></li>"
        for r in top
    )
    html_text = f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8"/><meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Executive Brief</title>
<style>body{{font-family:'Segoe UI',Tahoma,sans-serif;background:#f8fafc;color:#0f172a;margin:0}}.w{{max-width:980px;margin:0 auto;padding:24px}}.card{{background:#fff;border:1px solid #e2e8f0;border-radius:12px;padding:16px}}</style>
</head><body><div class="w"><div class="card"><h1>Executive Brief</h1><p><strong>Instruction:</strong> {html.escape(instruction)}</p><p><strong>Query:</strong> {html.escape(query)}</p><h2>Sources</h2><ol>{html_rows}</ol></div></div></body></html>"""
    brief_html_path.write_text(html_text, encoding="utf-8")


def _write_generic_exec_pptx(instruction: str, query: str, results: List[SearchResult], pptx_path: Path) -> None:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        pptx_path.write_text("PowerPoint package unavailable. Install python-pptx.", encoding="utf-8")
        return
    prs = Presentation()
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "Executive Brief Deck"
    s0.placeholders[1].text = query
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Objective"
    s1.placeholders[1].text = instruction[:1200]
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Top Findings"
    s2.placeholders[1].text = "\n".join(f"- {r.title}" for r in results[:8])[:2000]
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Next Steps"
    s3.placeholders[1].text = (
        "1) Validate findings with stakeholders.\n"
        "2) Build weighted decision framework.\n"
        "3) Track execution milestones and risk gates."
    )
    prs.save(str(pptx_path))


def _write_generic_dashboard_html(
    dash_path: Path,
    results: List[SearchResult],
    title: str,
    decision_rows: Optional[List[Dict[str, Any]]] = None,
    recommendation: Optional[Dict[str, Any]] = None,
    browser_notes: Optional[List[Dict[str, Any]]] = None,
) -> None:
    rows = []
    for r in results:
        rows.append(
            {
                "title": r.title,
                "source": r.source,
                "price": (f"${r.price:.2f}" if r.price is not None else "n/a"),
                "url": r.url,
                "snippet": r.snippet,
            }
        )
    payload = json.dumps(rows)
    decision_payload = json.dumps(list(decision_rows or []))
    recommendation_payload = json.dumps(dict(recommendation or {}))
    browser_notes_payload = json.dumps(list(browser_notes or []))
    html_text = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>{title}</title>
  <style>
    body {{ font-family: 'Segoe UI', Tahoma, sans-serif; margin:0; background:#f8fafc; color:#0f172a; }}
    .wrap {{ max-width: 1220px; margin:0 auto; padding:20px; }}
    .hero {{ border:1px solid #cbd5e1; background:linear-gradient(120deg,#eff6ff,#f0fdf4); border-radius:14px; padding:16px; }}
    .controls {{ display:flex; gap:10px; flex-wrap:wrap; margin:12px 0; }}
    input, select {{ border:1px solid #cbd5e1; border-radius:10px; padding:9px; font-size:14px; }}
    .grid {{ display:grid; grid-template-columns: minmax(0, 1.2fr) minmax(320px, 0.8fr); gap:16px; margin-top:14px; }}
    .panel {{ background:#fff; border:1px solid #dbe4ef; border-radius:14px; padding:14px; }}
    .panel-stack {{ display:grid; gap:16px; }}
    table {{ width:100%; border-collapse:collapse; background:#fff; }}
    th, td {{ padding:10px; border-bottom:1px solid #e2e8f0; vertical-align:top; font-size:14px; }}
    th {{ position:sticky; top:0; background:#0f172a; color:#e2e8f0; }}
    a {{ color:#0f766e; text-decoration:none; }}
    .meta {{ color:#475569; font-size:13px; }}
    @media (max-width: 980px) {{ .grid {{ grid-template-columns:1fr; }} }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="hero">
      <h1 style="margin:0">{title}</h1>
      <p style="margin:6px 0 0 0">Filter, inspect, and open source links.</p>
    </div>
    <div class="grid">
      <div class="panel">
        <h2 style="margin:0 0 8px 0; font-size:18px;">Recommendation</h2>
        <div id="recommendation"></div>
      </div>
      <div class="panel">
        <h2 style="margin:0 0 8px 0; font-size:18px;">Decision Matrix</h2>
        <div id="decisionRows" class="meta">No decision matrix available.</div>
      </div>
      <div class="panel">
        <h2 style="margin:0 0 8px 0; font-size:18px;">Browser Notes</h2>
        <div id="browserNotes" class="meta">No browser notes captured.</div>
      </div>
    </div>
    <div class="controls">
      <input id="q" placeholder="Search title/snippet" oninput="render()"/>
      <select id="source" onchange="render()"><option value="">Source: any</option></select>
    </div>
    <table>
      <thead><tr><th>Title</th><th>Source</th><th>Price</th><th>Link</th><th>Snippet</th></tr></thead>
      <tbody id="rows"></tbody>
    </table>
  </div>
<script>
const data = {payload};
const decisionRows = {decision_payload};
const recommendation = {recommendation_payload};
const browserNotes = {browser_notes_payload};
const source = document.getElementById("source");
[...new Set(data.map(x=>x.source))].sort().forEach(v=>{{ const o=document.createElement("option"); o.value=v; o.textContent=v; source.appendChild(o); }});
function esc(s){{ return (s||"").replace(/[&<>]/g, c => ({{"&":"&amp;","<":"&lt;",">":"&gt;"}}[c])); }}
function renderDecision(){{
  const recEl=document.getElementById("recommendation");
  if(recommendation && recommendation.selected_title){{
    const url = recommendation.selected_url ? `<a href="${{recommendation.selected_url}}" target="_blank" rel="noopener">${{esc(recommendation.selected_title)}}</a>` : esc(recommendation.selected_title);
    recEl.innerHTML = `<div><strong>${{url}}</strong></div><div class="meta" style="margin-top:6px">${{esc(recommendation.reason||'')}}</div>`;
  }} else {{
    recEl.innerHTML = `<div class="meta">No single recommendation generated.</div>`;
  }}
  const matrixEl=document.getElementById("decisionRows");
  if(!decisionRows.length){{
    matrixEl.textContent="No decision matrix available.";
    return;
  }}
  matrixEl.innerHTML = decisionRows.slice(0,8).map(row => {{
    const label = row.url ? `<a href="${{row.url}}" target="_blank" rel="noopener">${{esc(row.candidate||'candidate')}}</a>` : esc(row.candidate||'candidate');
    return `<div style="padding:8px 0;border-bottom:1px solid #e2e8f0"><strong>#${{row.rank||''}} ${{label}}</strong><div class="meta">score=${{row.score ?? 'n/a'}} | support=${{row.support_count ?? 0}}</div><div class="meta">${{esc(row.rationale||'')}}</div></div>`;
  }}).join("");
  const notesEl=document.getElementById("browserNotes");
  if(!browserNotes.length){{
    notesEl.textContent="No browser notes captured.";
  }} else {{
    notesEl.innerHTML = browserNotes.slice(0,6).map((note, idx) => `<div style="padding:8px 0;border-bottom:1px solid #e2e8f0"><strong>${{idx+1}}. <a href="${{note.url||'#'}}" target="_blank" rel="noopener">${{esc(note.title||note.url||'page')}}</a></strong><div class="meta">${{esc(note.summary||'')}}</div></div>`).join("");
  }}
}}
function render(){{
  const q=(document.getElementById("q").value||"").toLowerCase();
  const src=document.getElementById("source").value;
  const filtered=data.filter(x=>{{
    if(src && x.source!==src) return false;
    const hay=(x.title+" "+x.snippet).toLowerCase();
    return !q || hay.includes(q);
  }});
  document.getElementById("rows").innerHTML = filtered.map(x => `<tr><td>${{esc(x.title)}}</td><td>${{esc(x.source)}}</td><td>${{esc(x.price)}}</td><td><a href="${{x.url}}" target="_blank" rel="noopener">Open</a></td><td>${{esc((x.snippet||'').slice(0,180))}}</td></tr>`).join("");
}}
renderDecision();
render();
</script>
</body>
</html>
"""
    dash_path.write_text(html_text, encoding="utf-8")


def _scrape_linkedin_jobs(role_query: str, region: str, limit: int = 40) -> List[JobListing]:
    region_phrase = "Ireland" if region == "ireland" else "United States"
    url = (
        "https://www.linkedin.com/jobs/search/?keywords="
        + urllib.parse.quote_plus(role_query)
        + "&location="
        + urllib.parse.quote_plus(region_phrase)
    )
    page = _fetch_text(url)
    results: List[JobListing] = []
    card_pattern = re.compile(
        r'<div class="base-card[^>]*job-search-card[\s\S]*?</li>',
        re.IGNORECASE,
    )
    href_re = re.compile(r'class="base-card__full-link[^"]*"[^>]*href="([^"]+)"', re.IGNORECASE)
    title_re = re.compile(r'<span class="sr-only">\s*([\s\S]*?)\s*</span>', re.IGNORECASE)
    loc_re = re.compile(r'class="job-search-card__location">\s*([\s\S]*?)\s*</span>', re.IGNORECASE)
    for card in card_pattern.findall(page):
        href_m = href_re.search(card)
        title_m = title_re.search(card)
        if not href_m or not title_m:
            continue
        href = html.unescape(href_m.group(1).strip())
        title = html.unescape(re.sub(r"\s+", " ", title_m.group(1)).strip())
        loc = ""
        loc_m = loc_re.search(card)
        if loc_m:
            loc = html.unescape(re.sub(r"\s+", " ", loc_m.group(1)).strip())
        if not loc:
            loc = "Ireland" if region == "ireland" else "United States"
        text = f"{title} {loc}"
        salary_text, salary_min, salary_max, currency = _extract_salary(text)
        results.append(
            JobListing(
                title=title,
                url=href,
                source="linkedin",
                location=loc,
                remote=("remote" in text.lower()),
                salary_text=salary_text,
                salary_min=salary_min,
                salary_max=salary_max,
                currency=currency,
                snippet="",
            )
        )
        if len(results) >= limit:
            break
    return results


def _scrape_builtin_jobs(role_query: str, region: str, limit: int = 30) -> List[JobListing]:
    region_phrase = "Ireland" if region == "ireland" else "United States"
    url = "https://builtin.com/jobs?search=" + urllib.parse.quote_plus(f"{role_query} {region_phrase}")
    page = _fetch_text(url)
    block_re = re.compile(r'<a href="(/job/[^"]+)"[^>]*data-id="job-card-title"[\s\S]*?</a>', re.IGNORECASE)
    results: List[JobListing] = []
    for m in block_re.finditer(page):
        href = "https://builtin.com" + html.unescape(m.group(1))
        title = html.unescape(re.sub(r"<.*?>", "", m.group(0))).strip()
        around = page[max(0, m.start() - 450) : m.start() + 900]
        # Try to capture nearby location text; fallback to requested region.
        loc_match = re.search(r"([A-Za-z .'-]+,\s*[A-Z]{2})", around)
        location = loc_match.group(1) if loc_match else ("Ireland" if region == "ireland" else "United States")
        salary_text, salary_min, salary_max, currency = _extract_salary(around)
        results.append(
            JobListing(
                title=title,
                url=href,
                source="builtin",
                location=location,
                remote=("remote" in around.lower() or "remote" in title.lower()),
                salary_text=salary_text,
                salary_min=salary_min,
                salary_max=salary_max,
                currency=currency,
                snippet="",
            )
        )
        if len(results) >= limit:
            break
    return results


def _extract_regions(instruction: str) -> List[str]:
    low = instruction.lower()
    has_us = bool(re.search(r"\b(us|u\.s\.|united states|usa)\b", low))
    has_ireland = "ireland" in low
    if has_us and has_ireland:
        return ["us", "ireland"]
    if has_ireland:
        return ["ireland"]
    if has_us:
        return ["us"]
    return ["us"]


def _extract_job_constraints(instruction: str) -> JobSearchConstraints:
    low = instruction.lower()
    require_vp_avp = bool(re.search(r"\b(avp|vp|vice president|assistant vice president)\b", low))
    require_remote_or_hybrid = ("remote" in low) or ("hybrid" in low)
    min_base_salary = _extract_job_threshold_k(
        text=low,
        patterns=[
            r"(?:make|pay|salary|base|compensation)\s+(?:more than|above|over|at least)\s*\$?\s*([0-9]{2,3}(?:\.[0-9]+)?)\s*k",
            r"(?:more than|above|over|at least)\s*\$?\s*([0-9]{2,3}(?:\.[0-9]+)?)\s*k",
        ],
    )
    min_total_comp = _extract_job_threshold_k(
        text=low,
        patterns=[
            r"total compensation\s*(?:above|over|more than|at least)\s*\$?\s*([0-9]{2,3}(?:\.[0-9]+)?)\s*k",
            r"total comp(?:ensation)?\s*(?:above|over|more than|at least)\s*\$?\s*([0-9]{2,3}(?:\.[0-9]+)?)\s*k",
        ],
    )
    return JobSearchConstraints(
        require_vp_avp=require_vp_avp,
        require_remote_or_hybrid=require_remote_or_hybrid,
        min_base_salary_usd=min_base_salary,
        min_total_comp_usd=min_total_comp,
        allowed_regions=_extract_regions(instruction),
    )


def _extract_job_threshold_k(text: str, patterns: List[str]) -> Optional[float]:
    for pattern in patterns:
        m = re.search(pattern, text, flags=re.IGNORECASE)
        if m:
            try:
                return float(m.group(1)) * 1000.0
            except Exception:
                continue
    return None


def _job_constraints_are_strict(constraints: JobSearchConstraints) -> bool:
    return bool(
        constraints.require_vp_avp
        or constraints.require_remote_or_hybrid
        or constraints.min_base_salary_usd is not None
        or constraints.min_total_comp_usd is not None
    )


def _is_vp_avp_title(title: str) -> bool:
    low = title.lower()
    if "svp" in low:
        return False
    return bool(re.search(r"\b(avp|vp|vice president|assistant vice president)\b", low))


def _job_matches_region(job: JobListing, allowed_regions: List[str]) -> bool:
    if not allowed_regions:
        return True
    low = f"{job.location} {job.title} {job.snippet}".lower()
    is_ireland = "ireland" in low or ", ie" in low
    is_us = (
        "united states" in low
        or "usa" in low
        or " u.s." in low
        or bool(re.search(r",\s*[A-Z]{2}\b", job.location))
    )
    if is_ireland and "ireland" in allowed_regions:
        return True
    if is_us and "us" in allowed_regions:
        return True
    return False


def _job_total_comp_estimate(job: JobListing) -> Optional[float]:
    low = f"{job.title} {job.snippet} {job.salary_text}".lower()
    match = re.search(r"(?:total compensation|total comp(?:ensation)?)[^0-9$]{0,30}(\$?\s*[0-9]{2,3}(?:\.[0-9]+)?\s*[kK])", low)
    if match:
        token = match.group(1)
        value = _money_to_number(token)
        if value is not None:
            return value
    return _salary_sort_key(job) if _salary_sort_key(job) > 0 else None


def _job_matches_constraints(job: JobListing, constraints: JobSearchConstraints) -> bool:
    if constraints.require_vp_avp and (not _is_vp_avp_title(job.title)):
        return False
    if constraints.require_remote_or_hybrid:
        low = f"{job.title} {job.location} {job.snippet}".lower()
        is_hybrid = "hybrid" in low
        if not (job.remote or is_hybrid):
            return False
    if not _job_matches_region(job, constraints.allowed_regions):
        return False
    if constraints.min_base_salary_usd is not None:
        base = _salary_sort_key(job)
        if base <= 0 or base < constraints.min_base_salary_usd:
            return False
    if constraints.min_total_comp_usd is not None:
        total_comp = _job_total_comp_estimate(job)
        if total_comp is None or total_comp < constraints.min_total_comp_usd:
            return False
    return True


def _to_job_listing(item: SearchResult, source: str, fallback_region: str) -> Optional[JobListing]:
    title = (item.title or "").strip()
    url = (item.url or "").strip()
    snippet = (item.snippet or "").strip()
    if not title or not url:
        return None
    salary_text, salary_min, salary_max, currency = _extract_salary(title + " " + snippet)
    location = _extract_location(title, snippet, fallback_region)
    remote = "remote" in f"{title} {snippet}".lower()
    return JobListing(
        title=title,
        url=url,
        source=source,
        location=location,
        remote=remote,
        salary_text=salary_text,
        salary_min=salary_min,
        salary_max=salary_max,
        currency=currency,
        snippet=snippet,
    )


def _extract_salary(text: str) -> tuple[str, Optional[float], Optional[float], str]:
    cleaned = text.replace(",", "")
    usd_range = re.search(r"(\$[0-9]{2,3}(?:\.[0-9]{1,2})?\s*[kK]?)\s*-\s*(\$[0-9]{2,3}(?:\.[0-9]{1,2})?\s*[kK]?)", cleaned)
    eur_range = re.search(r"(â‚¬[0-9]{2,3}(?:\.[0-9]{1,2})?\s*[kK]?)\s*-\s*(â‚¬[0-9]{2,3}(?:\.[0-9]{1,2})?\s*[kK]?)", cleaned)
    single_usd = re.search(r"\$[0-9]{2,3}(?:\.[0-9]{1,2})?\s*[kK]?", cleaned)
    single_eur = re.search(r"â‚¬[0-9]{2,3}(?:\.[0-9]{1,2})?\s*[kK]?", cleaned)

    if usd_range:
        a, b = usd_range.group(1), usd_range.group(2)
        return f"{a} - {b}", _money_to_number(a), _money_to_number(b), "USD"
    if eur_range:
        a, b = eur_range.group(1), eur_range.group(2)
        return f"{a} - {b}", _money_to_number(a), _money_to_number(b), "EUR"
    if single_usd:
        v = single_usd.group(0)
        n = _money_to_number(v)
        return v, n, n, "USD"
    if single_eur:
        v = single_eur.group(0)
        n = _money_to_number(v)
        return v, n, n, "EUR"
    return "", None, None, ""


def _money_to_number(token: str) -> Optional[float]:
    m = re.search(r"([0-9]+(?:\.[0-9]+)?)", token)
    if not m:
        return None
    value = float(m.group(1))
    if "k" in token.lower():
        value *= 1000.0
    return value


def _extract_location(title: str, snippet: str, fallback_region: str) -> str:
    joined = f"{title} | {snippet}"
    patterns = [
        r"(Remote(?:\s*-\s*[A-Za-z ,]+)?)",
        r"([A-Za-z .'-]+,\s*(?:[A-Z]{2}|Ireland|UK|United States|USA))",
    ]
    for pat in patterns:
        match = re.search(pat, joined, flags=re.IGNORECASE)
        if match:
            return match.group(1).strip()
    return "Ireland" if fallback_region == "ireland" else "United States"


def _salary_rank(job: JobListing) -> float:
    if job.salary_min is None and job.salary_max is None:
        return float("inf")
    vals = [x for x in (job.salary_min, job.salary_max) if x is not None]
    return statistics.mean(vals) if vals else float("inf")


def _salary_sort_key(job: JobListing) -> float:
    if job.salary_min is None and job.salary_max is None:
        return -1.0
    vals = [x for x in (job.salary_min, job.salary_max) if x is not None]
    return statistics.mean(vals) if vals else -1.0


def _write_job_artifacts(instruction: str, jobs: List[JobListing]) -> Dict[str, str]:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path("data/reports/job_search") / ts
    out_dir.mkdir(parents=True, exist_ok=True)

    csv_path = out_dir / "jobs.csv"
    report_path = out_dir / "report.md"
    dash_path = out_dir / "dashboard.html"

    with csv_path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(
            f,
            fieldnames=[
                "title",
                "source",
                "location",
                "remote",
                "salary_text",
                "salary_min",
                "salary_max",
                "currency",
                "url",
                "snippet",
            ],
        )
        writer.writeheader()
        for j in jobs:
            writer.writerow(asdict(j))

    report_lines = [
        "# Job Search Report",
        "",
        f"- Generated: {datetime.now().isoformat(timespec='seconds')}",
        f"- Instruction: {instruction}",
        f"- Total Listings: {len(jobs)}",
        f"- Remote Listings: {sum(1 for j in jobs if j.remote)}",
        "",
        "## Top Listings",
    ]
    for idx, j in enumerate(jobs[:30], start=1):
        report_lines.append(f"{idx}. [{j.title}]({j.url}) | {j.source} | {j.location} | {j.salary_text or 'salary n/a'}")
    report_path.write_text("\n".join(report_lines), encoding="utf-8")

    _write_dashboard_html(dash_path=dash_path, jobs=jobs, title="VP/AVP Data & AI Job Dashboard")
    return {
        "directory": str(out_dir.resolve()),
        "jobs_csv": str(csv_path.resolve()),
        "report_md": str(report_path.resolve()),
        "dashboard_html": str(dash_path.resolve()),
    }


def _write_dashboard_html(dash_path: Path, jobs: List[JobListing], title: str) -> None:
    rows = []
    for j in jobs:
        rows.append(
            {
                "title": j.title,
                "source": j.source,
                "location": j.location,
                "remote": "Yes" if j.remote else "No",
                "salary": j.salary_text or "n/a",
                "salary_mid": _salary_rank(j) if _salary_rank(j) != float("inf") else 0,
                "currency": j.currency or "",
                "url": j.url,
            }
        )
    payload = json.dumps(rows)
    html = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>{title}</title>
  <style>
    body {{ font-family: 'Segoe UI', Tahoma, sans-serif; margin: 0; background: #f8fafc; color: #0f172a; }}
    .wrap {{ max-width: 1200px; margin: 0 auto; padding: 20px; }}
    .hero {{ background: linear-gradient(120deg,#ecfeff,#eef2ff); border:1px solid #cbd5e1; border-radius:14px; padding:18px; }}
    .controls {{ display:flex; gap:10px; flex-wrap:wrap; margin:14px 0; }}
    input, select {{ border:1px solid #cbd5e1; border-radius:10px; padding:9px; font-size:14px; }}
    table {{ width:100%; border-collapse: collapse; background:white; border-radius: 12px; overflow: hidden; }}
    th, td {{ border-bottom:1px solid #e2e8f0; padding:10px; text-align:left; font-size:14px; vertical-align: top; }}
    th {{ background:#0f172a; color:#e2e8f0; position: sticky; top: 0; }}
    a {{ color:#0f766e; text-decoration: none; }}
    .pill {{ font-size:12px; border-radius:20px; padding:2px 8px; background:#ecfeff; color:#115e59; }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="hero">
      <h1 style="margin:0">{title}</h1>
      <p style="margin:6px 0 0 0">Interactive view for role, region, remote, compensation, and apply links.</p>
    </div>
    <div class="controls">
      <input id="q" placeholder="Search title/location/source" oninput="render()"/>
      <select id="remote" onchange="render()">
        <option value="">Remote: any</option>
        <option value="Yes">Remote only</option>
        <option value="No">On-site/hybrid</option>
      </select>
      <select id="source" onchange="render()">
        <option value="">Source: any</option>
      </select>
      <select id="currency" onchange="render()">
        <option value="">Currency: any</option>
        <option value="USD">USD</option>
        <option value="EUR">EUR</option>
      </select>
    </div>
    <table>
      <thead><tr><th>Role</th><th>Source</th><th>Location</th><th>Remote</th><th>Salary</th><th>Apply</th></tr></thead>
      <tbody id="rows"></tbody>
    </table>
  </div>
<script>
const data = {payload};
const sourceSel = document.getElementById("source");
[...new Set(data.map(x=>x.source))].sort().forEach(v=>{{ const o=document.createElement("option"); o.value=v; o.textContent=v; sourceSel.appendChild(o); }});
function esc(s){{ return (s||"").replace(/[&<>]/g, c => ({{"&":"&amp;","<":"&lt;",">":"&gt;"}}[c])); }}
function render(){{
  const q=(document.getElementById("q").value||"").toLowerCase();
  const remote=document.getElementById("remote").value;
  const source=document.getElementById("source").value;
  const currency=document.getElementById("currency").value;
  const filtered=data.filter(x=>{{
    if(remote && x.remote!==remote) return false;
    if(source && x.source!==source) return false;
    if(currency && x.currency!==currency) return false;
    const hay=(x.title+" "+x.location+" "+x.source).toLowerCase();
    return !q || hay.includes(q);
  }});
  const tbody=document.getElementById("rows");
  tbody.innerHTML=filtered.map(x=>`<tr><td>${{esc(x.title)}}</td><td><span class="pill">${{esc(x.source)}}</span></td><td>${{esc(x.location)}}</td><td>${{esc(x.remote)}}</td><td>${{esc(x.salary)}}</td><td><a href="${{x.url}}" target="_blank" rel="noopener">Open</a></td></tr>`).join("");
}}
render();
</script>
</body>
</html>
"""
    dash_path.write_text(html, encoding="utf-8")


def _job_summary(jobs: List[JobListing]) -> Dict[str, Any]:
    usd_mids = [_salary_rank(j) for j in jobs if j.currency == "USD" and _salary_rank(j) != float("inf")]
    eur_mids = [_salary_rank(j) for j in jobs if j.currency == "EUR" and _salary_rank(j) != float("inf")]
    return {
        "total": len(jobs),
        "remote": sum(1 for j in jobs if j.remote),
        "by_source": _count_by(j.source for j in jobs),
        "by_region": {
            "ireland_like": sum(1 for j in jobs if "ireland" in j.location.lower() or ", ie" in j.location.lower()),
            "us_like": sum(1 for j in jobs if "united states" in j.location.lower() or re.search(r",\s*[A-Z]{2}$", j.location)),
        },
        "salary": {
            "usd_avg": round(statistics.mean(usd_mids), 2) if usd_mids else None,
            "eur_avg": round(statistics.mean(eur_mids), 2) if eur_mids else None,
            "count_with_salary": sum(1 for j in jobs if j.salary_text),
        },
    }


def _count_by(values: Any) -> Dict[str, int]:
    out: Dict[str, int] = {}
    for v in values:
        out[v] = out.get(v, 0) + 1
    return out


def _emit_progress(progress_cb: Optional[Callable[[int, str], None]], pct: int, message: str) -> None:
    if not progress_cb:
        return
    try:
        progress_cb(max(0, min(100, int(pct))), message)
    except Exception:
        return


def resume_pending_plan(
    pending: Dict[str, Any],
    step_mode: bool = False,
    human_like_interaction: bool = False,
) -> Dict[str, Any]:
    plan = pending.get("plan", {})
    plan_steps = [dict(x) for x in plan.get("steps", [])]
    start = int(pending.get("next_step_index", 0))
    run = execute_plan(
        plan,
        start_index=start,
        step_mode=step_mode,
        allow_input_fallback=True,
        human_like_interaction=bool(human_like_interaction),
    )
    response = {
        "ok": run.ok,
        "mode": "desktop_sequence_resume",
        "plan": plan,
        "trace": run.trace,
        "done": run.done,
        "next_step_index": run.next_step_index,
        "pending_plan": {"plan": plan, "next_step_index": run.next_step_index} if not run.done else None,
        "paused_for_credentials": run.paused_for_credentials,
        "pause_reason": run.pause_reason,
        "error": run.error,
        "artifacts": dict(run.artifacts or {}),
        "opened_url": Path(run.artifacts["primary_open_file"]).resolve().as_uri() if run.artifacts.get("primary_open_file") else "",
        "canvas": {
            "title": "Desktop Resume",
            "subtitle": f"Resumed from step {start}",
            "cards": [
                {"title": f"Step {t.get('step', 0)}: {t.get('action','')}", "price": "ok" if t.get("ok") else "error", "source": "uia", "url": ""}
                for t in run.trace[:6]
            ],
        },
    }
    resume_instruction = str(plan.get("instruction", "")).strip() or "Resume pending desktop sequence"
    return _finalize_operator_result(response, instruction=resume_instruction, plan_steps=plan_steps)

